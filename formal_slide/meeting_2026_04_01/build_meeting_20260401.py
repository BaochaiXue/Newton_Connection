#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
import subprocess
import textwrap
from pathlib import Path

import markdown
from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pygments import lex
from pygments.lexers import PythonLexer
from pygments.token import Comment, Keyword, Name, Number, Operator, Punctuation, String, Text, Token

MEETING_DIR = Path(__file__).resolve().parent
ROOT = MEETING_DIR.parents[1]
TEMPLATE_PPTX = MEETING_DIR / "templates" / "My Adjust.pptx"
OUT_PPTX = MEETING_DIR / "bridge_meeting_20260401.pptx"
DECK_GIF_DIR = MEETING_DIR / "gif"
DEFAULT_MAX_PPTX_MB = 100.0
DEFAULT_MAX_GIF_MB = 40.0

REF_TITLE_LEFT = 311760
REF_TITLE_TOP = 444960
REF_TITLE_W = 8520120
REF_TITLE_H = 572400

REF_BODY_LEFT = 311760
REF_BODY_TOP = 1152360
REF_BODY_W = 8520120
REF_BODY_H = 3416040

REF_TITLE_SLIDE_TOP = 744480
REF_TITLE_SLIDE_H = 2052360
REF_SUBTITLE_TOP = 2834280
REF_SUBTITLE_H = 792360

REF_GRID_COMMON_LEFT = 482040
REF_GRID_COMMON_TOP = 1170720
REF_GRID_COMMON_W = 8179560
REF_GRID_COMMON_H = 273960
REF_GRID_LABEL_W = 3983040
REF_GRID_LABEL_H = 243720
REF_GRID_PIC_W = 1735200
REF_GRID_PIC_H = 975960
REF_GRID_TL_LABEL_LEFT = 482040
REF_GRID_TR_LABEL_LEFT = 4678560
REF_GRID_TOP_LABEL_TOP = 1706760
REF_GRID_BOT_LABEL_TOP = 3129480
REF_GRID_TL_PIC_LEFT = 1605960
REF_GRID_TR_PIC_LEFT = 5802480
REF_GRID_TOP_PIC_TOP = 1999440
REF_GRID_BOT_PIC_TOP = 3421872

TABLE_LEFT = 482040
TABLE_TOP = 1524000
TABLE_W = 8179560
TABLE_H = 2857500
TABLE_NOTE_TOP = 1170720
TABLE_NOTE_H = 243720

OFFLINE_GIF_DIR = MEETING_DIR / "slides_assets" / "gif"
RESULTS_META_ROOT = ROOT / "results_meta" / "tasks"

RECALL_CLOTH_GIF_SRC = OFFLINE_GIF_DIR / "cloth_cmp2x3.gif"
RECALL_ZEBRA_GIF_SRC = OFFLINE_GIF_DIR / "zebra_cmp2x3.gif"
RECALL_SLOTH_GIF_SRC = OFFLINE_GIF_DIR / "sloth_cmp2x3.gif"
RECALL_ROPE_GIF_SRC = OFFLINE_GIF_DIR / "rope_drag_on_cmp2x3.gif"
RECALL_CLOTH_OVERLAY_GIF_SRC = OFFLINE_GIF_DIR / "cloth_overlay1x3.gif"
RECALL_ZEBRA_OVERLAY_GIF_SRC = OFFLINE_GIF_DIR / "zebra_overlay1x3.gif"
RECALL_SLOTH_OVERLAY_GIF_SRC = OFFLINE_GIF_DIR / "sloth_overlay1x3.gif"
RECALL_ROPE_OVERLAY_GIF_SRC = OFFLINE_GIF_DIR / "rope_drag_on_overlay1x3.gif"
RECALL_BUNNY_M5_GIF_SRC = OFFLINE_GIF_DIR / "bunny_drop_m5.gif"
RECALL_BUNNY_M500_GIF_SRC = OFFLINE_GIF_DIR / "bunny_drop_m500.gif"

RECALL_CLOTH_GIF = DECK_GIF_DIR / "cloth_cmp2x3_deck.gif"
RECALL_ZEBRA_GIF = DECK_GIF_DIR / "zebra_cmp2x3_deck.gif"
RECALL_SLOTH_GIF = DECK_GIF_DIR / "sloth_cmp2x3_deck.gif"
RECALL_ROPE_GIF = DECK_GIF_DIR / "rope_drag_on_cmp2x3_deck.gif"
RECALL_CLOTH_OVERLAY_GIF = DECK_GIF_DIR / "cloth_overlay1x3_deck.gif"
RECALL_ZEBRA_OVERLAY_GIF = DECK_GIF_DIR / "zebra_overlay1x3_deck.gif"
RECALL_SLOTH_OVERLAY_GIF = DECK_GIF_DIR / "sloth_overlay1x3_deck.gif"
RECALL_ROPE_OVERLAY_GIF = DECK_GIF_DIR / "rope_drag_on_overlay1x3_deck.gif"
RECALL_BUNNY_M5_GIF = DECK_GIF_DIR / "bunny_drop_m5_deck.gif"
RECALL_BUNNY_M500_GIF = DECK_GIF_DIR / "bunny_drop_m500_deck.gif"

PREV_GIF_DIR = ROOT / "formal_slide" / "meeting_2026_03_25" / "gif"
RECALL_ROPE_WEIGHT_1KG_GIF = PREV_GIF_DIR / "rope_bunny_total1kg_m5_v1.gif"
RECALL_ROPE_WEIGHT_5KG_GIF = PREV_GIF_DIR / "rope_bunny_total5kg_m5_v1.gif"
RECALL_BUNNY_SUPPORT_GIF = PREV_GIF_DIR / "cloth_rigid_compare_bunny_m5_v1.gif"
RECALL_BOX_SUPPORT_GIF = PREV_GIF_DIR / "cloth_rigid_compare_box_m5_v1.gif"
RECALL_THIN_EAR_5X_GIF = PREV_GIF_DIR / "thin_ear_ccd5x_v3.gif"
RECALL_THIN_EAR_10X_GIF = PREV_GIF_DIR / "thin_ear_ccd10x_v3.gif"

RENDER_BUNNY_BOARD_PATH = ROOT / "scripts" / "render_bunny_penetration_collision_board.py"
NEWTON_CORE_SPRING_PATH = ROOT / "Newton" / "newton" / "newton" / "_src" / "solvers" / "semi_implicit" / "kernels_particle.py"
NEWTON_CORE_SEMIIMPLICIT_PATH = ROOT / "Newton" / "newton" / "newton" / "_src" / "solvers" / "semi_implicit" / "solver_semi_implicit.py"
PHYSTWIN_SPRING_WARP_CODE_PATH = ROOT / "PhysTwin" / "qqtt" / "model" / "diff_simulator" / "spring_mass_warp.py"
IMAGE_DIR = MEETING_DIR / "images"
CODE_PERF_PHYSICS_NEWTON_PNG = IMAGE_DIR / "code_perf_physics_newton.png"
CODE_PERF_PHYSICS_PHYSTWIN_PNG = IMAGE_DIR / "code_perf_physics_phystwin.png"
CODE_PERF_EXECUTION_NEWTON_PNG = IMAGE_DIR / "code_perf_execution_newton.png"
CODE_PERF_EXECUTION_PHYSTWIN_PNG = IMAGE_DIR / "code_perf_execution_phystwin.png"
PERF_NEWTON_SYSTEM_PNG = IMAGE_DIR / "perf_newton_system_summary.png"
PERF_PHYSTWIN_SYSTEM_PNG = IMAGE_DIR / "perf_phystwin_system_summary.png"
CODE_SELFCOLLISION_OBJECT_PNG = IMAGE_DIR / "code_selfcollision_object.png"
CODE_SELFCOLLISION_GROUND_PNG = IMAGE_DIR / "code_selfcollision_ground.png"
CODE_SELFCOLLISION_TABLE_PHYSTWIN_PNG = IMAGE_DIR / "code_selfcollision_table_phystwin.png"
CODE_SELFCOLLISION_TABLE_BRIDGE_PNG = IMAGE_DIR / "code_selfcollision_table_bridge.png"
CODE_SELFCOLLISION_MATCHED_PHYSTWIN_PNG = IMAGE_DIR / "code_selfcollision_matched_phystwin.png"
CODE_SELFCOLLISION_MATCHED_BRIDGE_PNG = IMAGE_DIR / "code_selfcollision_matched_bridge.png"
CODE_SELFCOLLISION_FORCE_GROUND_PHYSTWIN_PNG = IMAGE_DIR / "code_selfcollision_force_ground_phystwin.png"
CODE_SELFCOLLISION_FORCE_GROUND_BRIDGE_PNG = IMAGE_DIR / "code_selfcollision_force_ground_bridge.png"
CODE_SELFCOLLISION_CONTROLLER_PHYSTWIN_PNG = IMAGE_DIR / "code_selfcollision_controller_phystwin.png"
CODE_SELFCOLLISION_CONTROLLER_BRIDGE_PNG = IMAGE_DIR / "code_selfcollision_controller_bridge.png"
FORCE_DIAG_CODE_PNG = IMAGE_DIR / "code_force_diag_capture.png"
FORCE_LAYOUT_CODE_PNG = IMAGE_DIR / "code_force_diag_layout.png"
PERF_ROPE_CASE_GIF = DECK_GIF_DIR / "rope_perf_case_anchor.gif"

ROPE_PERF_ROOT = ROOT / "results" / "rope_perf_apples_to_apples"
ROPE_PERF_SUMMARY_JSON = ROPE_PERF_ROOT / "summary.json"
ROPE_PERF_E0_VIEWER_SUMMARY_JSON = ROPE_PERF_ROOT / "newton" / "E0_viewer_baseline_end_to_end" / "summary.json"


def _load_results_meta(task_slug: str) -> dict:
    path = RESULTS_META_ROOT / f"{task_slug}.json"
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def _meta_local_root(task_slug: str) -> Path | None:
    obj = _load_results_meta(task_slug)
    local_root = (obj.get("authoritative_run") or {}).get("local_artifact_root")
    if not local_root:
        return None
    path = Path(str(local_root))
    return path if path.is_absolute() else (ROOT / path)


def _meta_superseded_root(task_slug: str, *, run_status: str) -> Path | None:
    obj = _load_results_meta(task_slug)
    for row in obj.get("superseded_runs", []):
        if str(row.get("status")) == run_status:
            path = Path(str(row.get("local_artifact_root")))
            return path if path.is_absolute() else (ROOT / path)
    return None


def _resolve_latest_bunny_run() -> Path:
    meta_root = _meta_local_root("bunny_penetration_force_diagnostic")
    if meta_root is not None and meta_root.exists():
        return meta_root.resolve()
    fallback = (
        ROOT
        / "results"
        / "bunny_force_visualization"
        / "runs"
        / "20260330_214500_meeting_ready_force_matrix"
    )
    return fallback.resolve()


CURRENT_BUNNY_RUN = _resolve_latest_bunny_run()
HISTORICAL_BUNNY_RUN = (
    _meta_superseded_root("bunny_penetration_force_diagnostic", run_status="historical")
    or CURRENT_BUNNY_RUN
).resolve()
HISTORICAL_BUNNY_MATRIX_DIR = HISTORICAL_BUNNY_RUN / "artifacts" / "matrix"
CURRENT_BUNNY_BOARD_MP4 = CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board" / "collision_force_board_2x2.mp4"
CURRENT_BUNNY_BOARD_SUMMARY = CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board" / "summary.json"
CURRENT_BUNNY_BOARD_PUBLISH_GIF = CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board" / "collision_force_board_2x2.gif"
CURRENT_BUNNY_BOARD_GIF = DECK_GIF_DIR / "bunny_collision_board_2x2.gif"
CURRENT_BUNNY_BOARD_FIRST_FRAME = (
    CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board" / "collision_force_board_2x2_first_frame.png"
)
CURRENT_BUNNY_SLOW_BOARD_DIR = CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board_slow4x"
CURRENT_BUNNY_SLOW_BOARD_STEM = "collision_force_board_2x2_slow4x"
CURRENT_BUNNY_SLOW_BOARD_MP4 = CURRENT_BUNNY_SLOW_BOARD_DIR / f"{CURRENT_BUNNY_SLOW_BOARD_STEM}.mp4"
CURRENT_BUNNY_SLOW_BOARD_SUMMARY = CURRENT_BUNNY_SLOW_BOARD_DIR / "summary.json"
CURRENT_BUNNY_SLOW_BOARD_PUBLISH_GIF = CURRENT_BUNNY_SLOW_BOARD_DIR / f"{CURRENT_BUNNY_SLOW_BOARD_STEM}.gif"
CURRENT_BUNNY_SLOW_BOARD_GIF = DECK_GIF_DIR / "bunny_collision_board_2x2_slow4x.gif"
CURRENT_BUNNY_PANEL_MP4 = {
    "box_penalty": CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board" / "panels" / "box_penalty.mp4",
    "box_total": CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board" / "panels" / "box_total.mp4",
    "bunny_penalty": CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board" / "panels" / "bunny_penalty.mp4",
    "bunny_total": CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board" / "panels" / "bunny_total.mp4",
}
CURRENT_BUNNY_PANEL_PUBLISH_GIF = {
    name: CURRENT_BUNNY_RUN / "artifacts" / "collision_force_board" / "panels" / f"{name}.gif"
    for name in CURRENT_BUNNY_PANEL_MP4
}
CURRENT_BUNNY_PANEL_GIF = {
    name: DECK_GIF_DIR / f"{name}.gif" for name in CURRENT_BUNNY_PANEL_MP4
}
CURRENT_BUNNY_BOARD_PUBLISH_GIF_PROFILES = [
    (1600, 15, 224),
    (1440, 12, 224),
    (1280, 12, 192),
    (1120, 10, 192),
    (960, 10, 160),
    (960, 8, 128),
]
CURRENT_BUNNY_SLOW_BOARD_PUBLISH_GIF_PROFILES = list(CURRENT_BUNNY_BOARD_PUBLISH_GIF_PROFILES)
CURRENT_BUNNY_PANEL_PUBLISH_GIF_PROFILES = [
    (960, 15, 224),
    (960, 12, 192),
    (896, 12, 192),
    (832, 10, 160),
    (768, 10, 160),
    (640, 8, 128),
]
CURRENT_BUNNY_BOARD_GIF_PROFILES = [
    (1280, 12, 192),
    (1120, 10, 160),
    (960, 10, 160),
    (960, 8, 128),
]
CURRENT_BUNNY_SLOW_BOARD_GIF_PROFILES = list(CURRENT_BUNNY_BOARD_GIF_PROFILES)
CURRENT_BUNNY_PANEL_GIF_PROFILES = [
    (832, 10, 160),
    (768, 10, 160),
    (640, 8, 128),
]
ACCEPTED_BUNNY_BOARD_PNG = (
    HISTORICAL_BUNNY_MATRIX_DIR / "bunny_penetration_summary_board.png"
    if (HISTORICAL_BUNNY_MATRIX_DIR / "bunny_penetration_summary_board.png").exists()
    else CURRENT_BUNNY_BOARD_FIRST_FRAME
)
SELF_COLLISION_CAMPAIGN_DIR = (_meta_local_root("self_collision_transfer") or (ROOT / "Newton" / "phystwin_bridge" / "results" / "final_self_collision_campaign_20260331_033636_533f3d0")).resolve()
SELF_COLLISION_SLIDES_DIR = SELF_COLLISION_CAMPAIGN_DIR / "slides_update"


def _resolve_case_video_maps() -> tuple[dict[str, Path], dict[str, Path]]:
    summary_path = HISTORICAL_BUNNY_RUN / "summary.json"
    if summary_path.exists():
        obj = json.loads(summary_path.read_text(encoding="utf-8"))
        pheno: dict[str, Path] = {}
        force: dict[str, Path] = {}
        for case in obj.get("cases", []):
            slug = str(case.get("slug") or case.get("case"))
            if not slug:
                continue
            render_video = case.get("render_video")
            force_video = case.get("force_video")
            if render_video:
                pheno[slug] = Path(str(render_video))
            if force_video:
                force[slug] = Path(str(force_video))
        if pheno and force:
            return pheno, force

    pheno = {
        "bunny_baseline": HISTORICAL_BUNNY_MATRIX_DIR / "bunny_baseline" / "phenomenon" / "self_off" / "cloth_bunny_drop_off_m0p5.mp4",
        "box_control": HISTORICAL_BUNNY_MATRIX_DIR / "box_control" / "phenomenon" / "self_off" / "cloth_bunny_drop_off_m0p5.mp4",
        "bunny_low_inertia": HISTORICAL_BUNNY_MATRIX_DIR / "bunny_low_inertia" / "phenomenon" / "self_off" / "cloth_bunny_drop_off_m0p5.mp4",
        "bunny_larger_scale": HISTORICAL_BUNNY_MATRIX_DIR / "bunny_larger_scale" / "phenomenon" / "self_off" / "cloth_bunny_drop_off_m0p5.mp4",
    }
    force = {
        "bunny_baseline": HISTORICAL_BUNNY_MATRIX_DIR / "bunny_baseline" / "force_mechanism" / "self_off" / "force_diagnostic" / "force_diag_trigger_window.mp4",
        "box_control": HISTORICAL_BUNNY_MATRIX_DIR / "box_control" / "force_mechanism" / "self_off" / "force_diagnostic" / "force_diag_trigger_window.mp4",
        "bunny_low_inertia": HISTORICAL_BUNNY_MATRIX_DIR / "bunny_low_inertia" / "force_mechanism" / "self_off" / "force_diagnostic" / "force_diag_trigger_window.mp4",
        "bunny_larger_scale": HISTORICAL_BUNNY_MATRIX_DIR / "bunny_larger_scale" / "force_mechanism" / "self_off" / "force_diagnostic" / "force_diag_trigger_window.mp4",
    }
    return pheno, force


ACCEPTED_PHENO_MP4, ACCEPTED_FORCE_MP4 = _resolve_case_video_maps()
ACCEPTED_PHENO_GIF = {name: DECK_GIF_DIR / f"{name}_phenomenon.gif" for name in ACCEPTED_PHENO_MP4}
ACCEPTED_FORCE_GIF = {name: DECK_GIF_DIR / f"{name}_force.gif" for name in ACCEPTED_FORCE_MP4}
SLIDE_QUESTION_CLAIM_PNG = SELF_COLLISION_SLIDES_DIR / "01_question_and_claim.png"
SLIDE_SOURCE_EVIDENCE_PNG = SELF_COLLISION_SLIDES_DIR / "02_native_source_code_evidence.png"
SLIDE_NATIVE_FAILURE_PNG = SELF_COLLISION_SLIDES_DIR / "03_native_failure_matrix.png"
SLIDE_EXACTNESS_PNG = SELF_COLLISION_SLIDES_DIR / "04_phystwin_exactness.png"
SLIDE_FINAL_DEMO_PNG = SELF_COLLISION_SLIDES_DIR / "05_final_demo_frames.png"
SLIDE_STRICT_PARITY_PNG = SELF_COLLISION_SLIDES_DIR / "06_strict_parity.png"
ROBOT_DROP_OFF_ROOT = (_meta_local_root("native_robot_rope_drop_release") or (ROOT / "results" / "native_robot_rope_drop_release" / "runs" / "20260331_232106_native_franka_recoilfix_drag_off_w5")).resolve()
ROBOT_DROP_ON_ROOT = (_meta_superseded_root("native_robot_rope_drop_release", run_status="validated_ab_compare") or (ROOT / "results" / "native_robot_rope_drop_release" / "runs" / "20260331_232459_native_franka_recoilfix_drag_on_w5")).resolve()
ROBOT_DROP_BASELINE_OFF_MP4 = ROBOT_DROP_OFF_ROOT / "final_presentation.mp4"
ROBOT_DROP_BASELINE_ON_MP4 = ROBOT_DROP_ON_ROOT / "final_presentation.mp4"
ROBOT_DROP_BASELINE_OFF_GIF = DECK_GIF_DIR / "robot_drop_release_drag_off.gif"
ROBOT_DROP_BASELINE_ON_GIF = DECK_GIF_DIR / "robot_drop_release_drag_on.gif"
ROBOT_TABLETOP_ROOT = (
    _meta_local_root("robot_rope_franka_tabletop_push_hero")
    or (ROOT / "Newton" / "phystwin_bridge" / "results" / "robot_rope_franka" / "BEST_RUN")
).resolve()
ROBOT_TABLETOP_HERO_MP4 = ROBOT_TABLETOP_ROOT / "hero_presentation.mp4"
ROBOT_TABLETOP_VALIDATION_MP4 = ROBOT_TABLETOP_ROOT / "validation_camera.mp4"
ROBOT_TABLETOP_HERO_GIF = DECK_GIF_DIR / "robot_tabletop_push_hero.gif"
ROBOT_TABLETOP_VALIDATION_GIF = DECK_GIF_DIR / "robot_tabletop_push_validation.gif"
SELF_PARITY_SUPPORT_MP4 = (
    ROOT
    / "Newton"
    / "phystwin_bridge"
    / "results"
    / "tmp_off_vs_phystwin_302_compare_20260401"
    / "parity_support_demo"
    / "parity_support_demo.mp4"
)
SELF_PARITY_SUPPORT_GIF = DECK_GIF_DIR / "self_collision_parity_support.gif"

RECALL_DIRECT_GIF_SPECS = [
    (RECALL_CLOTH_GIF_SRC, RECALL_CLOTH_GIF, 800, 6, 80),
    (RECALL_ZEBRA_GIF_SRC, RECALL_ZEBRA_GIF, 800, 6, 80),
    (RECALL_SLOTH_GIF_SRC, RECALL_SLOTH_GIF, 800, 6, 80),
    (RECALL_ROPE_GIF_SRC, RECALL_ROPE_GIF, 800, 6, 80),
    (RECALL_CLOTH_OVERLAY_GIF_SRC, RECALL_CLOTH_OVERLAY_GIF, 800, 6, 80),
    (RECALL_ZEBRA_OVERLAY_GIF_SRC, RECALL_ZEBRA_OVERLAY_GIF, 800, 6, 80),
    (RECALL_SLOTH_OVERLAY_GIF_SRC, RECALL_SLOTH_OVERLAY_GIF, 800, 6, 80),
    (RECALL_ROPE_OVERLAY_GIF_SRC, RECALL_ROPE_OVERLAY_GIF, 800, 6, 80),
    (RECALL_BUNNY_M5_GIF_SRC, RECALL_BUNNY_M5_GIF, 800, 8, 96),
    (RECALL_BUNNY_M500_GIF_SRC, RECALL_BUNNY_M500_GIF, 800, 8, 96),
]


def _load_rope_perf_summary() -> dict:
    if not ROPE_PERF_SUMMARY_JSON.exists():
        return {"rows": []}
    return json.loads(ROPE_PERF_SUMMARY_JSON.read_text(encoding="utf-8"))


def _load_optional_json(path: Path) -> dict:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


ROPE_PERF_SUMMARY = _load_rope_perf_summary()
E0_VIEWER_SUMMARY = _load_optional_json(ROPE_PERF_E0_VIEWER_SUMMARY_JSON)
ROPE_PERF_ROWS = {row["stage"]: row for row in ROPE_PERF_SUMMARY.get("rows", [])}


def _rope_row(stage: str) -> dict:
    return ROPE_PERF_ROWS.get(stage, {})


def _fmt(row: dict, key: str, fmt: str) -> str:
    value = row.get(key)
    if value is None:
        return "n/a"
    return format(value, fmt)


E1_ROW = _rope_row("E1_viewer_end_to_end")
A0_ROW = _rope_row("A0_baseline_throughput")
A1_ROW = _rope_row("A1_precomputed_throughput")
A2_ROW = _rope_row("A2_baseline_attribution")
A3_ROW = _rope_row("A3_precomputed_attribution")
B0_ROW = _rope_row("B0_headless_throughput")
B1_ROW = _rope_row("B1_headless_attribution")


def _rope_payload(stage: str) -> dict:
    return _rope_row(stage).get("payload", {})
E1_PAYLOAD = _rope_payload("E1_viewer_end_to_end")
A0_PAYLOAD = _rope_payload("A0_baseline_throughput")
A1_PAYLOAD = _rope_payload("A1_precomputed_throughput")
A2_PAYLOAD = _rope_payload("A2_baseline_attribution")
A3_PAYLOAD = _rope_payload("A3_precomputed_attribution")
B0_PAYLOAD = _rope_payload("B0_headless_throughput")
B1_PAYLOAD = _rope_payload("B1_headless_attribution")

A3_AGG = A3_PAYLOAD.get("aggregate", {})
B1_AGG = B1_PAYLOAD.get("aggregate", {})

A3_BRIDGE_MS = (
    float(A3_AGG.get("write_kinematic_state", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("drag_correction", {}).get("mean_of_run_means_ms", 0.0))
)
A3_INTERNAL_MS = (
    float(A3_AGG.get("spring_forces", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("triangle_forces", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("bending_forces", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("tetra_forces", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("body_joint_forces", {}).get("mean_of_run_means_ms", 0.0))
)
A3_COLLISION_MS = (
    float(A3_AGG.get("particle_grid_build", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("model_collide", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("particle_contact_forces", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("triangle_contact_forces", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("body_contact_forces", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("particle_body_contact_forces", {}).get("mean_of_run_means_ms", 0.0))
)
A3_INTEGRATION_MS = (
    float(A3_AGG.get("integrate_particles", {}).get("mean_of_run_means_ms", 0.0))
    + float(A3_AGG.get("integrate_bodies", {}).get("mean_of_run_means_ms", 0.0))
)
A3_TOTAL_SUBSTEP_MS = float(A3_AGG.get("total_substep", {}).get("mean_of_run_means_ms", 0.0))
A3_UNEXPLAINED_MS = max(
    A3_TOTAL_SUBSTEP_MS - A3_BRIDGE_MS - A3_INTERNAL_MS - A3_COLLISION_MS - A3_INTEGRATION_MS,
    0.0,
)

B1_CONTROLLER_FRAME_MS = float(B1_AGG.get("controller_target", {}).get("mean_of_run_means_ms", 0.0))
B1_SIM_LAUNCH_MS = float(B1_AGG.get("simulator_launch", {}).get("mean_of_run_means_ms", 0.0))
B1_STATE_RESET_MS = float(B1_AGG.get("state_reset", {}).get("mean_of_run_means_ms", 0.0))
E0_VIEWER_FPS = float(E0_VIEWER_SUMMARY.get("viewer_fps_mean", 0.0))
E0_VIEWER_RTF = float(E0_VIEWER_SUMMARY.get("rtf_mean", 0.0))
E0_VIEWER_WALL_MS = float(E0_VIEWER_SUMMARY.get("wall_ms_mean", 0.0))
E1_VIEWER_FPS = float(E1_PAYLOAD.get("viewer_fps_mean", 0.0))
A0_VS_B0 = float(A0_ROW.get("ms_per_substep_mean", 0.0)) / max(float(B0_ROW.get("ms_per_substep_mean", 1.0e-12)), 1.0e-12)
A1_VS_B0 = float(A1_ROW.get("ms_per_substep_mean", 0.0)) / max(float(B0_ROW.get("ms_per_substep_mean", 1.0e-12)), 1.0e-12)
A0_TO_A1 = float(A0_ROW.get("ms_per_substep_mean", 0.0)) / max(float(A1_ROW.get("ms_per_substep_mean", 1.0e-12)), 1.0e-12)
A1_TO_E1 = float(E1_ROW.get("wall_ms_mean", 0.0)) / max(float(A1_ROW.get("wall_ms_mean", 1.0e-12)), 1.0e-12)
E0_TO_E1 = E0_VIEWER_WALL_MS / max(float(E1_ROW.get("wall_ms_mean", 1.0e-12)), 1.0e-12)
A3_POST_BRIDGE_MS = A3_INTERNAL_MS + A3_INTEGRATION_MS + A3_UNEXPLAINED_MS

RECALL_SLIDES: list[dict] = [
    {
        "kind": "title",
        "title": "Xinjie Zhang",
        "subtitle": [
            "April 1 meeting.",
            "1. Recall",
            "2. Performance analysis",
            "3. Penetration analysis",
            "4. Self-collision / ground-contact laws",
            "5. Robotic with deformable objects",
        ],
        "transcript": [
            "这一页是 opening speaker page。",
            "这次 meeting 我会按五段讲：1. recall，2. performance analysis，3. penetration analysis，4. controlled self-collision and ground-contact laws，5. robotic with deformable objects。",
            "所以 opening 要做的不是再把会议压成“short recall + two results”，而是先把这五段主线直接说清楚。",
        ],
    },
    {
        "kind": "grid",
        "title": "Recall: The Bridge Baseline Was Already Established",
        "common_settings": "Bridge existence is not the question anymore. These four baseline imports were already established.",
        "items": [
            ("Cloth baseline", RECALL_CLOTH_GIF),
            ("Zebra baseline", RECALL_ZEBRA_GIF),
            ("Sloth baseline", RECALL_SLOTH_GIF),
            ("Rope baseline", RECALL_ROPE_GIF),
        ],
        "transcript": [
            "第一页 recall 只讲一句话：import baseline 早就已经站住了。",
            "cloth、zebra、sloth、rope 这四个 baseline 都不是今天新证明的内容，所以后面的问题不再是 PhysTwin deformable 能不能 load 到 Newton。",
            "这页的 takeaway 很明确：bridge baseline already works。",
        ],
    },
    {
        "kind": "twocol",
        "title": "Recall: The Remaining Issue Is Contact Behavior, Not Bridge Existence",
        "common_settings": "Interaction already exists. The remaining question is contact behavior, especially geometry-sensitive penetration.",
        "left_label": "Same cloth + bunny support",
        "left_path": RECALL_BUNNY_SUPPORT_GIF,
        "right_label": "Same cloth + thick-box support",
        "right_path": RECALL_BOX_SUPPORT_GIF,
        "transcript": [
            "第二页 recall 只把问题进一步收窄。",
            "同样是 rigid support，bunny 明显比 thick box 更难，这说明当前 open question 更像是 geometry-sensitive contact behavior，而不是 bridge 根本没有工作。",
            "所以从这里开始，我们后面讨论的不是 bridge existence，而是 contact semantics 和 penetration mechanism。",
        ],
    },
    {
        "kind": "table_gif",
        "title": "H1: Replay Feeding Was The First Realtime Bottleneck",
        "note": "Same rope case. Same replay. Only replay feeding or rendering changes.",
        "gif_label": "Same rope replay object",
        "gif_path": PERF_ROPE_CASE_GIF,
        "columns": ["Measurement", "Value", "What it means"],
        "rows": [
            ["Old viewer path, render ON", f"FPS={E0_VIEWER_FPS:.2f} / RTF={E0_VIEWER_RTF:.3f}x", "baseline controller replay path misses realtime"],
            ["Current viewer path, render ON", f"FPS={E1_VIEWER_FPS:.2f} / RTF={_fmt(E1_ROW, 'rtf_mean', '.3f')}x", "precomputed controller replay path reaches realtime"],
            ["Current path, render OFF", f"RTF={_fmt(A1_ROW, 'rtf_mean', '.3f')}x", "rendering adds only a modest extra cost"],
        ],
        "transcript": [
            "这一页先直接回答你刚才提的 practical question：为什么我们原来的 rope viewer 到不了 realtime，现在又为什么能到 realtime。",
            "这里我专门补了一个旧路径对照实验。三行都还是同一个 `rope_double_hand`、同一条 replay trajectory、同样的 `dt` 和 `667` 个 substeps。区别只在于 Newton 侧 replay feeding 的实现，以及有没有把 rendering 打开。",
            f"旧的 visible viewer 路径，也就是 baseline controller replay path，`viewer FPS` 大约是 `{E0_VIEWER_FPS:.2f}`，`RTF` 只有 `{E0_VIEWER_RTF:.3f}x`，所以它确实达不到 realtime。现在的 visible viewer 路径换成 precomputed controller replay 以后，`viewer FPS` 大约是 `{E1_VIEWER_FPS:.2f}`，`RTF` 变成 `{_fmt(E1_ROW, 'rtf_mean', '.3f')}x`，所以它已经能过 realtime。",
            f"这两条 visible-viewer row 的核心区别，不是 physics 变了，也不是场景变了，而是 controller replay feeding 变了。我们把原来每个 substep 都要重复做的 interpolation 和 state write 预先展开，所以 viewer end-to-end wall time 大约改善了 `{E0_TO_E1:.2f}x`。",
            f"再和当前 path 的 render OFF 行放在一起看，A1 的 `RTF` 还能到 `{_fmt(A1_ROW, 'rtf_mean', '.3f')}x`。这说明现在 viewer 能过 realtime，不是因为突然把 rendering 变得特别轻了，而是因为 replay path 本身已经先被减负。",
            "所以 profiling 的价值就非常直接：它不是在绕开 real viewer，而是在解释 real viewer 为什么以前慢，以及为什么现在终于能过 realtime。边界是，这个结论只针对当前 clean rope replay case，不自动推广到所有更复杂的 contact scene。",
        ],
    },
    {
        "kind": "table_gif",
        "title": "H2: Fair Rope Benchmark = Same Replay, Physics, And GPU",
        "note": "Fairness controls only.",
        "gif_label": "Controlled rope replay",
        "gif_path": PERF_ROPE_CASE_GIF,
        "columns": ["Control", "Setting", "Why it matters"],
        "rows": [
            ["Case", "`rope_double_hand`", "same rope object and spring graph"],
            ["Replay", "same controller trajectory", "same input history, not two different tasks"],
            ["Physics", "`dt=5e-05`, `667` substeps", "same time resolution on both sides"],
            ["Hardware", "same RTX 4090", "same GPU / same workstation"],
            ["Primary comparison", "render OFF", "viewer/UI cost does not pollute the simulator benchmark"],
        ],
        "transcript": [
            "这一页先只讲 fairness controls，不急着解释 E1、A0、A1、B0 这些名字。",
            "我要先让观众知道：后面所有 rope profiling row 都是同一个 `rope_double_hand`，同一条 controller trajectory，同样的 `dt=5e-05` 和 `667` 个 substeps，而且都在同一张 RTX 4090 上跑。",
            f"这里最关键的一句是：主 apples-to-apples comparison 把 rendering 关掉，不是因为我们不关心 viewer，而是因为要先隔离 simulator replay 本身。与此同时，same-trajectory 也不是口头假设，当前 benchmark 记录的 IR 与 PhysTwin controller trajectory max abs diff 是 {B0_PAYLOAD.get('trajectory_parity', {}).get('controller_traj_max_abs_diff', 0.0):.1e}。",
            "这页对 real viewer 的意义是，它告诉听众后面那个 no-render benchmark 到底是在控制什么变量。边界是：这里只定义 fairness，不解释每个 benchmark row 的具体含义。",
        ],
    },
    {
        "kind": "table_gif",
        "title": "H3: E1, A0, A1, And B0 Are Four Rows Of One Benchmark",
        "note": "Plain-language row definitions.",
        "gif_label": "Same rope replay object",
        "gif_path": PERF_ROPE_CASE_GIF,
        "columns": ["Name", "Plain-language meaning", "Question answered"],
        "rows": [
            ["Newton viewer (E1)", "same rope replay, visible render ON", "how the practical rope viewer behaves end to end"],
            ["Newton A0", "same replay, render OFF, baseline controller replay path", "how fast Newton is before replay-overhead reduction"],
            ["Newton A1", "same replay, render OFF, precomputed controller replay path", "how much replay overhead we can remove on Newton"],
            ["PhysTwin B0", "same replay, render OFF, headless reference path", "the apples-to-apples throughput target"],
        ],
        "transcript": [
            "这一页再把 benchmark rows 的名字讲成人话，避免听众只看到一堆内部缩写。",
            "E1 是 visible render ON 的 real viewer row，所以它回答 practical question：viewer 本身现在跑得怎么样。A0 和 A1 都是同一个 no-render rope replay，只是 Newton 侧 feeding 同一条 trajectory 的方式不同。B0 则是同 case、同 trajectory 的 PhysTwin headless reference。",
            "所以 A0/A1/B0 不是三个不同任务，而是同一个 controlled rope benchmark 上的三个 measurement row。真正不同的是：有没有打开 rendering，以及 Newton 这边 controller replay overhead 有没有被预先压掉。",
            "这页对 real viewer 的意义是，把 practical row 和 simulator-only rows 的角色分清楚。这样后面看结果表时，观众才不会把 viewer FPS 和 ms/substep 混成一回事。",
            "边界同样要讲清楚：E1 不是拿来和 PhysTwin B0 直接比 wall time 的；它是 viewer-facing context row。真正 apples-to-apples 的 throughput 对照还是 A0/A1 对 B0。",
        ],
    },
    {
        "kind": "table_gif",
        "title": "Result P1: Newton A1 Is Still 3.30x Slower Than PhysTwin B0",
        "note": "Viewer row = practical speed. A0/A1/B0 = apples-to-apples simulator speed.",
        "gif_label": "Same rope replay",
        "gif_path": PERF_ROPE_CASE_GIF,
        "columns": ["Config", "Main metric", "RTF", "Plain-language reading"],
        "rows": [
            ["Newton viewer (E1)", f"FPS={E1_VIEWER_FPS:.2f}", _fmt(E1_ROW, "rtf_mean", ".3f"), "visible viewer, render ON"],
            ["Newton A0", _fmt(A0_ROW, "ms_per_substep_mean", ".6f") + " ms/substep", _fmt(A0_ROW, "rtf_mean", ".3f"), "baseline no-render replay"],
            ["Newton A1", _fmt(A1_ROW, "ms_per_substep_mean", ".6f") + " ms/substep", _fmt(A1_ROW, "rtf_mean", ".3f"), "precomputed no-render replay"],
            ["PhysTwin B0", _fmt(B0_ROW, "ms_per_substep_mean", ".6f") + " ms/substep", _fmt(B0_ROW, "rtf_mean", ".3f"), "same-case headless reference"],
        ],
        "transcript": [
            "这一页只回答数字本身，不先解释原因。",
            f"先看 practical row：E1 这个 visible rope viewer 大约是 `{E1_VIEWER_FPS:.2f} FPS`，`RTF` 是 `{_fmt(E1_ROW, 'rtf_mean', '.3f')}x`。所以它回答的是『viewer 现在能不能跑』这个问题。",
            f"再看 apples-to-apples rows：Newton A0 是 `{_fmt(A0_ROW, 'ms_per_substep_mean', '.6f')} ms/substep`，Newton A1 是 `{_fmt(A1_ROW, 'ms_per_substep_mean', '.6f')} ms/substep`，PhysTwin B0 是 `{_fmt(B0_ROW, 'ms_per_substep_mean', '.6f')} ms/substep`。",
            f"所以最核心的一句结论必须讲成人话：在同一个 clean rope replay benchmark 下，就算把 Newton 的 replay feeding 改成 A1 这种更轻的路径，它还是比 PhysTwin B0 慢大约 `{A1_VS_B0:.2f}x`。",
            "这页对 real viewer 的意义在于，它把两个问题分开了：viewer 今天能不能跑，是一回事；如果我们要更大 headroom，Newton 和 PhysTwin 在相同 replay benchmark 下还差多远，又是另一回事。",
            "边界也要讲清楚：viewer row 和 no-render rows 不应该被混成一个数。viewer row讲 end-to-end practical speed；A0/A1/B0 讲的是 simulator throughput。",
        ],
    },
    {
        "kind": "table_gif",
        "title": "H4: A0 -> A1 Isolates Replay Overhead, Not The Whole Gap",
        "note": "A0 and A1 differ only in how the same replay is fed into Newton.",
        "gif_label": "Same rope replay",
        "gif_path": PERF_ROPE_CASE_GIF,
        "columns": ["Evidence", "Value", "Human meaning"],
        "rows": [
            ["A0 -> A1 speedup", f"{A0_TO_A1:.2f}x", "controller replay overhead is real"],
            ["Newton A3 controller replay overhead", f"{A3_BRIDGE_MS:.3f} ms/substep", "feeding the same trajectory into Newton costs measurable time"],
            ["A1 vs B0 remaining gap", f"{A1_VS_B0:.2f}x slower", "replay overhead is not the whole answer"],
            ["Viewer ON vs A1 render OFF", f"{A1_TO_E1:.2f}x wall-time slowdown", "render adds less overhead than the remaining replay gap"],
        ],
        "transcript": [
            "这一页专门把 A0 和 A1 讲清楚，因为如果这里讲不明白，后面的推导都会很含糊。",
            "A0 和 A1 不是两个随便取的配置名。它们测的是同一条 rope replay、同一套 physics、同一组 `dt` 和 substeps。唯一想隔离出来的问题是：把同一条 replay trajectory 喂进 Newton，本身要花多少额外时间。",
            f"结果是 A0 到 A1 有 `{A0_TO_A1:.2f}x` 的 speedup，所以 controller replay overhead 的确存在。这里所谓的方法，不是换 physics，而是把 controller target 和 velocity 先按 substep 预计算好，避免每一步都重复做 interpolation 和写状态。同步 attribution 里，这部分大约是 `{A3_BRIDGE_MS:.3f} ms/substep`。",
            f"但这不是全部答案，因为把这部分降下来以后，Newton A1 相对 PhysTwin B0 还是慢 `{A1_VS_B0:.2f}x`。而且 viewer ON 相对 A1 render OFF 只慢 `{A1_TO_E1:.2f}x`，这说明 replay overhead 和 runtime organization 的问题，量级上比单纯 render cost 更值得先看。",
            "这页对 real viewer 的价值很直接：如果只优化 controller feeding，viewer 会变好一些，但不会把同 case 的 Newton-vs-PhysTwin headroom gap 自动消掉。它的边界是，这个结论只针对 clean rope replay，不是对所有 scene 的一刀切判断。",
        ],
    },
    {
        "kind": "code_twocol_large",
        "title": "H5: Residual Gap Points To Execution Structure",
        "note": "Core code + Nsight point to execution structure; this is not full root-cause proof.",
        "left_label": "[Newton/newton/newton/_src/solvers/semi_implicit/solver_semi_implicit.py : 141-145, 160-165, 172-177]\nNewton still advances the replay by visiting many solver stages in sequence.",
        "left_path": CODE_PERF_EXECUTION_NEWTON_PNG,
        "right_label": "[PhysTwin/qqtt/model/diff_simulator/spring_mass_warp.py : 768-782, 800-802]\nPhysTwin captures and replays the rope step through a forward graph.",
        "right_path": CODE_PERF_EXECUTION_PHYSTWIN_PNG,
        "transcript": [
            "这一页回答最后一个最容易被说过头的问题：残余差距到底说明了什么，又不说明什么。",
            "",
            "[Newton/newton/newton/_src/solvers/semi_implicit/solver_semi_implicit.py : 141-145, 160-165, 172-177]",
            "[What | Newton]",
            "L141-L145 这段先后调用 spring、triangle、bending 等不同 force stage，说明 Newton 这一条 replay path 不是一个单一的整步更新，而是把每一类力和接触拆成多个阶段依次执行。",
            "L160-L165 又继续进入不同 contact stage，所以就算当前 rope baseline 本身 contact 很弱，solver 组织方式仍然是“多段顺序推进”的结构，而不是一次性回放完一整步。",
            "L172-L177 最后才做 `integrate_particles` 和下一步切换，因此从 runtime 角度看，这条 path 更像 many separated execution stages，而不是一个天然的 graph replay。",
            "[Why | Newton]",
            "这段代码不能直接给出性能比例，但它确实回答了一个结构性问题：Newton rope replay 的 core path 现在就是按多个 stage 组织起来的。结合 A1 之后残余差距还在、以及 Nsight 里 `cuLaunchKernel` 仍然占主导，这个方向上的解释是有支撑的。",
            "[Risk | Newton]",
            "这段源码不能单独证明“Newton 一定就是因为 launch 太多才慢”。它只能说明 residual gap 至少和当前这种 staged execution organization 是相容的，而不是无条件把原因精确锁死。",
            "",
            "[PhysTwin/qqtt/model/diff_simulator/spring_mass_warp.py : 768-782, 800-802]",
            "[What | PhysTwin]",
            "L768-L782 这里在 `use_graph` 打开时直接 capture `self.step()`，并把 capture 得到的 graph 存起来，所以 PhysTwin 不是在 wrapper 外面临时做 batching，而是 rope core 自己就暴露了一条 graph-based replay path。",
            "L800-L802 又单独保留了 `forward_graph`，说明它不只是一次性 capture 一个训练图，而是前向 replay 也明确走 graph replay 语义。",
            "[Why | PhysTwin]",
            "这段代码把 PhysTwin 一侧的快路径讲得很清楚：同一个 clean rope replay，不只是 benchmark 上看起来快，core source 里也真的有 graph-based replay 这个结构。再加上 Nsight 里 `cudaGraphLaunch_v10000` 占主导，实验、源码和系统证据是对齐的。",
            "[Risk | PhysTwin]",
            "这也不能推出『Newton 只要照抄 graph 就一定追平』。它只支持一个更谨慎的结论：在这个 controlled rope benchmark 里，残余差距更像 runtime organization 的差异，而不只是 controller replay overhead。",
            "",
            "所以这页最后的人话结论是：在 clean rope replay 里，残余 gap 更像 Newton 这边还保留了 many separated launches，而 PhysTwin 这边已经有 graph-based replay path。",
            "但这页也必须明确说不证明什么：它不证明 full physics parity；它不消除 self-collision 或 bunny-contact mismatch 这些别的层次的问题；它也不意味着 collision 在所有场景里都不重要。它只说明在这个 controlled rope replay benchmark 里，controller replay overhead 不是全部答案。",
        ],
    },
    {
        "kind": "body",
        "title": "H6: Optimize Replay Organization Before Render Tuning",
        "bullets": [
            f"**Baseline first:** A1 should replace A0 for viewer diagnosis because A0 -> A1 already saves `{A0_TO_A1:.2f}x`.",
            f"**Render is not first:** on this rope case, viewer ON is only `{A1_TO_E1:.2f}x` slower than the same replay with render OFF.",
            "**Next target:** reduce replay-organization overhead before any kernel micro-optimization.",
        ],
        "transcript": [
            "最后一页只讲这整段 profiling 对 real viewer 到底有什么实际价值。",
            f"第一，A0 到 A1 已经告诉我们：如果目标是 viewer diagnosis，就不应该继续拿更重的 baseline feed path 当默认基线，因为光这一项就已经差了大约 `{A0_TO_A1:.2f}x`。",
            f"第二，E1 对 A1 的比较又告诉我们：在这个 rope case 上，render ON 相对同 replay 的 render OFF 只多了大约 `{A1_TO_E1:.2f}x` wall time，所以如果我们要争取更多 viewer headroom，优先级不应该只放在 rendering。",
            "第三，A1 对 B0、再加上 Newton core staged step path 和 PhysTwin core graph replay path，一起给出的方向是：下一步更值得先研究 replay organization，也就是能不能把 Newton 这条 path 做得更 batched，或者更 graph-like。",
            "第四，这一步做完以后，才值得认真判断 kernel micro-optimization 还有没有足够大的回报。",
            "所以 profiling 这一段最后服务的不是一个抽象 benchmark，而是一个很实际的判断：如果 real viewer 还想要更多余量，我们下一步应该把工程时间先花在 replay path organization 上。",
            "它的边界仍然不变：这里说的是 controlled rope replay benchmark，不是 robot，不是 self-collision，也不是 bunny penetration。",
        ],
    },
    {
        "kind": "body",
        "title": "F1: Force Video Must Show Global Motion And Local Mechanism",
        "bullets": [
            "**Hypothesis:** a local patch alone is not enough; the whole cloth must stay visible.",
            "**Method:** pair a global phenomenon view with a local force zoom.",
            "**Goal:** one force video should explain both motion and contact mechanism.",
        ],
        "transcript": [
            "这里开始进入第三段 penetration analysis。",
            "这一页先把新的 hypothesis 说清楚。",
            "教授要的不是一个只剩局部 probe 的 force patch，也不是一个 static trigger snapshot。",
            "真正需要的是两件事同时成立：第一，看到整块 cloth 相对于 bunny 的整体 penetration 过程；第二，在 local patch 里看到 outward normal、external force、internal force 和 acceleration 的机制解释。",
            "所以这次我把输出正式拆成 phenomenon video 和 force mechanism video，但 force mechanism 本身仍然保留 global panel 加 local zoom panel，而不是只放局部 close-up。",
        ],
    },
    {
        "kind": "grid",
        "title": "F2: Four Phenomenon Videos Show The Full Penetration Process",
        "common_settings": None,
        "items": [
            ("Bunny baseline\nfull process", ACCEPTED_PHENO_GIF["bunny_baseline"]),
            ("Box control\nfull process", ACCEPTED_PHENO_GIF["box_control"]),
            ("Bunny low inertia\nfull process", ACCEPTED_PHENO_GIF["bunny_low_inertia"]),
            ("Bunny larger scale\nfull process", ACCEPTED_PHENO_GIF["bunny_larger_scale"]),
        ],
        "transcript": [
            "这一页只讲 phenomenon，不讲 force。",
            "现在四个 case 都有 accepted 的 global process video，而且都过了 black-screen、motion 和 temporal-density 的 QA。",
            "从这里可以清楚看到四种情况里，pre-contact、first contact、penetration growth 和后续 settle 的整体行为。",
            "也就是说，我们现在已经不再依赖单帧截图来讲 penetration，而是真正有了 meeting-ready 的过程视频。",
        ],
    },
    {
        "kind": "grid",
        "title": "F3: Split Panels Match The Current 2x2 Board",
        "common_settings": None,
        "items": [
            ("Box penalty\nsingle panel", CURRENT_BUNNY_PANEL_GIF["box_penalty"]),
            ("Box total\nsingle panel", CURRENT_BUNNY_PANEL_GIF["box_total"]),
            ("Bunny penalty\nsingle panel", CURRENT_BUNNY_PANEL_GIF["bunny_penalty"]),
            ("Bunny total\nsingle panel", CURRENT_BUNNY_PANEL_GIF["bunny_total"]),
        ],
        "transcript": [
            "这一页现在不再放旧的四个历史 force mechanism 视频了。",
            "这里也把实验设定说清楚：cloth total mass 是 0.1 kg，rigid target mass 是 0.5 kg。",
            "现在 F2 里的四个单视频，都是直接从当前 canonical `2 x 2` board 裁出来的单面板。",
            "也就是说，这一页只负责把四个 panel 单独放大给老师看：box penalty、box total、bunny penalty、bunny total。",
            "这样 F2 和下一页 F3 就不会再出现“这一页还是旧结果、下一页才是新结果”的割裂了。",
        ],
    },
    {
        "kind": "image",
        "title": "F4: The New 2x2 Board Makes Box-vs-Bunny Immediate",
        "note": "TL box penalty, TR box total, BL bunny penalty, BR bunny total.",
        "path": CURRENT_BUNNY_BOARD_GIF,
        "transcript": [
            "这一页就是新的 `2 x 2` board video，我把它直接放进 PPT 里了。",
            "这里必须把实验设定写明：self-collision 是 OFF，cloth total mass 是 0.1 kg，rigid target mass 是 0.5 kg。",
            "它的结构非常直接：左上是 box penalty，右上是 box total，左下是 bunny penalty，右下是 bunny total。",
            "所以这一页的价值不是再讲一个局部 mechanism patch，而是把 box versus bunny、penalty versus total 这两个对比同时压进一个 meeting-readable clip。",
            "如果现场只想停一页讲 penetration，我现在会优先停这一页，因为它最接近老师要的最终 visual comparison。",
        ],
    },
    {
        "kind": "image",
        "title": "F5: 4x Slow Motion Makes Contact Development Easier To Read",
        "note": "Same board as F4, but slowed 4x.",
        "path": CURRENT_BUNNY_SLOW_BOARD_GIF,
        "transcript": [
            "这一页是 F3 的补充版本，不改实验设定，只改播放节奏。",
            "也就是说，self-collision 还是 OFF，cloth total mass 还是 0.1 kg，rigid target mass 还是 0.5 kg。",
            "这里放的是同一套 `2 x 2` board，但是视频整体放慢四倍，而且每个 panel 里都显式写了 `4x slow motion`。",
            "这页的作用不是替代正常速度版本，而是让老师在会议里更容易看清 pre-contact、first contact 和 penetration growth 的过渡。",
        ],
    },
    {
        "kind": "body",
        "title": "S1: Separate Self-Collision Law From Ground-Contact Law",
        "bullets": [
            "**Self-collision axis:** native Newton vs PhysTwin-style law.",
            "**Ground-contact axis:** native Newton vs PhysTwin-style law.",
            "**Fixed scene:** one cloth reference case with one implicit ground plane.",
            "**Takeaway:** analyze the two laws separately.",
        ],
        "transcript": [
            "这里开始我先统一术语，不再说 `phystwin mode` 或者 `Newton way` 这种容易混淆的词。",
            "从现在开始只讲四个术语：native Newton self-collision，PhysTwin-style self-collision，native Newton ground-contact，PhysTwin-style ground-contact。",
            "这样做的目的，是把 self-collision law 和 ground-contact law 明确分开。因为这次真正的问题不是某个 mode 名字，而是这两条 law 分别会怎样影响同一个 cloth reference case 的 RMSE。",
        ],
    },
    {
        "kind": "table",
        "title": "S2: A 2x2 Matrix Isolates The Two Laws",
        "note": "One cloth reference case. One 2x2 law matrix.",
        "title_size": 24,
        "note_font_size": 10,
        "cell_font_size": 13,
        "columns": ["Case", "Self-collision law", "Ground-contact law", "Purpose"],
        "rows": [
            ["case 1", "native OFF", "native Newton", "baseline"],
            ["case 2", "native OFF", "PhysTwin-style", "ground-law effect only"],
            ["case 3", "PhysTwin-style", "native Newton", "self-collision effect only"],
            ["case 4", "PhysTwin-style", "PhysTwin-style", "fully PhysTwin-style pair"],
        ],
        "transcript": [
            "这一页是这次 self-collision 段最重要的结构页。",
            "现在我们不再把问题讲成一个整体 mode，而是做 controlled `2 x 2` matrix：一条轴是 self-collision law，另一条轴是 ground-contact law。",
            "这样同一个 cloth scene、同一个 IR、同一个 dt 和 substeps 下，我们就能回答到底是哪一条 law 改变了 RMSE，以及两条 law 之间有没有 interaction effect。",
        ],
    },
    {
        "kind": "table",
        "title": "S3: Operator Exactness Passes; Rollout Parity Is Still Blocked",
        "note": "Best full-rollout case is mixed: PhysTwin-style self-collision + native ground.",
        "title_size": 24,
        "note_font_size": 10,
        "cell_font_size": 13,
        "columns": ["Evidence", "Result", "Interpretation"],
        "rows": [
            ["Operator-level exactness", "`max_abs_dv = 1.13e-6`, `median_rel_dv = 4.07e-8`", "local PhysTwin-style self-collision operator is already strong"],
            ["Fair 2x2 full rollout", "best case = self PhysTwin-style + ground native Newton", "fully PhysTwin-style pair is not the best full-rollout pair"],
            ["Strict gate", "all four cases remain above `1e-5`", "rollout-level parity is still blocked"],
            ["Current reading", "self-collision axis matters more than ground axis", "remaining gap is interaction-level, not just local contact-law mismatch"],
        ],
        "transcript": [
            "最后这一页只讲当前能够 defend 的结论，不讲过多过程。",
            "第一，operator-level exactness 已经很强，所以我们不能再把问题简单说成 self-collision kernel 还没有对齐。",
            "第二，公平 `2 x 2` full-rollout matrix 已经说明，当前最好的组合不是 fully PhysTwin-style pair，而是 PhysTwin-style self-collision 加 native Newton ground-contact。",
            "第三，四个 case 全部仍然没有过 strict `1e-5` gate。所以当前最准确的表述是：local operator evidence is strong, but rollout-level parity is still blocked。",
            "也就是说，剩余问题更像 whole-step interaction mismatch，而不是 isolated self-collision law mismatch。",
        ],
    },
    {
        "kind": "twocol",
        "title": "R0: Native Robot Release/Drop Is A Sanity Baseline",
        "common_settings": "`drop_release_baseline` | native Franka | semi-implicit rope | real ground | 1:1 video.",
        "left_label": "Authoritative OFF run\nrecoil-fixed stage-0 baseline",
        "left_path": ROBOT_DROP_BASELINE_OFF_GIF,
        "right_label": "Matched ON run\nA/B drag comparison",
        "right_path": ROBOT_DROP_BASELINE_ON_GIF,
        "transcript": [
            "在进入 tabletop push 之前，我先把更窄的 robot stage-0 baseline 单独点出来。",
            "这页不是 final manipulation claim，而是 `drop_release_baseline` 的 sanity baseline：native Newton Franka 在场，rope 先被 support，再 release，然后在 semi-implicit pipeline 里自由下落并撞到 real ground。",
            "当前 authoritative OFF run 的关键数字是：release time 大约 `0.40 s`，first ground contact 大约 `0.7168 s`，impact speed 大约 `3.11 m/s`，early-fall acceleration 拟合约 `-9.80 m/s^2`。",
            "右侧的 drag ON 是 matched A/B。当前 repo 里的结论是 drag 影响是 minor，不是这个 baseline 的主问题来源。",
            "所以这页的作用很窄：它证明 native robot integration、semi-implicit rope free fall、real ground contact 和 1:1 readable video 这四件事已经成立，但不把它说成 full two-way coupling。",
        ],
    },
    {
        "kind": "twocol",
        "title": "R1: Tabletop Push Is A Contact Baseline, Not Full 2-Way Coupling",
        "common_settings": "`tabletop_push_hero` | fixed `sim_dt=5e-5`, `substeps=667` | native finger contact | open-loop robot motion.",
        "left_label": "Hero view\npromoted `c10` contact-fix run",
        "left_path": ROBOT_TABLETOP_HERO_GIF,
        "right_label": "Validation view\nsame promoted run",
        "right_path": ROBOT_TABLETOP_VALIDATION_GIF,
        "transcript": [
            "最后一段是 robotic with deformable objects。",
            "这一章今天的可 defend 结论不再是 release/drop baseline，而是 native tabletop push baseline 已经成立，所以 robot-deformable chapter 至少有了一个更直接的 contact story。",
            "更具体地说，它用 `demo_robot_rope_franka.py` 里的 `tabletop_push_hero` 证明了：native Franka、native tabletop、PhysTwin rope 同时可见，rope 在 visible clip 开始前已经 settle，然后 robot 的 own finger / claw 在桌面高度接近、接触、再推动 rope lateral motion。",
            "但这页也必须明确写清楚：它还不是 full two-way coupling。当前 robot motion 仍然是 demo-side commanded open-loop joint trajectory，rope 会在 Newton 里因为真实接触而响应，但 rope 还不会反过来改变 robot command。",
            "这次比旧版更关键的一点，是 contact-causality 做过修复。新的 promoted `c10` 把 visible first contact 提前，所以“rope 先动、手指后到”的观感明显减弱。",
            "所以这页的结论是，robotic with deformable objects 这一章现在至少有了一个 native tabletop finger-push baseline，可以保守 defend 为 real contact baseline，但不能 overclaim 成 full manipulation 或 full bidirectional coupling。",
        ],
    },
]


def _build_transcript_md(slides: list[dict] | None = None) -> str:
    lines = [
        "# Meeting Transcript — PhysTwin -> Newton Bridge",
        "",
        "语言：中文主讲 + English terminology  ",
        "形式：five-part agenda  ",
        "结构：1. recall  2. performance analysis  3. penetration analysis  4. controlled self-collision / ground-contact laws  5. robotic with deformable objects  ",
        "目标：让 opening framing、transcript framing 和实际 deck 章节结构完全对齐",
        "",
        "---",
        "",
    ]
    active_slides = slides or list(RECALL_SLIDES)
    for idx, slide in enumerate(active_slides, start=1):
        lines.append(f"## Slide {idx} — {slide['title']}")
        for paragraph in slide["transcript"]:
            lines.append(paragraph)
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def _prepare_generated_assets() -> None:
    IMAGE_DIR.mkdir(parents=True, exist_ok=True)
    DECK_GIF_DIR.mkdir(parents=True, exist_ok=True)
    for src, out, width, fps, max_colors in RECALL_DIRECT_GIF_SPECS:
        _ensure_resized_gif(src, out, width=width, fps=fps, max_colors=max_colors)
    _ensure_resized_gif(RECALL_ROPE_GIF_SRC, PERF_ROPE_CASE_GIF, width=720, fps=8, max_colors=96)
    _code_excerpt_image(
        NEWTON_CORE_SPRING_PATH,
        "newton_rope_spring_force",
        _extract_code_segments(
            NEWTON_CORE_SPRING_PATH,
            [(37, 53)],
            highlight_lines={44, 46, 47, 50, 52},
        ),
        CODE_PERF_PHYSICS_NEWTON_PNG,
    )
    _code_excerpt_image(
        PHYSTWIN_SPRING_WARP_CODE_PATH,
        "phystwin_rope_force_update",
        _extract_code_segments(
            PHYSTWIN_SPRING_WARP_CODE_PATH,
            [(121, 132), (156, 160)],
            highlight_lines={123, 129, 132, 157, 159},
        ),
        CODE_PERF_PHYSICS_PHYSTWIN_PNG,
    )
    _code_excerpt_image(
        NEWTON_CORE_SEMIIMPLICIT_PATH,
        "newton_semimplicit_step_path",
        _extract_code_segments(
            NEWTON_CORE_SEMIIMPLICIT_PATH,
            [(141, 145), (160, 165), (172, 177)],
            highlight_lines={142, 145, 161, 173, 177},
        ),
        CODE_PERF_EXECUTION_NEWTON_PNG,
    )
    _code_excerpt_image(
        PHYSTWIN_SPRING_WARP_CODE_PATH,
        "phystwin_graph_replay_path",
        _extract_code_segments(
            PHYSTWIN_SPRING_WARP_CODE_PATH,
            [(769, 782), (800, 802)],
            highlight_lines={769, 772, 780, 782, 802},
        ),
        CODE_PERF_EXECUTION_PHYSTWIN_PNG,
    )
    _system_summary_image(
        PERF_NEWTON_SYSTEM_PNG,
        "System Evidence",
        "results/rope_perf_apples_to_apples/newton/A3_precomputed_attribution + nsight/newton_A1",
        [
            "A3 bridge: 0.037 ms/substep",
            f"A3 post-bridge residual: {A3_POST_BRIDGE_MS:.3f} ms/substep",
            f"A3 collision: {A3_COLLISION_MS:.3f} ms/substep",
            "Nsight API: 77.2% cuLaunchKernel",
            "Interpretation: many staged launches still survive after bridge precompute",
        ],
        font_size=22,
        max_chars=78,
    )
    _system_summary_image(
        PERF_PHYSTWIN_SYSTEM_PNG,
        "System Evidence",
        "results/rope_perf_apples_to_apples/phystwin/B0_headless_throughput + nsight/phystwin_B0",
        [
            "B0 use_graph: True",
            "B0 object_collision_flag: False",
            "B1 simulator_launch: 6.601 ms/frame",
            "Nsight API: 92.6% cudaGraphLaunch_v10000",
            "Interpretation: replay path is graph-driven in the PhysTwin core",
        ],
        font_size=22,
        max_chars=78,
    )
    _code_excerpt_image(
        PHYSTWIN_SPRING_WARP_CODE_PATH,
        "phystwin_object_collision",
        _extract_code_segments(
            PHYSTWIN_SPRING_WARP_CODE_PATH,
            [(196, 200), (203, 205), (212, 220), (294, 296)],
            highlight_lines={196, 203, 205, 220, 295},
        ),
        CODE_SELFCOLLISION_OBJECT_PNG,
    )
    _code_excerpt_image(
        PHYSTWIN_SPRING_WARP_CODE_PATH,
        "phystwin_force_ground",
        _extract_code_segments(
            PHYSTWIN_SPRING_WARP_CODE_PATH,
            [(156, 160), (323, 323), (329, 330), (332, 339), (343, 344)],
            highlight_lines={157, 160, 323, 332, 344},
        ),
        CODE_SELFCOLLISION_FORCE_GROUND_PHYSTWIN_PNG,
    )
    _code_excerpt_image(
        PHYSTWIN_SPRING_WARP_CODE_PATH,
        "phystwin_matched_contact_scope",
        _extract_code_segments(
            PHYSTWIN_SPRING_WARP_CODE_PATH,
            [(260, 268), (294, 296)],
            highlight_lines={261, 268, 294, 295, 296},
        ),
        CODE_SELFCOLLISION_MATCHED_PHYSTWIN_PNG,
    )
    _code_excerpt_image(
        PHYSTWIN_SPRING_WARP_CODE_PATH,
        "phystwin_collision_graph",
        _extract_code_segments(
            PHYSTWIN_SPRING_WARP_CODE_PATH,
            [(928, 930), (934, 939), (246, 247), (255, 257)],
            highlight_lines={928, 930, 246, 255, 256},
        ),
        CODE_SELFCOLLISION_TABLE_PHYSTWIN_PNG,
    )
    bridge_self_contact_code = ROOT / "Newton" / "phystwin_bridge" / "demos" / "self_contact_bridge_kernels.py"
    _code_excerpt_image(
        bridge_self_contact_code,
        "bridge_matched_contact_scope",
        _extract_code_segments(
            bridge_self_contact_code,
            [(537, 540), (542, 544), (549, 553), (556, 557)],
            highlight_lines={540, 542, 544, 553, 557},
        ),
        CODE_SELFCOLLISION_MATCHED_BRIDGE_PNG,
    )
    _code_excerpt_image(
        bridge_self_contact_code,
        "bridge_force_ground",
        _extract_code_segments(
            bridge_self_contact_code,
            [(263, 267), (660, 660), (666, 667), (669, 676), (680, 681)],
            highlight_lines={264, 267, 660, 669, 681},
        ),
        CODE_SELFCOLLISION_FORCE_GROUND_BRIDGE_PNG,
    )
    bridge_phystwin_stack = ROOT / "Newton" / "phystwin_bridge" / "tools" / "core" / "phystwin_contact_stack.py"
    _code_excerpt_image(
        bridge_phystwin_stack,
        "bridge_phystwin_collision_table",
        _extract_code_segments(
            bridge_phystwin_stack,
            [(302, 307), (311, 315), (318, 324)],
            highlight_lines={305, 306, 312, 315, 319},
        ),
        CODE_SELFCOLLISION_TABLE_BRIDGE_PNG,
    )
    _code_excerpt_image(
        PHYSTWIN_SPRING_WARP_CODE_PATH,
        "phystwin_controller_springs",
        _extract_code_segments(
            PHYSTWIN_SPRING_WARP_CODE_PATH,
            [(82, 86), (103, 105), (109, 111)],
            highlight_lines={85, 86, 104, 105, 111},
        ),
        CODE_SELFCOLLISION_CONTROLLER_PHYSTWIN_PNG,
    )
    bridge_import_code = ROOT / "Newton" / "phystwin_bridge" / "tools" / "core" / "newton_import_ir.py"
    _code_excerpt_image(
        bridge_import_code,
        "bridge_controller_particles",
        _extract_code_segments(
            bridge_import_code,
            [(1606, 1610), (1613, 1616), (1619, 1623)],
            highlight_lines={1610, 1613, 1614, 1619, 1620},
        ),
        CODE_SELFCOLLISION_CONTROLLER_BRIDGE_PNG,
    )
    for name, src in ACCEPTED_PHENO_MP4.items():
        _ensure_gif(src, ACCEPTED_PHENO_GIF[name], width=640, fps=8, max_colors=96)
    if CURRENT_BUNNY_BOARD_MP4.exists():
        _ensure_current_bunny_publish_gifs()
        _ensure_current_bunny_slow_board()
        _ensure_gif(
            CURRENT_BUNNY_BOARD_MP4,
            CURRENT_BUNNY_BOARD_GIF,
            max_size_mb=DEFAULT_MAX_GIF_MB,
            quality_profiles=CURRENT_BUNNY_BOARD_GIF_PROFILES,
        )
        if CURRENT_BUNNY_SLOW_BOARD_MP4.exists():
            _ensure_gif(
                CURRENT_BUNNY_SLOW_BOARD_MP4,
                CURRENT_BUNNY_SLOW_BOARD_GIF,
                max_size_mb=DEFAULT_MAX_GIF_MB,
                quality_profiles=CURRENT_BUNNY_SLOW_BOARD_GIF_PROFILES,
            )
        for name, src in CURRENT_BUNNY_PANEL_MP4.items():
            _ensure_gif(
                src,
                CURRENT_BUNNY_PANEL_GIF[name],
                max_size_mb=DEFAULT_MAX_GIF_MB,
                quality_profiles=CURRENT_BUNNY_PANEL_GIF_PROFILES,
            )
    for name, src in ACCEPTED_FORCE_MP4.items():
        _ensure_gif(src, ACCEPTED_FORCE_GIF[name], width=640, fps=8, max_colors=96)
    if SELF_PARITY_SUPPORT_MP4.exists():
        _ensure_gif(SELF_PARITY_SUPPORT_MP4, SELF_PARITY_SUPPORT_GIF, width=640, fps=8, max_colors=96)
    _ensure_gif(ROBOT_DROP_BASELINE_OFF_MP4, ROBOT_DROP_BASELINE_OFF_GIF, width=640, fps=8, max_colors=96)
    _ensure_gif(ROBOT_DROP_BASELINE_ON_MP4, ROBOT_DROP_BASELINE_ON_GIF, width=640, fps=8, max_colors=96)
    if ROBOT_TABLETOP_HERO_MP4.exists():
        _ensure_gif(ROBOT_TABLETOP_HERO_MP4, ROBOT_TABLETOP_HERO_GIF, width=720, fps=8, max_colors=96)
    if ROBOT_TABLETOP_VALIDATION_MP4.exists():
        _ensure_gif(ROBOT_TABLETOP_VALIDATION_MP4, ROBOT_TABLETOP_VALIDATION_GIF, width=720, fps=8, max_colors=96)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the 2026-04-01 recall-only meeting deck.")
    parser.add_argument("--out-dir", type=Path, default=MEETING_DIR)
    parser.add_argument(
        "--out-pptx",
        type=Path,
        default=None,
        help="Optional explicit PPTX output path. Defaults to bridge_meeting_20260401.pptx in --out-dir.",
    )
    parser.add_argument(
        "--slide-range",
        type=str,
        default=None,
        help="Optional 1-based inclusive slide range like 8-12. When omitted, build the full deck.",
    )
    parser.add_argument(
        "--max-pptx-mb",
        type=float,
        default=DEFAULT_MAX_PPTX_MB,
        help="Hard size budget for the generated PPTX in MB. Set <= 0 to disable the size gate.",
    )
    return parser.parse_args()


def _layout(prs: Presentation):
    blank_idx = 6 if len(prs.slide_layouts) > 6 else 0
    return prs.slide_layouts[blank_idx]


def _delete_all_slides(prs: Presentation) -> None:
    sld_id_list = prs.slides._sldIdLst
    for sld_id in list(sld_id_list):
        prs.part.drop_rel(sld_id.rId)
        sld_id_list.remove(sld_id)


def _clear_placeholders(slide) -> None:
    for shape in list(slide.placeholders):
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except Exception:
            pass


def _set_marked_paragraph(
    paragraph,
    text: str,
    *,
    font_name: str,
    font_size: int,
    color: RGBColor,
    accent_color: RGBColor | None = None,
    bold_default: bool = False,
) -> None:
    paragraph.clear()
    parts = text.split("**")
    for idx, part in enumerate(parts):
        if not part:
            continue
        highlighted = idx % 2 == 1
        run = paragraph.add_run()
        run.text = part
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = True if highlighted else bold_default
        run.font.color.rgb = accent_color if highlighted and accent_color else color


def _set_lines(shape, lines: list[str]) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    lines = [line for line in lines if str(line).strip()]
    for i, raw in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        stripped = str(raw).lstrip(" ")
        indent = len(str(raw)) - len(stripped)
        p.level = 1 if indent >= 2 else 0
        _set_marked_paragraph(
            p,
            stripped,
            font_name="Arial",
            font_size=18 if p.level == 0 else 16,
            color=RGBColor(0x22, 0x22, 0x22),
            accent_color=RGBColor(0x1F, 0x4E, 0x79),
        )


def _set_title_textbox(box, title: str, *, size_pt: int = 28) -> None:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    _set_marked_paragraph(
        p,
        title,
        font_name="Arial",
        font_size=size_pt,
        color=RGBColor(0x00, 0x00, 0x00),
        accent_color=RGBColor(0x1F, 0x4E, 0x79),
    )


def _add_label(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: int = 12,
    bold: bool = True,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    _set_marked_paragraph(
        p,
        text,
        font_name="Arial",
        font_size=font_size,
        color=RGBColor(0x22, 0x22, 0x22),
        accent_color=RGBColor(0x1F, 0x4E, 0x79),
        bold_default=bold,
    )


def _fit_box(iw: int, ih: int, bw: int, bh: int) -> tuple[int, int, int, int]:
    scale = min(bw / max(iw, 1), bh / max(ih, 1))
    w, h = max(1, int(iw * scale)), max(1, int(ih * scale))
    return (bw - w) // 2, (bh - h) // 2, w, h


def _load_mono_font(size: int = 18, *, bold: bool = False):
    candidates = (
        [
            "/usr/share/fonts/truetype/cascadia/CascadiaCode-Bold.ttf",
            "/usr/share/fonts/truetype/cascadia/CascadiaMono-Bold.ttf",
            "/usr/share/fonts/truetype/jetbrains-mono/JetBrainsMono-Bold.ttf",
            "/usr/share/fonts/truetype/msttcorefonts/Consolas Bold.ttf",
            "/usr/share/fonts/truetype/noto/NotoSansMono-Bold.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSansMono-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationMono-Bold.ttf",
        ]
        if bold
        else [
            "/usr/share/fonts/truetype/cascadia/CascadiaCode.ttf",
            "/usr/share/fonts/truetype/cascadia/CascadiaMono.ttf",
            "/usr/share/fonts/truetype/jetbrains-mono/JetBrainsMono-Regular.ttf",
            "/usr/share/fonts/truetype/msttcorefonts/Consolas.ttf",
            "/usr/share/fonts/truetype/noto/NotoSansMono-Regular.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf",
        ]
    )
    for p in candidates:
        if Path(p).exists():
            return ImageFont.truetype(p, size=size)
    return ImageFont.load_default()


def _middle_truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    keep_left = max(8, (max_chars - 3) // 2)
    keep_right = max(8, max_chars - keep_left - 3)
    return f"{text[:keep_left]}...{text[-keep_right:]}"


def _trim_code_line(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    if max_chars <= 3:
        return text[:max_chars]
    return f"{text[: max_chars - 3]}..."


def _render_code_png(
    out_path: Path,
    header_title: str,
    header_subtitle: str,
    lines: list[str],
    *,
    font_size: int = 24,
    max_chars: int = 76,
) -> Path:
    panel_bg = (30, 30, 30)
    chrome_bg = (37, 37, 38)
    gutter_bg = (37, 37, 38)
    border = (58, 58, 58)
    text_color = (212, 212, 212)
    path_color = (156, 220, 254)
    gutter_color = (133, 133, 133)
    highlight_fill = (47, 93, 143, 228)
    highlight_accent = (55, 148, 255, 255)
    shadow_color = (0, 0, 0, 110)
    token_colors = {
        Comment: (106, 153, 85),
        Keyword: (86, 156, 214),
        Keyword.Constant: (86, 156, 214),
        Keyword.Namespace: (86, 156, 214),
        Name.Function: (220, 220, 170),
        Name.Class: (78, 201, 176),
        Name.Builtin: (78, 201, 176),
        Name.Builtin.Pseudo: (78, 201, 176),
        Name.Decorator: (197, 134, 192),
        Name.Exception: (78, 201, 176),
        String: (206, 145, 120),
        Number: (181, 206, 168),
        Operator.Word: (86, 156, 214),
        Token.Error: (244, 71, 71),
        Text: text_color,
    }

    font = _load_mono_font(font_size)
    font_bold = _load_mono_font(font_size, bold=True)
    chrome_font = _load_mono_font(max(14, int(font_size * 0.78)), bold=True)
    chrome_meta_font = _load_mono_font(max(12, int(font_size * 0.62)))
    probe = Image.new("RGBA", (64, 64), (0, 0, 0, 0))
    probe_draw = ImageDraw.Draw(probe)
    glyph_box = probe_draw.textbbox((0, 0), "Mg", font=font)
    char_w = max(10, glyph_box[2] - glyph_box[0])
    line_h = max(24, int(font_size * 1.48))

    gutter_w = 118
    panel_pad_x = 34
    panel_pad_bottom = 30
    title_bar_h = 64
    code_top_pad = 24
    max_line_px = 0
    visible_lines = [_trim_code_line(ln, max_chars) for ln in lines]
    for raw in visible_lines:
        highlight = raw.lstrip().startswith(">>>")
        ln = raw.replace(">>>", "", 1).lstrip() if highlight else raw
        stripped = ln.lstrip()
        code = ln
        if stripped and stripped[0].isdigit():
            parts = stripped.split(" ", 1)
            if len(parts) == 2 and parts[0].isdigit():
                leading = len(ln) - len(stripped)
                code = " " * leading + parts[1]
        width_px = int(probe_draw.textlength(code, font=font))
        max_line_px = max(max_line_px, width_px)

    width = min(max(1720, gutter_w + panel_pad_x * 2 + max_line_px + 120), 2100)
    height = min(max(980, title_bar_h + code_top_pad + len(visible_lines) * line_h + panel_pad_bottom + 46), 1480)

    img = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    panel_left, panel_top = 22, 18
    panel_right, panel_bottom = width - 22, height - 18

    shadow = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    shadow_draw = ImageDraw.Draw(shadow, "RGBA")
    shadow_draw.rounded_rectangle(
        [panel_left + 10, panel_top + 14, panel_right + 8, panel_bottom + 14],
        radius=28,
        fill=shadow_color,
    )
    shadow = shadow.filter(ImageFilter.GaussianBlur(radius=14))
    img.alpha_composite(shadow)

    draw = ImageDraw.Draw(img, "RGBA")
    draw.rounded_rectangle(
        [panel_left, panel_top, panel_right, panel_bottom],
        radius=24,
        fill=panel_bg + (255,),
        outline=border + (255,),
        width=2,
    )
    draw.rounded_rectangle(
        [panel_left, panel_top, panel_right, panel_top + title_bar_h],
        radius=24,
        fill=chrome_bg + (255,),
    )
    draw.rectangle(
        [panel_left, panel_top + 26, panel_right, panel_top + title_bar_h],
        fill=chrome_bg + (255,),
    )

    for idx, color in enumerate(((255, 95, 86), (255, 189, 46), (39, 201, 63))):
        cx = panel_left + 28 + idx * 26
        cy = panel_top + title_bar_h // 2
        draw.ellipse([cx - 7, cy - 7, cx + 7, cy + 7], fill=color + (240,))

    display_title = _middle_truncate(header_title, 34)
    display_subtitle = _middle_truncate(header_subtitle, 56)
    draw.text((panel_left + 120, panel_top + 17), display_title, font=chrome_font, fill=(230, 230, 230, 255))
    subtitle_w = int(draw.textlength(display_subtitle, font=chrome_meta_font))
    draw.text((panel_right - 26 - subtitle_w, panel_top + 20), display_subtitle, font=chrome_meta_font, fill=path_color + (225,))

    code_top = panel_top + title_bar_h + code_top_pad
    gutter_left = panel_left + 14
    gutter_right = gutter_left + gutter_w
    draw.rounded_rectangle(
        [gutter_left, code_top - 8, gutter_right, panel_bottom - 18],
        radius=16,
        fill=gutter_bg + (255,),
    )
    code_x = gutter_right + 26

    def _color_for_token(tok):
        for ttype, color in token_colors.items():
            if tok in ttype:
                return color
        return text_color

    def _boost(color: tuple[int, int, int], *, amount: int = 26) -> tuple[int, int, int]:
        return tuple(min(255, channel + amount) for channel in color)

    lexer = PythonLexer(stripnl=False)
    for row_idx, raw_line in enumerate(visible_lines):
        y = code_top + row_idx * line_h
        highlight = raw_line.lstrip().startswith(">>>")
        ln = raw_line.replace(">>>", "", 1).lstrip() if highlight else raw_line
        line_no = ""
        code = ln
        stripped = ln.lstrip()
        leading = len(ln) - len(stripped)
        if stripped and stripped[0].isdigit():
            parts = stripped.split(" ", 1)
            if len(parts) == 2 and parts[0].isdigit():
                line_no = parts[0]
                code = " " * leading + parts[1]
        code = _trim_code_line(code, max_chars)

        row_top = y - 4
        row_bottom = y + line_h - 2
        if highlight:
            draw.rounded_rectangle(
                [panel_left + 14, row_top, panel_right - 14, row_bottom],
                radius=10,
                fill=highlight_fill,
                outline=(118, 179, 245, 255),
                width=2,
            )
            draw.rounded_rectangle(
                [panel_left + 14, row_top, panel_left + 20, row_bottom],
                radius=4,
                fill=highlight_accent,
            )

        if line_no:
            line_no_w = int(draw.textlength(f"{line_no:>4}", font=font))
            draw.text(
                (gutter_right - 18 - line_no_w, y),
                f"{line_no:>4}",
                font=font_bold if highlight else font,
                fill=(_boost(gutter_color, amount=40) if highlight else gutter_color) + (255,),
            )

        x = code_x
        for tok, value in lex(code, lexer):
            clean = value.replace("\n", "")
            if not clean:
                continue
            color = _color_for_token(tok)
            if highlight:
                color = _boost(color, amount=30)
            draw.text((x, y), clean, font=font_bold if highlight else font, fill=color + (255,))
            x += int(draw.textlength(clean, font=font))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)
    return out_path


def _extract_code_segments(path: Path, segments: list[tuple[int, int]], *, highlight_lines: set[int] | None = None) -> list[str]:
    source_lines = path.read_text(encoding="utf-8").splitlines()
    highlight_lines = highlight_lines or set()
    if len(highlight_lines) > 5:
        raise ValueError(
            f"{path}: code excerpt highlight exceeds slide rule ({len(highlight_lines)} > 5)"
        )
    out: list[str] = []
    for seg_i, (start, end) in enumerate(segments):
        for lineno in range(start, end + 1):
            prefix = ">>> " if lineno in highlight_lines else ""
            out.append(f"{prefix}{lineno:4d} {source_lines[lineno - 1]}")
        if seg_i != len(segments) - 1:
            out.append("...")
    return out


def _code_excerpt_image(path: Path, title: str, lines: list[str], out_path: Path) -> Path:
    try:
        rel_path = path.relative_to(ROOT)
    except ValueError:
        rel_path = path
    return _render_code_png(
        out_path,
        path.name,
        str(rel_path),
        lines,
    )


def _system_summary_image(
    out_path: Path,
    title: str,
    subtitle: str,
    lines: list[str],
    *,
    font_size: int = 24,
    max_chars: int = 72,
) -> Path:
    return _render_code_png(
        out_path,
        title,
        subtitle,
        lines,
        font_size=font_size,
        max_chars=max_chars,
    )


def _ensure_transcoded_gif(
    src_path: Path,
    gif_path: Path,
    *,
    width: int = 640,
    fps: int = 8,
    max_colors: int = 96,
    max_size_mb: float | None = None,
    quality_profiles: list[tuple[int, int, int]] | None = None,
) -> Path:
    gif_path.parent.mkdir(parents=True, exist_ok=True)
    profile_rows = [
        [int(w), int(f), int(c)]
        for w, f, c in (quality_profiles or [(width, fps, max_colors)])
    ]
    meta_path = gif_path.with_suffix(f"{gif_path.suffix}.meta.json")
    signature = {
        "profiles": profile_rows,
        "max_size_mb": None if max_size_mb is None else float(max_size_mb),
        "palette_stats_mode": "full",
        "palette_dither": "sierra2_4a",
    }
    if gif_path.exists() and meta_path.exists() and gif_path.stat().st_mtime >= src_path.stat().st_mtime:
        try:
            meta_obj = json.loads(meta_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            meta_obj = {}
        if meta_obj.get("signature") == signature:
            return gif_path
    budget_bytes = None
    if max_size_mb is not None and float(max_size_mb) > 0.0:
        budget_bytes = int(round(float(max_size_mb) * 1_000_000.0))
    candidate_path = gif_path.with_name(f"{gif_path.stem}.tmp{gif_path.suffix}")
    chosen_profile: list[int] | None = None
    chosen_size = 0
    for render_width, render_fps, render_colors in profile_rows:
        subprocess.run(
            [
                "ffmpeg",
                "-y",
                "-loglevel",
                "error",
                "-i",
                str(src_path),
                "-vf",
                (
                    f"fps={int(render_fps)},scale={int(render_width)}:-1:flags=lanczos,"
                    f"split[s0][s1];[s0]palettegen=max_colors={int(render_colors)}:stats_mode=full[p];"
                    "[s1][p]paletteuse=dither=sierra2_4a"
                ),
                "-loop",
                "0",
                str(candidate_path),
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        candidate_size = candidate_path.stat().st_size
        chosen_profile = [int(render_width), int(render_fps), int(render_colors)]
        chosen_size = int(candidate_size)
        if budget_bytes is None or candidate_size <= budget_bytes:
            break
    candidate_path.replace(gif_path)
    meta_path.write_text(
        json.dumps(
            {
                "signature": signature,
                "chosen_profile": chosen_profile,
                "size_bytes": chosen_size,
                "source": str(src_path),
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    return gif_path


def _ensure_gif(
    mp4_path: Path,
    gif_path: Path,
    *,
    width: int = 640,
    fps: int = 8,
    max_colors: int = 96,
    max_size_mb: float | None = None,
    quality_profiles: list[tuple[int, int, int]] | None = None,
) -> Path:
    return _ensure_transcoded_gif(
        mp4_path,
        gif_path,
        width=width,
        fps=fps,
        max_colors=max_colors,
        max_size_mb=max_size_mb,
        quality_profiles=quality_profiles,
    )


def _ensure_current_bunny_panel_mp4s() -> None:
    if not CURRENT_BUNNY_BOARD_SUMMARY.exists():
        raise FileNotFoundError(f"missing bunny board summary: {CURRENT_BUNNY_BOARD_SUMMARY}")
    board_mtime = 0.0
    if CURRENT_BUNNY_BOARD_MP4.exists():
        board_mtime = CURRENT_BUNNY_BOARD_MP4.stat().st_mtime
    expected = list(CURRENT_BUNNY_PANEL_MP4.values()) + list(CURRENT_BUNNY_PANEL_PUBLISH_GIF.values())
    stale_or_missing = [
        path
        for path in expected
        if (not path.exists()) or (board_mtime > 0.0 and path.stat().st_mtime < board_mtime)
    ]
    if not stale_or_missing:
        return
    helper = ROOT / "scripts" / "export_bunny_collision_board_panels.py"
    subprocess.run(
        [
            "python",
            str(helper),
            "--board-summary",
            str(CURRENT_BUNNY_BOARD_SUMMARY),
        ],
        check=True,
    )


def _ensure_current_bunny_publish_gifs() -> None:
    if not CURRENT_BUNNY_BOARD_MP4.exists():
        return
    _ensure_gif(
        CURRENT_BUNNY_BOARD_MP4,
        CURRENT_BUNNY_BOARD_PUBLISH_GIF,
        max_size_mb=DEFAULT_MAX_GIF_MB,
        quality_profiles=CURRENT_BUNNY_BOARD_PUBLISH_GIF_PROFILES,
    )
    _ensure_current_bunny_panel_mp4s()
    for name, src in CURRENT_BUNNY_PANEL_MP4.items():
        _ensure_gif(
            src,
            CURRENT_BUNNY_PANEL_PUBLISH_GIF[name],
            max_size_mb=DEFAULT_MAX_GIF_MB,
            quality_profiles=CURRENT_BUNNY_PANEL_PUBLISH_GIF_PROFILES,
        )


def _ensure_current_bunny_slow_board() -> None:
    if not CURRENT_BUNNY_BOARD_SUMMARY.exists():
        return
    main_summary = json.loads(CURRENT_BUNNY_BOARD_SUMMARY.read_text(encoding="utf-8"))
    box_summary = Path(str(main_summary["cases"]["box_control"]["detector_summary_path"])).expanduser().resolve()
    bunny_summary = Path(str(main_summary["cases"]["bunny_baseline"]["detector_summary_path"])).expanduser().resolve()
    slow_inputs = [CURRENT_BUNNY_BOARD_SUMMARY, box_summary, bunny_summary, RENDER_BUNNY_BOARD_PATH]
    newest_input_mtime = max(path.stat().st_mtime for path in slow_inputs if path.exists())
    slow_outputs = [CURRENT_BUNNY_SLOW_BOARD_MP4, CURRENT_BUNNY_SLOW_BOARD_PUBLISH_GIF, CURRENT_BUNNY_SLOW_BOARD_SUMMARY]
    needs_render = any((not path.exists()) or path.stat().st_mtime < newest_input_mtime for path in slow_outputs)
    if needs_render:
        CURRENT_BUNNY_SLOW_BOARD_DIR.mkdir(parents=True, exist_ok=True)
        subprocess.run(
            [
                "python",
                str(RENDER_BUNNY_BOARD_PATH),
                "--box-summary",
                str(box_summary),
                "--bunny-summary",
                str(bunny_summary),
                "--out-dir",
                str(CURRENT_BUNNY_SLOW_BOARD_DIR),
                "--output-stem",
                CURRENT_BUNNY_SLOW_BOARD_STEM,
                "--slowdown-factor",
                "4",
                "--playback-label",
                "4x slow motion",
            ],
            check=True,
        )
    elif CURRENT_BUNNY_SLOW_BOARD_MP4.exists():
        _ensure_gif(
            CURRENT_BUNNY_SLOW_BOARD_MP4,
            CURRENT_BUNNY_SLOW_BOARD_PUBLISH_GIF,
            max_size_mb=DEFAULT_MAX_GIF_MB,
            quality_profiles=CURRENT_BUNNY_SLOW_BOARD_PUBLISH_GIF_PROFILES,
        )


def _ensure_resized_gif(
    src_path: Path,
    out_path: Path,
    *,
    width: int = 720,
    fps: int = 8,
    max_colors: int = 96,
) -> Path:
    return _ensure_transcoded_gif(src_path, out_path, width=width, fps=fps, max_colors=max_colors)


def _mb(byte_count: int) -> float:
    return float(byte_count) / 1_000_000.0


def _validate_pptx_size(pptx_path: Path, *, max_pptx_mb: float) -> None:
    if max_pptx_mb <= 0.0:
        return
    actual_mb = _mb(pptx_path.stat().st_size)
    if actual_mb > max_pptx_mb:
        raise RuntimeError(
            f"PPTX size gate failed: {pptx_path} is {actual_mb:.1f} MB, above the {max_pptx_mb:.1f} MB budget."
        )


def _add_pic(slide, path: Path, left: int, top: int, width: int, height: int):
    with Image.open(path) as img:
        dx, dy, w, h = _fit_box(img.size[0], img.size[1], width, height)
    return slide.shapes.add_picture(str(path), left + dx, top + dy, width=w, height=h)


def _title_slide(prs: Presentation, title: str, sub: list[str]) -> None:
    slide = prs.slides.add_slide(_layout(prs))
    _clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(
        REF_TITLE_LEFT, REF_TITLE_SLIDE_TOP, REF_TITLE_W, REF_TITLE_SLIDE_H
    )
    _set_title_textbox(title_box, title, size_pt=28)
    body_box = slide.shapes.add_textbox(
        REF_BODY_LEFT, REF_SUBTITLE_TOP, REF_BODY_W, REF_SUBTITLE_H
    )
    _set_lines(body_box, sub)


def _body(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(_layout(prs))
    _clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(
        REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H
    )
    _set_title_textbox(title_box, title, size_pt=28)
    body_box = slide.shapes.add_textbox(
        REF_BODY_LEFT, REF_BODY_TOP, REF_BODY_W, REF_BODY_H
    )
    _set_lines(body_box, bullets)


def _table(
    prs: Presentation,
    title: str,
    columns: list[str],
    rows: list[list[str]],
    note: str | None = None,
    *,
    title_size: int = 28,
    note_font_size: int = 11,
    cell_font_size: int = 15,
) -> None:
    slide = prs.slides.add_slide(_layout(prs))
    _clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(
        REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H
    )
    _set_title_textbox(title_box, title, size_pt=title_size)
    if note:
        _add_label(
            slide,
            TABLE_LEFT,
            TABLE_NOTE_TOP,
            TABLE_W,
            TABLE_NOTE_H,
            note,
            font_size=note_font_size,
            bold=False,
        )
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(columns),
        TABLE_LEFT,
        TABLE_TOP,
        TABLE_W,
        TABLE_H,
    )
    table = table_shape.table
    if len(columns) == 4:
        widths = [
            int(TABLE_W * 0.14),
            int(TABLE_W * 0.24),
            int(TABLE_W * 0.18),
            int(TABLE_W * 0.44),
        ]
    else:
        widths = [
            int(TABLE_W * 0.12),
            int(TABLE_W * 0.32),
            int(TABLE_W * 0.18),
            int(TABLE_W * 0.14),
            int(TABLE_W * 0.24),
        ]
    for idx, width in enumerate(widths[: len(columns)]):
        table.columns[idx].width = width

    def set_cell_text(
        cell,
        text: str,
        *,
        bold: bool,
        align_center: bool = False,
        fill_color: RGBColor | None = None,
    ) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color or RGBColor(0xFF, 0xFF, 0xFF)
        cell.text_frame.clear()
        cell.text_frame.word_wrap = True
        p = cell.text_frame.paragraphs[0]
        _set_marked_paragraph(
            p,
            text,
            font_name="Arial",
            font_size=cell_font_size,
            color=RGBColor(0x22, 0x22, 0x22),
            accent_color=RGBColor(0x1F, 0x4E, 0x79),
            bold_default=bold,
        )
        if align_center:
            p.alignment = 1

    for c, heading in enumerate(columns):
        header = table.cell(0, c)
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(0xD9, 0xE2, 0xF3)
        set_cell_text(header, heading, bold=True, align_center=True)

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            fill = None
            if len(columns) == 4 and c == 3:
                if "1.00x" in str(value):
                    fill = RGBColor(0xE2, 0xF0, 0xD9)
                else:
                    fill = RGBColor(0xF4, 0xCC, 0xCC)
            elif len(columns) == 4 and c == 2:
                try:
                    fill = RGBColor(0xE2, 0xF0, 0xD9) if float(value) >= 1.0 else RGBColor(0xFF, 0xF2, 0xCC)
                except Exception:
                    fill = None
            set_cell_text(
                table.cell(r, c),
                str(value),
                bold=(c == 0),
                align_center=(c != 0),
                fill_color=fill,
            )


def _table_with_gif(
    prs: Presentation,
    title: str,
    columns: list[str],
    rows: list[list[str]],
    gif_path: Path,
    gif_label: str,
    *,
    note: str | None = None,
) -> None:
    slide = prs.slides.add_slide(_layout(prs))
    _clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(title_box, title, size_pt=28)
    if note:
        _add_label(slide, 420000, 1180000, 5600000, 220000, note, font_size=11, bold=False)
    _add_label(slide, 6670000, 1500000, 1900000, 220000, gif_label, font_size=11, bold=False)
    _add_pic(slide, gif_path, 6420000, 1780000, 2250000, 1380000)

    table_left = 420000
    table_top = 1540000
    table_w = 5750000
    table_h = 3150000
    table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), table_left, table_top, table_w, table_h)
    table = table_shape.table
    widths = [
        int(table_w * 0.14),
        int(table_w * 0.34),
        int(table_w * 0.16),
        int(table_w * 0.12),
        int(table_w * 0.24),
    ]
    for idx, width in enumerate(widths[: len(columns)]):
        table.columns[idx].width = width

    def set_cell_text(
        cell,
        text: str,
        *,
        bold: bool,
        align_center: bool = False,
        fill_color: RGBColor | None = None,
    ) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color or RGBColor(0xFF, 0xFF, 0xFF)
        cell.text_frame.clear()
        cell.text_frame.word_wrap = True
        p = cell.text_frame.paragraphs[0]
        _set_marked_paragraph(
            p,
            text,
            font_name="Arial",
            font_size=13,
            color=RGBColor(0x22, 0x22, 0x22),
            accent_color=RGBColor(0x1F, 0x4E, 0x79),
            bold_default=bold,
        )
        if align_center:
            p.alignment = 1

    for c, heading in enumerate(columns):
        header = table.cell(0, c)
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(0xD9, 0xE2, 0xF3)
        set_cell_text(header, heading, bold=True, align_center=True)

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            set_cell_text(
                table.cell(r, c),
                str(value),
                bold=(c == 0),
                align_center=(c != 0),
            )


def _gif_grid_2x2(
    prs: Presentation,
    title: str,
    items: list[tuple[str, Path]],
    *,
    common_settings: str | None = None,
) -> None:
    slide = prs.slides.add_slide(_layout(prs))
    _clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(
        REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H
    )
    _set_title_textbox(title_box, title, size_pt=28)
    if common_settings:
        _add_label(
            slide,
            REF_GRID_COMMON_LEFT,
            REF_GRID_COMMON_TOP,
            REF_GRID_COMMON_W,
            REF_GRID_COMMON_H,
            common_settings,
            font_size=11,
            bold=False,
        )
    label_w = 3000000
    label_h = 250000
    pic_w = 2350000
    pic_h = 1320000
    label_positions = [
        (520000, 1700000),
        (4720000, 1700000),
        (520000, 3380000),
        (4720000, 3380000),
    ]
    pic_positions = [
        (1280000, 1980000),
        (5780000, 1980000),
        (1280000, 3660000),
        (5780000, 3660000),
    ]
    for idx, (label, path) in enumerate(items[:4]):
        lx, ly = label_positions[idx]
        px, py = pic_positions[idx]
        _add_label(slide, lx, ly, label_w, label_h, label, font_size=11, bold=False)
        _add_pic(slide, path, px, py, pic_w, pic_h)


def _gif_twocol(
    prs: Presentation,
    title: str,
    left_label: str,
    left_path: Path,
    right_label: str,
    right_path: Path,
    *,
    common_settings: str | None = None,
) -> None:
    slide = prs.slides.add_slide(_layout(prs))
    _clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(
        REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H
    )
    _set_title_textbox(title_box, title, size_pt=28)
    if common_settings:
        _add_label(
            slide,
            REF_GRID_COMMON_LEFT,
            REF_GRID_COMMON_TOP,
            REF_GRID_COMMON_W,
            REF_GRID_COMMON_H,
            common_settings,
            font_size=11,
            bold=False,
        )
    label_y = 1730000
    pic_y = 2010000
    label_w = 3300000
    label_h = 250000
    pic_w = 3000000
    pic_h = 1760000
    left_label_x = 520000
    right_label_x = 4670000
    left_pic_x = 920000
    right_pic_x = 5090000
    _add_label(
        slide,
        left_label_x,
        label_y,
        label_w,
        label_h,
        left_label,
        font_size=11,
        bold=False,
    )
    _add_pic(
        slide,
        left_path,
        left_pic_x,
        pic_y,
        pic_w,
        pic_h,
    )
    _add_label(
        slide,
        right_label_x,
        label_y,
        label_w,
        label_h,
        right_label,
        font_size=11,
        bold=False,
    )
    _add_pic(
        slide,
        right_path,
        right_pic_x,
        pic_y,
        pic_w,
        pic_h,
    )


def _perf_gif_bullets(
    prs: Presentation,
    title: str,
    gif_label: str,
    gif_path: Path,
    bullets: list[str],
    *,
    note: str | None = None,
) -> None:
    slide = prs.slides.add_slide(_layout(prs))
    _clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(title_box, title, size_pt=28)
    if note:
        _add_label(
            slide,
            REF_GRID_COMMON_LEFT,
            REF_GRID_COMMON_TOP,
            REF_GRID_COMMON_W,
            REF_GRID_COMMON_H,
            note,
            font_size=12,
            bold=False,
        )
    _add_label(slide, 540000, 1620000, 3920000, 260000, gif_label, font_size=11, bold=False)
    _add_pic(slide, gif_path, 540000, 1930000, 3920000, 2360000)
    body_box = slide.shapes.add_textbox(4540000, 1680000, 3800000, 2580000)
    _set_lines(body_box, bullets)


def _code_twocol_large(
    prs: Presentation,
    title: str,
    left_label: str,
    left_path: Path,
    right_label: str,
    right_path: Path,
    *,
    note: str | None = None,
) -> None:
    slide = prs.slides.add_slide(_layout(prs))
    _clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(title_box, title, size_pt=24)
    if note:
        _add_label(
            slide,
            REF_GRID_COMMON_LEFT,
            REF_GRID_COMMON_TOP,
            REF_GRID_COMMON_W,
            REF_GRID_COMMON_H,
            note,
            font_size=10,
            bold=False,
        )
    left_x = 420000
    right_x = 4550000
    label_y = 1560000
    pic_y = 1815000
    pic_w = 3900000
    pic_h = 2950000
    _add_label(slide, left_x, label_y, pic_w, 200000, left_label, font_size=10, bold=False)
    _add_pic(slide, left_path, left_x, pic_y, pic_w, pic_h)
    _add_label(slide, right_x, label_y, pic_w, 200000, right_label, font_size=10, bold=False)
    _add_pic(slide, right_path, right_x, pic_y, pic_w, pic_h)


def _image_slide(prs: Presentation, title: str, path: Path, *, note: str | None = None) -> None:
    slide = prs.slides.add_slide(_layout(prs))
    _clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(
        REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H
    )
    _set_title_textbox(title_box, title, size_pt=28)
    if note:
        _add_label(
            slide,
            REF_GRID_COMMON_LEFT,
            REF_GRID_COMMON_TOP,
            REF_GRID_COMMON_W,
            REF_GRID_COMMON_H,
            note,
            font_size=12,
            bold=False,
        )
    pic_left = 420000
    pic_top = 1650000
    pic_w = 8550000
    pic_h = 4300000
    _add_pic(slide, path, pic_left, pic_top, pic_w, pic_h)


def _parse_slide_range(spec: str | None) -> list[dict]:
    if not spec:
        return list(RECALL_SLIDES)
    match = re.fullmatch(r"\s*(\d+)\s*-\s*(\d+)\s*", spec)
    if not match:
        raise ValueError(f"Unsupported --slide-range format: {spec!r}. Expected START-END, e.g. 8-12.")
    start = int(match.group(1))
    end = int(match.group(2))
    if start < 1 or end < start or end > len(RECALL_SLIDES):
        raise ValueError(
            f"Slide range {spec!r} is out of bounds for a {len(RECALL_SLIDES)}-slide deck."
        )
    return list(RECALL_SLIDES[start - 1 : end])


def build_recall_only_deck(prs: Presentation, slides: list[dict] | None = None) -> None:
    _delete_all_slides(prs)
    active_slides = slides or list(RECALL_SLIDES)
    for slide in active_slides:
        kind = slide["kind"]
        if kind == "title":
            _title_slide(prs, slide["title"], slide["subtitle"])
        elif kind == "grid":
            _gif_grid_2x2(
                prs,
                slide["title"],
                slide["items"],
                common_settings=slide.get("common_settings"),
            )
        elif kind == "twocol":
            _gif_twocol(
                prs,
                slide["title"],
                slide["left_label"],
                slide["left_path"],
                slide["right_label"],
                slide["right_path"],
                common_settings=slide.get("common_settings"),
            )
        elif kind == "code_twocol":
            _code_twocol_large(
                prs,
                slide["title"],
                slide["left_label"],
                slide["left_path"],
                slide["right_label"],
                slide["right_path"],
                note=slide.get("common_settings") or slide.get("note"),
            )
        elif kind == "code_twocol_large":
            _code_twocol_large(
                prs,
                slide["title"],
                slide["left_label"],
                slide["left_path"],
                slide["right_label"],
                slide["right_path"],
                note=slide.get("note"),
            )
        elif kind == "perf_gif_bullets":
            _perf_gif_bullets(
                prs,
                slide["title"],
                slide["gif_label"],
                slide["gif_path"],
                slide["bullets"],
                note=slide.get("note"),
            )
        elif kind == "body":
            _body(prs, slide["title"], slide["bullets"])
        elif kind == "table":
            _table(
                prs,
                slide["title"],
                slide["columns"],
                slide["rows"],
                note=slide.get("note"),
                title_size=slide.get("title_size", 28),
                note_font_size=slide.get("note_font_size", 11),
                cell_font_size=slide.get("cell_font_size", 15),
            )
        elif kind == "table_gif":
            _table_with_gif(
                prs,
                slide["title"],
                slide["columns"],
                slide["rows"],
                slide["gif_path"],
                slide["gif_label"],
                note=slide.get("note"),
            )
        elif kind == "image":
            _image_slide(
                prs,
                slide["title"],
                slide["path"],
                note=slide.get("note"),
            )
        else:
            raise ValueError(f"Unknown recall slide kind: {kind}")


def _markdown_to_pdf(md_text: str, html_path: Path, pdf_path: Path) -> None:
    html_body = markdown.markdown(md_text, extensions=["tables", "fenced_code"])
    html = f"""<!doctype html>
<html>
<head>
  <meta charset="utf-8">
  <style>
    body {{ font-family: 'DejaVu Sans', sans-serif; line-height: 1.55; margin: 36px 48px; color: #222; }}
    h1, h2, h3 {{ color: #1f3d5a; }}
    code {{ background: #f2f4f7; padding: 1px 4px; border-radius: 4px; }}
    pre {{ background: #f2f4f7; padding: 10px; overflow-x: auto; }}
    blockquote {{ color: #555; border-left: 4px solid #ccd5e0; padding-left: 12px; }}
    table {{ border-collapse: collapse; width: 100%; }}
    th, td {{ border: 1px solid #ccd5e0; padding: 6px 8px; }}
    a {{ color: #1f5d99; text-decoration: none; }}
  </style>
</head>
<body>{html_body}</body>
</html>"""
    html_path.write_text(html, encoding="utf-8")
    try:
        from weasyprint import HTML

        HTML(filename=str(html_path)).write_pdf(str(pdf_path))
        return
    except Exception:
        pass

    subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_path.parent),
            str(html_path),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def main() -> int:
    args = parse_args()
    out_dir = args.out_dir.resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    _prepare_generated_assets()
    active_slides = _parse_slide_range(args.slide_range)

    if not TEMPLATE_PPTX.exists():
        raise FileNotFoundError(f"Missing local template: {TEMPLATE_PPTX}")

    prs = Presentation(str(TEMPLATE_PPTX))
    build_recall_only_deck(prs, slides=active_slides)
    out_pptx = args.out_pptx.resolve() if args.out_pptx else (out_dir / OUT_PPTX.name)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))

    transcript_text = _build_transcript_md(active_slides)
    transcript_md = out_dir / "transcript.md"
    transcript_md.write_text(transcript_text, encoding="utf-8")
    transcript_html = out_dir / "transcript.html"
    transcript_pdf = out_dir / "transcript.pdf"
    _markdown_to_pdf(transcript_text, transcript_html, transcript_pdf)
    _validate_pptx_size(out_pptx, max_pptx_mb=float(args.max_pptx_mb))

    print(f"PPTX: {out_pptx}")
    print(f"PPTX size: {_mb(out_pptx.stat().st_size):.1f} MB (budget {float(args.max_pptx_mb):.1f} MB)")
    print(f"Transcript MD: {transcript_md}")
    print(f"Transcript PDF: {transcript_pdf}")
    print(f"Slides: {len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
