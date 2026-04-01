#!/usr/bin/env python3
from __future__ import annotations

import argparse
import csv
import json
import re
import shutil
import statistics
import subprocess
import textwrap
from datetime import date
from pathlib import Path

import markdown
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont
from pygments import lex
from pygments.lexers import PythonLexer
from pygments.token import (
    Comment,
    Keyword,
    Name,
    Number,
    Operator,
    Punctuation,
    String,
    Text,
    Token,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

MEETING_DIR = Path(__file__).resolve().parent
ROOT = MEETING_DIR.parents[1]
LOCAL_TEMPLATE_DIR = MEETING_DIR / "templates"
MEETING_TEMPLATE_PPTX = LOCAL_TEMPLATE_DIR / "My Adjust.pptx"
SLIDE_W = 9144000
SLIDE_H = 5143500
TODAY_TAG = date.today().strftime("%Y%m%d")
FORMAL_SLIDE_DIR = MEETING_DIR
FORMAL_SLIDE_TAG = "20260325"

REF_TITLE_LEFT = 311760
REF_TITLE_TOP = 444960
REF_TITLE_W = 8520120
REF_TITLE_H = 572400

REF_BODY_LEFT = 311760
REF_BODY_TOP = 1152360
REF_BODY_W = 8520120
REF_BODY_H = 3416040

REF_SECTION_TOP = 2151000
REF_SECTION_H = 841320

REF_TITLE_SLIDE_TOP = 744480
REF_TITLE_SLIDE_H = 2052360
REF_SUBTITLE_TOP = 2834280
REF_SUBTITLE_H = 792360

REF_MEDIA_PIC_LEFT = 1394640
REF_MEDIA_PIC_TOP = 1152360
REF_MEDIA_PIC_W = 6354000
REF_MEDIA_PIC_H = 2117880
REF_MEDIA_FOOTER_TOP = 3373200
REF_MEDIA_FOOTER_H = 639720

REF_SUMMARY_BOX_LEFT = 482040
REF_SUMMARY_BOX_TOP = 1323360
REF_SUMMARY_BOX_W = 8179560
REF_SUMMARY_BOX_H = 3074400
REF_SUMMARY_TEXT_LEFT = 683280
REF_SUMMARY_TEXT_TOP = 1433160
REF_SUMMARY_TEXT_W = 7777080
REF_SUMMARY_TEXT_H = 2132640

REF_GRID_COMMON_LEFT = 482040
REF_GRID_COMMON_TOP = 1170720
REF_GRID_COMMON_W = 8179560
REF_GRID_COMMON_H = 273960
REF_GRID_LABEL_W = 3983040
REF_GRID_LABEL_H = 243720
REF_GRID_PIC_W = 1735200
REF_GRID_PIC_H = 975960
REF_GRID_TL_LABEL_LEFT = 482040
REF_GRID_TR_LABEL_LEFT = 4678560
REF_GRID_TOP_LABEL_TOP = 1706760
REF_GRID_BOT_LABEL_TOP = 3129480
REF_GRID_TL_PIC_LEFT = 1605960
REF_GRID_TR_PIC_LEFT = 5802480
REF_GRID_TOP_PIC_TOP = 1999440
REF_GRID_BOT_PIC_TOP = 3421872

GRID_CSV = (
    ROOT / "tmp" / "cloth_bunny_off_grid_20260317" / "aggregate_penetration_metrics.csv"
)
BEST_OFF_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s3p0_h0p1_bm5000"
    / "self_off"
    / "cloth_bunny_drop_off_m5000.mp4"
)
BASELINE_OFF_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s1p0_h0p1_bm5000"
    / "self_off"
    / "cloth_bunny_drop_off_m5000.mp4"
)
MASS_LIGHT_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s1p0_h0p1_bm1"
    / "self_off"
    / "cloth_bunny_drop_off_m1.mp4"
)
MASS_MID1_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s1p0_h0p1_bm10"
    / "self_off"
    / "cloth_bunny_drop_off_m10.mp4"
)
MASS_MID2_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s1p0_h0p1_bm1000"
    / "self_off"
    / "cloth_bunny_drop_off_m1000.mp4"
)
MASS_HEAVY_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s1p0_h0p1_bm5000"
    / "self_off"
    / "cloth_bunny_drop_off_m5000.mp4"
)
HEIGHT_LOW_MP4 = BASELINE_OFF_MP4
HEIGHT_MID1_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s1p0_h0p25_bm5000"
    / "self_off"
    / "cloth_bunny_drop_off_m5000.mp4"
)
HEIGHT_MID2_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s1p0_h0p5_bm5000"
    / "self_off"
    / "cloth_bunny_drop_off_m5000.mp4"
)
HEIGHT_HIGH_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s1p0_h1p0_bm5000"
    / "self_off"
    / "cloth_bunny_drop_off_m5000.mp4"
)
CCD_SMALL_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s0p5_h0p1_bm5000"
    / "self_off"
    / "cloth_bunny_drop_off_m5000.mp4"
)
CCD_MID1_MP4 = BASELINE_OFF_MP4
CCD_MID2_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_grid_20260317"
    / "ccd_s2p0_h0p1_bm5000"
    / "self_off"
    / "cloth_bunny_drop_off_m5000.mp4"
)
CCD_LARGE_MP4 = BEST_OFF_MP4
MPM_SLOTH_MP4 = (
    ROOT
    / "tmp"
    / "sloth_oldrender_vis0005_full_300_20260317"
    / "sloth_mpm_sand_one_way_wide.mp4"
)
CURRENT_ROPE_BUNNY_MP4 = (
    ROOT
    / "tmp"
    / "rope_bunny_meeting_full_close_20260318"
    / "rope_bunny_drop_m5.mp4"
)
ROPE_BUNNY_FULL_M5_MP4 = (
    ROOT
    / "tmp"
    / "rope_bunny_close_1m_manualcam_20260325"
    / "m5"
    / "rope_bunny_drop_m5.mp4"
)
ROPE_BUNNY_FULL_M500_MP4 = (
    ROOT
    / "tmp"
    / "rope_bunny_close_1m_manualcam_20260325"
    / "m500"
    / "rope_bunny_drop_m500.mp4"
)
ROPE_BUNNY_WEIGHT_1KG_MP4 = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_rope_bunny_weight_compare_20260325"
    / "rope_total_1kg"
    / "rope_bunny_total1kg_m5.mp4"
)
ROPE_BUNNY_WEIGHT_5KG_MP4 = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_rope_bunny_weight_compare_20260325"
    / "rope_total_5kg"
    / "rope_bunny_total5kg_m5.mp4"
)
ROPE_BUNNY_WEIGHT_1KG_SUMMARY = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_rope_bunny_weight_compare_20260325"
    / "rope_total_1kg"
    / "rope_bunny_total1kg_m5_summary.json"
)
ROPE_BUNNY_WEIGHT_5KG_SUMMARY = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_rope_bunny_weight_compare_20260325"
    / "rope_total_5kg"
    / "rope_bunny_total5kg_m5_summary.json"
)
CLOTH_BUNNY_WEIGHT_1KG_MP4 = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_weight_compare_20260325"
    / "cloth_total_1kg"
    / "self_off"
    / "cloth_bunny_total1kg_off_m5.mp4"
)
CLOTH_BUNNY_WEIGHT_5KG_MP4 = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_weight_compare_20260325"
    / "cloth_total_5kg"
    / "self_off"
    / "cloth_bunny_total5kg_off_m5.mp4"
)
CLOTH_BUNNY_WEIGHT_1KG_SUMMARY = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_weight_compare_20260325"
    / "cloth_total_1kg"
    / "self_off"
    / "cloth_bunny_total1kg_off_m5_summary.json"
)
CLOTH_BUNNY_WEIGHT_5KG_SUMMARY = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_weight_compare_20260325"
    / "cloth_total_5kg"
    / "self_off"
    / "cloth_bunny_total5kg_off_m5_summary.json"
)
CLOTH_BUNNY_WEIGHT_0P1KG_MP4 = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_lightweight_compare_20260325"
    / "cloth_total_0p1kg"
    / "self_off"
    / "cloth_bunny_total0p1kg_off_m5.mp4"
)
CLOTH_BUNNY_WEIGHT_0P5KG_MP4 = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_lightweight_compare_20260325"
    / "cloth_total_0p5kg"
    / "self_off"
    / "cloth_bunny_total0p5kg_off_m5.mp4"
)
CLOTH_BUNNY_WEIGHT_0P1KG_SUMMARY = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_lightweight_compare_20260325"
    / "cloth_total_0p1kg"
    / "self_off"
    / "cloth_bunny_total0p1kg_off_m5_summary.json"
)
CLOTH_BUNNY_WEIGHT_0P5KG_SUMMARY = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_lightweight_compare_20260325"
    / "cloth_total_0p5kg"
    / "self_off"
    / "cloth_bunny_total0p5kg_off_m5_summary.json"
)
CLOTH_BUNNY_RADIUS_COMPARE_DIR = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_radius_compare_20260325"
)
CLOTH_BUNNY_SIZE_COMPARE_DIR = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_size_compare_20260325"
)
CLOTH_BUNNY_DT_COMPARE_DIR = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_dt_compare_20260325"
)
CLOTH_BUNNY_ULTRALIGHT_COMPARE_DIR = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_ultralight_compare_20260325"
)
TWO_ROPES_WEIGHT_COMPARE_DIR = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_two_ropes_weight_compare_20260325"
)
TWO_ROPES_BOX_STACK_MP4 = (
    ROOT
    / "tmp"
    / "two_ropes_boxshape50_fullvideo_20260325"
    / "rope_box_stack_drop.mp4"
)
TWO_ROPES_BOX_STACK_SUMMARY = (
    ROOT
    / "tmp"
    / "two_ropes_boxshape50_fullvideo_20260325"
    / "rope_box_stack_drop_summary.json"
)
THREE_BOX_STACK_MP4 = (
    ROOT
    / "tmp"
    / "three_box_stack_fullvideo_20260325"
    / "rope_box_stack_drop.mp4"
)
THREE_BOX_STACK_SUMMARY = (
    ROOT
    / "tmp"
    / "three_box_stack_fullvideo_20260325"
    / "rope_box_stack_drop_summary.json"
)
ROBOT_ROPE_FRANKA_MP4 = (
    ROOT
    / "results"
    / "robot_deformable_demo"
    / "runs"
    / "20260330_213045_native_franka_lift_release_presentation"
    / "media"
    / "final.mp4"
)
ROBOT_ROPE_FRANKA_SUMMARY = (
    ROOT
    / "results"
    / "robot_deformable_demo"
    / "runs"
    / "20260330_213045_native_franka_lift_release_presentation"
    / "summary.json"
)
CLOTH_RIGID_COMPARE_BUNNY_MP4 = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_rigid_compare_20260325"
    / "bunny"
    / "self_off"
    / "cloth_rigid_compare_bunny_off_m5.mp4"
)
CLOTH_RIGID_COMPARE_BOX_MP4 = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_rigid_compare_20260325"
    / "box"
    / "self_off"
    / "cloth_rigid_compare_box_off_m5.mp4"
)
CLOTH_RIGID_COMPARE_BUNNY_SUMMARY = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_rigid_compare_20260325"
    / "bunny"
    / "self_off"
    / "cloth_rigid_compare_bunny_off_m5_summary.json"
)
CLOTH_RIGID_COMPARE_BOX_SUMMARY = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_rigid_compare_20260325"
    / "box"
    / "self_off"
    / "cloth_rigid_compare_box_off_m5_summary.json"
)
CURRENT_CLOTH_ON_MP4 = (
    ROOT
    / "tmp"
    / "cloth_on_hypothesis_20260318"
    / "h6_long_horizon"
    / "on1"
    / "self_on"
    / "cloth_bunny_drop_on_m5000.mp4"
)
CURRENT_CLOTH_OFF_MP4 = (
    ROOT
    / "tmp"
    / "cloth_on_hypothesis_20260318"
    / "h1_radius_scale"
    / "off_baseline"
    / "self_off"
    / "cloth_bunny_drop_off_m5000.mp4"
)
P0_BUNNY_OFF_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_decoupled_render_20260321"
    / "self_off"
    / "cloth_bunny_drop_off_m0p5.mp4"
)
H1_BUNNY_RADIUS_2X_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_off_center_contact2x_render_20260322"
    / "self_off"
    / "cloth_bunny_drop_off_off_m0p5.mp4"
)
P0_BUNNY_FINAL_MP4 = (
    ROOT
    / "tmp"
    / "bunny_mesh_render_compare_20260321"
    / "bodycenter_contact8x"
    / "self_off"
    / "cloth_bunny_drop_off_m0p5.mp4"
)
P0_BUNNY_FINAL_SUMMARY = (
    ROOT
    / "tmp"
    / "bunny_mesh_render_compare_20260321"
    / "bodycenter_contact8x"
    / "self_off"
    / "cloth_bunny_drop_off_m0p5_summary.json"
)
P0_BOX_OFF_MP4 = (
    ROOT
    / "tmp"
    / "cloth_rigid_box_final_20260321"
    / "self_off"
    / "cloth_rigid_drop_off_m0p5.mp4"
)
P0_BOX_BASELINE_MP4 = (
    ROOT
    / "tmp"
    / "cloth_rigid_box_baseline_render_20260322"
    / "self_off"
    / "cloth_rigid_drop_off_m0p5.mp4"
)
P0_BOX_OFF_SUMMARY = (
    ROOT
    / "tmp"
    / "cloth_rigid_box_final_20260321"
    / "self_off"
    / "cloth_rigid_drop_off_m0p5_summary.json"
)
P0_BUNNY_SWEEP_SUMMARY = (
    ROOT / "tmp" / "bunny_mesh_local_sweep_20260321" / "summary_final.md"
)
P0_BUNNY_SWEEP_CSV = (
    ROOT / "tmp" / "bunny_mesh_local_sweep_20260321" / "summary_final.csv"
)
ADVISOR_BM0P1_MP4 = (
    ROOT
    / " result_for_slides"
    / "h1c_alpha_contact_bm0p1"
    / "self_off"
    / "cloth_bunny_drop_off_m0p1.mp4"
)
ADVISOR_BM0P5_MP4 = (
    ROOT
    / " result_for_slides"
    / "h1c_alpha_contact_bm0p5"
    / "self_off"
    / "cloth_bunny_drop_off_m0p5.mp4"
)
H2_ON_FAILURE_MP4 = (
    ROOT
    / "tmp"
    / "cloth_bunny_on_contactdist_render_20260321"
    / "s0p5"
    / "self_on"
    / "cloth_bunny_drop_on_m0p5.mp4"
)
THIN_EAR_COUNTEREXAMPLE_DIR = (
    ROOT
    / "formal_slide"
    / "meeting_2026_03_25"
    / "exp_cloth_bunny_thin_ear_particle_radius_5x_10x_20260324"
    / "self_off"
)
THIN_EAR_5X_MP4 = THIN_EAR_COUNTEREXAMPLE_DIR / "cloth_bunny_closeup_ccd5x_off_m0p5.mp4"
THIN_EAR_10X_MP4 = THIN_EAR_COUNTEREXAMPLE_DIR / "cloth_bunny_closeup_ccd10x_off_m0p5.mp4"
THIN_EAR_5X_SUMMARY = THIN_EAR_COUNTEREXAMPLE_DIR / "cloth_bunny_closeup_ccd5x_off_m0p5_summary.json"
THIN_EAR_10X_SUMMARY = THIN_EAR_COUNTEREXAMPLE_DIR / "cloth_bunny_closeup_ccd10x_off_m0p5_summary.json"
FORCE_DIAG_BOX_DIR = ROOT / "tmp" / "force_diag_box_smoke_20260327" / "self_off" / "force_diagnostic"
FORCE_DIAG_BOX_SUMMARY = FORCE_DIAG_BOX_DIR / "force_diag_trigger_summary.json"
FORCE_DIAG_BOX_SNAPSHOT = FORCE_DIAG_BOX_DIR / "force_diag_trigger_snapshot.png"
ROPE_PROFILE_FULL_JSON = ROOT / "tmp" / "rope_control_profile_full_20260328" / "rope_control_realtime_profile.json"

PHYSTWIN_SIM = (
    ROOT / "PhysTwin" / "qqtt" / "model" / "diff_simulator" / "spring_mass_warp.py"
)
CLOTH_BUNNY_COMMON = (
    ROOT / "Newton" / "phystwin_bridge" / "demos" / "demo_cloth_bunny_common.py"
)
NEWTON_KERNELS_CONTACT = (
    ROOT
    / "Newton"
    / "newton"
    / "newton"
    / "_src"
    / "solvers"
    / "semi_implicit"
    / "kernels_contact.py"
)
NEWTON_GEOM_KERNELS = (
    ROOT / "Newton" / "newton" / "newton" / "_src" / "geometry" / "kernels.py"
)
NEWTON_IMPORT_IR = (
    ROOT / "Newton" / "phystwin_bridge" / "tools" / "core" / "newton_import_ir.py"
)
NEWTON_KERNELS_PARTICLE = (
    ROOT
    / "Newton"
    / "newton"
    / "newton"
    / "_src"
    / "solvers"
    / "semi_implicit"
    / "kernels_particle.py"
)
NEWTON_SOLVER_BASE = (
    ROOT / "Newton" / "newton" / "newton" / "_src" / "solvers" / "solver.py"
)

FALLBACK_CODE_LINE_MAP: dict[Path, dict[int, str]] = {
    NEWTON_GEOM_KERNELS: {
        1013: "def create_soft_contacts(model, state_in, state_out, margin, soft_contacts):",
        1014: "    # broad phase + narrow phase for particle-shape contact generation",
        1112: "        max_dist = margin + radius / min_scale",
        1113: "        face_index, face_u, face_v, sign = wp.mesh_query_point_sign_normal(",
        1114: "            shape_geo.source, x_local, max_dist, face_index, face_u, face_v, sign",
        1115: "        )",
        1132: "        if d < margin + radius:",
        1133: "            soft_contacts[soft_contact_count] = (particle_i, body_i, bx, bn, d)",
        1134: "            soft_contact_count += 1",
    },
    NEWTON_KERNELS_CONTACT: {
        79: "if not model.particle_has_self_collision:",
        80: "    return",
        83: "radius_i = particle_radius[i]",
        84: "radius_j = particle_radius[index]",
        85: "err = d - radius_i - radius_j",
        88: "if err <= soft_contact_margin:",
        90: "    n = wp.normalize(delta)",
        94: "    particle_f[i] += n * err * model.soft_contact_ke",
        215: "radius = model.particle_radius[particle_index]",
        216: "c = wp.dot(n, px - bx) - radius",
        217: "ke = 0.5 * (model.soft_contact_ke + shape_material_ke[shape_index])",
        218: "kd = 0.5 * (model.soft_contact_kd + shape_material_kd[shape_index])",
        219: "kf = 0.5 * (model.soft_contact_kf + shape_material_kf[shape_index])",
        220: "mu = 0.5 * (model.soft_contact_mu + shape_material_mu[shape_index])",
        221: "fn = n * c * ke",
        222: "fd = n * wp.min(vn, 0.0) * kd",
        223: "ft = vt * wp.min(kf, mu * wp.length(fn + fd) / (wp.length(vt) + 1.0e-6))",
        244: "particle_f[particle_index] -= (fn + fd + ft)",
        245: "body_f[body_index] += fn + fd + ft",
        246: "body_t[body_index] += wp.cross(bx - body_q[body_index], fn + fd + ft)",
        247: "if body_index >= 0:",
        248: "    body_contact_count[body_index] += 1",
        263: "particle_f[particle_index] += contact_force",
        265: "body_f[body_index] -= contact_force",
        267: "body_t[body_index] -= wp.cross(r, contact_force)",
        269: "particle_f[particle_index] -= total_contact",
        271: "body_f[body_index] += total_contact",
        273: "body_t[body_index] += wp.cross(r, total_contact)",
    },
}

LY_TITLE = 0
LY_SECTION = 1
LY_BODY = 2
LY_TWO_COL = 3


def _layout(prs: Presentation, idx: int):
    blank_idx = 6 if len(prs.slide_layouts) > 6 else 0
    return prs.slide_layouts[blank_idx]


def _default_body_rect(slide):
    return REF_BODY_LEFT, REF_BODY_TOP, REF_BODY_W, REF_BODY_H


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Generate meeting slides + transcript.")
    p.add_argument(
        "--out-dir",
        type=Path,
        default=FORMAL_SLIDE_DIR,
        help=f"Compatibility-only flag. Slides are always written to {FORMAL_SLIDE_DIR}.",
    )
    args = p.parse_args()

    requested = args.out_dir.expanduser()
    if requested.is_absolute():
        requested = requested.resolve()
    else:
        requested = (Path.cwd() / requested).resolve()

    formal = FORMAL_SLIDE_DIR.resolve()
    if requested != formal:
        raise SystemExit(f"Slides must be exported to {formal}; refusing --out-dir {requested}")

    args.out_dir = formal
    return args


def _delete_all_slides(prs: Presentation) -> None:
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        prs.part.drop_rel(sldId.rId)
        sldIdLst.remove(sldId)


def _ph(slide, idx: int):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    raise KeyError(f"Placeholder idx={idx} missing")


def _clear_placeholders(slide) -> None:
    for shape in list(slide.placeholders):
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except Exception:
            pass


def _set_lines(shape, lines) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    lines = [l for l in lines if str(l).strip()]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        raw = str(line)
        stripped = raw.lstrip(" ")
        indent = len(raw) - len(stripped)
        p.level = 1 if indent >= 2 else 0
        p.alignment = PP_ALIGN.LEFT
        _set_marked_paragraph(
            p,
            stripped,
            font_name="Arial",
            font_size=18 if p.level == 0 else 16,
            color=RGBColor(0x22, 0x22, 0x22),
            accent_color=RGBColor(0x1F, 0x4E, 0x79),
        )


def _set_marked_paragraph(
    paragraph,
    text: str,
    *,
    font_name: str,
    font_size: int,
    color: RGBColor,
    accent_color: RGBColor | None = None,
    bold_default: bool = False,
) -> None:
    paragraph.clear()
    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if not part:
            continue
        highlighted = part.startswith("**") and part.endswith("**")
        content = part[2:-2] if highlighted else part
        run = paragraph.add_run()
        run.text = content
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = True if highlighted else bold_default
        run.font.color.rgb = accent_color if highlighted and accent_color else color


def _set_title_textbox(box, title: str, *, size_pt: int = 28) -> None:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    plain_title = re.sub(r"\*\*(.*?)\*\*", r"\1", title)
    effective_size = size_pt
    if size_pt >= 36:
        if len(plain_title) > 40:
            effective_size = 30
        if len(plain_title) > 55:
            effective_size = 24
        if len(plain_title) > 70:
            effective_size = 20
    else:
        if len(plain_title) > 42:
            effective_size = 24
        if len(plain_title) > 56:
            effective_size = 20
        if len(plain_title) > 72:
            effective_size = 18
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_marked_paragraph(
        p,
        title,
        font_name="Arial",
        font_size=effective_size,
        color=RGBColor(0x00, 0x00, 0x00),
        accent_color=RGBColor(0x1F, 0x4E, 0x79),
    )


def _title_slide(prs: Presentation, title: str, sub: list[str]) -> None:
    s = prs.slides.add_slide(_layout(prs, LY_TITLE))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_SLIDE_TOP, REF_TITLE_W, REF_TITLE_SLIDE_H)
    _set_title_textbox(box, title, size_pt=28)
    box = s.shapes.add_textbox(REF_BODY_LEFT, REF_SUBTITLE_TOP, REF_BODY_W, REF_SUBTITLE_H)
    _set_lines(box, sub)


def _section(prs: Presentation, heading: str) -> None:
    s = prs.slides.add_slide(_layout(prs, LY_SECTION))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_SECTION_TOP, REF_TITLE_W, REF_SECTION_H)
    _set_title_textbox(box, heading, size_pt=36)


def _body(prs: Presentation, title: str, bullets: list[str]):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=28)
    l, t, w, h = _default_body_rect(s)
    target = s.shapes.add_textbox(l, t, w, h)
    _set_lines(target, bullets)
    return s


def _twocol(prs: Presentation, title: str, left: list[str], right: list[str]):
    s = prs.slides.add_slide(_layout(prs, LY_TWO_COL))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=28)
    l, t, w, h = _default_body_rect(s)
    left_target = s.shapes.add_textbox(l, t, int((w - Inches(0.3)) / 2), h)
    right_target = s.shapes.add_textbox(
        l + int((w + Inches(0.3)) / 2), t, int((w - Inches(0.3)) / 2), h
    )
    _set_lines(left_target, left)
    _set_lines(right_target, right)
    return s


def _load_rows(csv_path: Path) -> list[dict]:
    rows: list[dict] = []
    with csv_path.open(newline="", encoding="utf-8") as fh:
        reader = csv.DictReader(fh)
        for row in reader:
            for k in (
                "contact_collision_dist_used",
                "contact_dist_scale",
                "drop_height_m",
                "rigid_mass",
                "total_object_mass",
                "mass_ratio",
                "body_speed_max",
                "min_clearance_min_to_bunny_top",
                "min_clearance_mean_to_bunny_top",
                "min_clearance_max_to_bunny_top",
                "final_clearance_min_to_bunny_top",
                "final_clearance_mean_to_bunny_top",
                "final_clearance_max_to_bunny_top",
                "dynamic_bunny_top_final_z",
            ):
                row[k] = float(row[k])
            for k in (
                "first_min_below_top_frame",
                "first_mean_below_top_frame",
                "first_all_below_top_frame",
            ):
                row[k] = None if row[k] in ("", "None", None) else int(row[k])
            rows.append(row)
    return rows


def _ensure_gif(
    mp4: Path, gif_dir: Path, max_mb: int = 40, stem_override: str | None = None
) -> Path:
    gif_name = f"{stem_override}.gif" if stem_override else mp4.with_suffix(".gif").name
    gif = gif_dir / gif_name
    max_bytes = max_mb * 1024 * 1024
    if gif.exists() and not mp4.exists():
        return gif
    if (
        gif.exists()
        and gif.stat().st_size < max_bytes
        and mp4.exists()
        and gif.stat().st_mtime >= mp4.stat().st_mtime
    ):
        return gif
    gif.parent.mkdir(parents=True, exist_ok=True)
    for fps, w, colors in [(10, 800, 192), (8, 720, 128), (6, 640, 96), (5, 560, 64)]:
        vf = (
            f"fps={fps},scale={w}:-1:flags=lanczos,split[s0][s1];"
            f"[s0]palettegen=max_colors={colors}:stats_mode=diff[p];"
            f"[s1][p]paletteuse=dither=bayer:bayer_scale=4"
        )
        subprocess.run(
            [
                "ffmpeg",
                "-y",
                "-loglevel",
                "error",
                "-i",
                str(mp4),
                "-vf",
                vf,
                "-loop",
                "0",
                str(gif),
            ],
            check=True,
        )
        if gif.exists() and gif.stat().st_size < max_bytes:
            return gif
    raise RuntimeError(f"Failed to compress {mp4} under {max_mb}MB")


def _fit_box(iw: int, ih: int, bw: int, bh: int) -> tuple[int, int, int, int]:
    s = min(bw / max(iw, 1), bh / max(ih, 1))
    w, h = max(1, int(iw * s)), max(1, int(ih * s))
    return (bw - w) // 2, (bh - h) // 2, w, h


def _add_pic(slide, path: Path, left: int, top: int, width: int, height: int):
    img = Image.open(path)
    dx, dy, w, h = _fit_box(img.size[0], img.size[1], width, height)
    return slide.shapes.add_picture(str(path), left + dx, top + dy, width=w, height=h)


def _add_label(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: int = 12,
    bold: bool = True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_marked_paragraph(
        p,
        text,
        font_name="Arial",
        font_size=font_size,
        color=RGBColor(0x22, 0x22, 0x22),
        accent_color=RGBColor(0x1F, 0x4E, 0x79),
        bold_default=bold,
    )


def _detach_body_placeholder(slide):
    try:
        body = _ph(slide, 1)
        left, top, width, height = (
            int(body.left),
            int(body.top),
            int(body.width),
            int(body.height),
        )
        sp = body.element
        sp.getparent().remove(sp)
        return left, top, width, height
    except KeyError:
        return _default_body_rect(slide)


def _gif_twocol(
    prs: Presentation,
    title: str,
    ll: str,
    lg: Path,
    rl: str,
    rg: Path,
    common_settings: str | None = None,
):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)
    common_h = REF_GRID_COMMON_H if common_settings else 0
    lh = REF_GRID_LABEL_H
    cw = REF_GRID_LABEL_W
    ch = REF_MEDIA_PIC_H
    y0 = REF_GRID_TOP_LABEL_TOP
    x0 = REF_GRID_TL_LABEL_LEFT
    x1 = REF_GRID_TR_LABEL_LEFT
    if common_settings:
        _add_label(
            s,
            REF_GRID_COMMON_LEFT,
            REF_GRID_COMMON_TOP,
            REF_GRID_COMMON_W,
            common_h,
            f"Common settings: {common_settings}",
            font_size=12,
            bold=False,
        )
    _add_label(s, x0, y0, cw, lh, ll, font_size=11)
    _add_pic(s, lg, REF_GRID_TL_PIC_LEFT, REF_GRID_TOP_PIC_TOP, REF_GRID_PIC_W, REF_GRID_PIC_H)
    _add_label(s, x1, y0, cw, lh, rl, font_size=11)
    _add_pic(s, rg, REF_GRID_TR_PIC_LEFT, REF_GRID_TOP_PIC_TOP, REF_GRID_PIC_W, REF_GRID_PIC_H)
    return s


def _gif_twocol_large(
    prs: Presentation,
    title: str,
    ll: str,
    lg: Path,
    rl: str,
    rg: Path,
    common_settings: str | None = None,
):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)

    gap = int(Inches(0.22))
    common_h = int(Inches(0.34)) if common_settings else 0
    label_h = int(Inches(0.42))
    pad_y = int(Inches(0.04))
    pic_w = int((width - gap) / 2)
    x0 = left
    x1 = left + pic_w + gap
    label_y = top + common_h
    pic_y = label_y + label_h + pad_y
    pic_h = max(1, height - common_h - label_h - pad_y)

    if common_settings:
        _add_label(
            s,
            left,
            top,
            width,
            common_h,
            f"Common settings: {common_settings}",
            font_size=12,
            bold=False,
        )
    _add_label(s, x0, label_y, pic_w, label_h, ll, font_size=12, bold=True)
    _add_pic(s, lg, x0, pic_y, pic_w, pic_h)
    _add_label(s, x1, label_y, pic_w, label_h, rl, font_size=12, bold=True)
    _add_pic(s, rg, x1, pic_y, pic_w, pic_h)
    return s


def _gif_grid_2x2(
    prs: Presentation,
    title: str,
    items: list[tuple[str, Path]],
    common_settings: str | None = None,
):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)
    label_h = REF_GRID_LABEL_H
    common_h = REF_GRID_COMMON_H if common_settings else 0
    if common_settings:
        _add_label(
            s,
            REF_GRID_COMMON_LEFT,
            REF_GRID_COMMON_TOP,
            REF_GRID_COMMON_W,
            common_h,
            f"Common settings: {common_settings}",
            font_size=12,
            bold=False,
        )
    label_positions = [
        (REF_GRID_TL_LABEL_LEFT, REF_GRID_TOP_LABEL_TOP),
        (REF_GRID_TR_LABEL_LEFT, REF_GRID_TOP_LABEL_TOP),
        (REF_GRID_TL_LABEL_LEFT, REF_GRID_BOT_LABEL_TOP),
        (REF_GRID_TR_LABEL_LEFT, REF_GRID_BOT_LABEL_TOP),
    ]
    pic_positions = [
        (REF_GRID_TL_PIC_LEFT, REF_GRID_TOP_PIC_TOP),
        (REF_GRID_TR_PIC_LEFT, REF_GRID_TOP_PIC_TOP),
        (REF_GRID_TL_PIC_LEFT, REF_GRID_BOT_PIC_TOP),
        (REF_GRID_TR_PIC_LEFT, REF_GRID_BOT_PIC_TOP),
    ]
    for idx, (label, gif) in enumerate(items[:4]):
        lx, ly = label_positions[idx]
        px, py = pic_positions[idx]
        _add_label(s, lx, ly, REF_GRID_LABEL_W, label_h, label, font_size=10, bold=False)
        _add_pic(s, gif, px, py, REF_GRID_PIC_W, REF_GRID_PIC_H)
    return s


def _load_mono_font(size: int = 18, *, bold: bool = False):
    candidates = (
        [
            "/usr/share/fonts/truetype/dejavu/DejaVuSansMono-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationMono-Bold.ttf",
        ]
        if bold
        else [
            "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf",
        ]
    )
    for p in candidates:
        if Path(p).exists():
            return ImageFont.truetype(p, size=size)
    return ImageFont.load_default()


def _load_sans_font(size: int = 22, *, bold: bool = False):
    candidates = (
        [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        ]
        if bold
        else [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        ]
    )
    for p in candidates:
        if Path(p).exists():
            return ImageFont.truetype(p, size=size)
    return ImageFont.load_default()


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> list[str]:
    words = text.split()
    if not words:
        return [""]
    lines: list[str] = []
    cur = words[0]
    for word in words[1:]:
        trial = cur + " " + word
        if draw.textlength(trial, font=font) <= max_width:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    lines.append(cur)
    return lines


def _draw_round_box(draw, xy, fill, outline, radius=22, width=3):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def _draw_boxed_text(draw, xy, title: str, body: list[str], title_font, body_font, *, fill, outline, title_fill, body_fill):
    _draw_round_box(draw, xy, fill, outline)
    x0, y0, x1, y1 = xy
    pad_x = 18
    pad_y = 14
    draw.text((x0 + pad_x, y0 + pad_y), title, font=title_font, fill=title_fill)
    y = y0 + pad_y + 34
    max_width = int(x1 - x0 - 2 * pad_x)
    for line in body:
        wrapped = _wrap_text(draw, line, body_font, max_width)
        for seg in wrapped:
            draw.text((x0 + pad_x, y), seg, font=body_font, fill=body_fill)
            y += 22
        y += 2


def _draw_arrow(draw, start, end, fill, width=5, head=12):
    x0, y0 = start
    x1, y1 = end
    draw.line([x0, y0, x1, y1], fill=fill, width=width)
    import math
    ang = math.atan2(y1 - y0, x1 - x0)
    left = (x1 - head * math.cos(ang - math.pi / 6), y1 - head * math.sin(ang - math.pi / 6))
    right = (x1 - head * math.cos(ang + math.pi / 6), y1 - head * math.sin(ang + math.pi / 6))
    draw.polygon([(x1, y1), left, right], fill=fill)


def _text_height(font, extra: int = 0) -> int:
    bbox = font.getbbox("Ag")
    return (bbox[3] - bbox[1]) + extra


def _draw_wrapped_lines(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    lines: list[str],
    font,
    fill,
    max_width: int,
    *,
    line_gap: int = 6,
) -> int:
    x, y = xy
    line_h = _text_height(font, line_gap)
    for line in lines:
        wrapped = _wrap_text(draw, line, font, max_width)
        for seg in wrapped:
            draw.text((x, y), seg, font=font, fill=fill)
            y += line_h
    return y


def _draw_centered_label(draw, center, text: str, font, fill):
    bbox = draw.textbbox((0, 0), text, font=font)
    x = int(center[0] - (bbox[2] - bbox[0]) / 2)
    y = int(center[1] - (bbox[3] - bbox[1]) / 2)
    draw.text((x, y), text, font=font, fill=fill)


def _draw_dashed_line(draw, start, end, fill, *, width: int = 4, dash: int = 16, gap: int = 12):
    import math

    x0, y0 = start
    x1, y1 = end
    length = math.hypot(x1 - x0, y1 - y0)
    if length == 0:
        return
    dx = (x1 - x0) / length
    dy = (y1 - y0) / length
    pos = 0.0
    while pos < length:
        seg_end = min(length, pos + dash)
        draw.line(
            [
                x0 + dx * pos,
                y0 + dy * pos,
                x0 + dx * seg_end,
                y0 + dy * seg_end,
            ],
            fill=fill,
            width=width,
        )
        pos += dash + gap


def _make_bridge_pipeline_diagram(out_path: Path) -> Path:
    W, H = 2400, 1180
    img = Image.new("RGB", (W, H), (248, 250, 252))
    draw = ImageDraw.Draw(img)
    title_font = _load_sans_font(38, bold=True)
    sec_font = _load_sans_font(28, bold=True)
    body_font = _load_sans_font(21, bold=False)
    small_font = _load_sans_font(18, bold=False)
    chip_font = _load_sans_font(22, bold=True)
    arrow_font = _load_sans_font(16, bold=True)

    navy = (24, 58, 96)
    slate = (71, 85, 105)
    line = (181, 194, 209)
    blue = (235, 244, 255)
    green = (236, 253, 245)
    amber = (255, 247, 237)
    purple = (245, 243, 255)
    mint = (236, 253, 245)
    cyan = (236, 254, 255)
    peach = (255, 237, 213)
    white = (255, 255, 255)

    draw.text((70, 36), "PhysTwin -> Newton -> Interaction Pipeline", font=title_font, fill=navy)
    draw.text(
        (70, 88),
        "Hypothesis H1: the bridge rebuilds a real Newton object instead of replaying a baked animation.",
        font=small_font,
        fill=slate,
    )

    # Main pipeline row
    phys_box = (70, 180, 560, 470)
    ir_box = (650, 180, 1140, 470)
    recon_box = (1230, 180, 1720, 470)
    interact_box = (1810, 180, 2300, 470)

    _draw_boxed_text(
        draw,
        phys_box,
        "1. PhysTwin",
        [
            "topology",
            "mass + springs",
            "contact params",
        ],
        sec_font,
        body_font,
        fill=blue,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    _draw_boxed_text(
        draw,
        ir_box,
        "2. Bridge IR",
        [
            "same semantics",
            "single source of truth",
            "no baked motion",
        ],
        sec_font,
        body_font,
        fill=green,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    _draw_boxed_text(
        draw,
        recon_box,
        "3. Newton Rebuild",
        [
            "particles + springs",
            "rigid body + mesh",
            "native solver state",
        ],
        sec_font,
        body_font,
        fill=amber,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    _draw_boxed_text(
        draw,
        interact_box,
        "4. Interaction",
        [
            "deformable <-> rigid",
            "real force exchange",
            "not replay only",
        ],
        sec_font,
        body_font,
        fill=cyan,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )

    _draw_arrow(draw, (560, 325), (650, 325), navy, width=6, head=18)
    _draw_arrow(draw, (1140, 325), (1230, 325), navy, width=6, head=18)
    _draw_arrow(draw, (1720, 325), (1810, 325), navy, width=6, head=18)
    _draw_centered_label(draw, (605, 292), "export_ir.py", arrow_font, navy)
    _draw_centered_label(draw, (1185, 292), "import", arrow_font, navy)
    _draw_centered_label(draw, (1765, 292), "contact", arrow_font, navy)

    draw.text((70, 560), "Newton side: native worlds stay separate", font=sec_font, fill=navy)

    spring_box = (120, 635, 660, 860)
    rigid_box = (930, 635, 1470, 860)
    mpm_box = (1740, 635, 2280, 860)

    _draw_boxed_text(
        draw,
        spring_box,
        "Spring-Mass",
        [
            "SemiImplicit",
            "rope / cloth / sloth / zebra",
        ],
        sec_font,
        body_font,
        fill=white,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    _draw_boxed_text(
        draw,
        rigid_box,
        "Rigid",
        [
            "body + shape",
            "box / bunny / plane",
        ],
        sec_font,
        body_font,
        fill=white,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    _draw_boxed_text(
        draw,
        mpm_box,
        "MPM",
        [
            "implicit MPM",
            "sand / water / granular",
        ],
        sec_font,
        body_font,
        fill=white,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )

    draw.text((70, 930), "Key claim", font=sec_font, fill=navy)
    inter_cards = [
        (
            (120, 990, 800, 1100),
            mint,
            "Spring-mass <-> rigid uses native particle-shape contact",
            [
                "This is the bridge path we validate first.",
            ],
        ),
        (
            (860, 990, 1540, 1100),
            cyan,
            "Rigid <-> MPM is a different native path",
            [
                "So MPM questions are separate from H1.",
            ],
        ),
        (
            (1600, 990, 2280, 1100),
            peach,
            "Demo-side proxy exists only for spring-mass -> MPM",
            [
                "It is not the same claim as deformable-rigid H1.",
            ],
        ),
    ]
    for box_xy, fill, title, body in inter_cards:
        _draw_boxed_text(
            draw,
            box_xy,
            title,
            body,
            chip_font,
            small_font,
            fill=fill,
            outline=line,
            title_fill=navy,
            body_fill=(31, 41, 55),
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)
    return out_path


def _make_body_shape_mesh_diagram(out_path: Path) -> Path:
    W, H = 2200, 1080
    img = Image.new("RGB", (W, H), (248, 250, 252))
    draw = ImageDraw.Draw(img)
    title_font = _load_sans_font(38, bold=True)
    sec_font = _load_sans_font(28, bold=True)
    body_font = _load_sans_font(21, bold=False)
    small_font = _load_sans_font(18, bold=False)

    navy = (24, 58, 96)
    slate = (71, 85, 105)
    line = (181, 194, 209)
    blue = (235, 244, 255)
    amber = (255, 247, 237)
    mint = (236, 253, 245)
    purple = (245, 243, 255)

    draw.text((70, 34), "How Bunny Contact Works in Newton", font=title_font, fill=navy)
    draw.text(
        (70, 86),
        "Hypothesis H1: particles hit the bunny mesh, and the rigid state receives the reaction.",
        font=small_font,
        fill=slate,
    )

    body_box = (70, 180, 760, 520)
    shape_box = (1440, 180, 2130, 520)
    contact_box = (860, 660, 1340, 930)
    particle_box = (160, 700, 560, 930)

    _draw_boxed_text(
        draw,
        body_box,
        "1. Rigid state",
        [
            "mass + inertia",
            "pose + velocity",
            "integrated in time",
        ],
        sec_font,
        body_font,
        fill=blue,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    _draw_boxed_text(
        draw,
        shape_box,
        "2. Bunny mesh",
        [
            "closest point",
            "signed normal",
            "contact geometry only",
        ],
        sec_font,
        body_font,
        fill=amber,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    _draw_boxed_text(
        draw,
        contact_box,
        "3. Contact",
        [
            "query -> penetration",
            "elastic + damping + friction",
        ],
        sec_font,
        body_font,
        fill=mint,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    _draw_boxed_text(
        draw,
        particle_box,
        "Deformable particles",
        [
            "cloth / rope nodes",
            "SemiImplicit state",
        ],
        _load_sans_font(24, bold=True),
        _load_sans_font(19, bold=False),
        fill=purple,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )

    # Particle dots
    particle_dots = [
        (220, 800), (270, 760), (320, 810), (375, 770), (425, 825), (480, 780),
        (245, 860), (300, 905), (360, 865), (420, 900), (485, 850),
    ]
    for x, y in particle_dots:
        draw.ellipse((x - 12, y - 12, x + 12, y + 12), fill=(59, 130, 246), outline=(37, 99, 235), width=2)

    # Mini shape icons
    draw.rounded_rectangle((1590, 320, 1725, 430), radius=14, outline=(180, 83, 9), fill=(255, 237, 213), width=3)
    _draw_centered_label(draw, (1657, 375), "box", small_font, navy)
    mesh_pts = [(1795, 420), (1850, 320), (1920, 385), (1975, 295), (2045, 400), (1980, 470), (1865, 475)]
    draw.polygon(mesh_pts, outline=(124, 45, 18), fill=(254, 215, 170))
    draw.line(mesh_pts + [mesh_pts[0]], fill=(124, 45, 18), width=3)
    _draw_centered_label(draw, (1925, 505), "mesh", small_font, navy)

    # Contact arrows
    _draw_arrow(draw, (560, 810), (860, 795), (39, 113, 73), width=7, head=18)
    _draw_centered_label(draw, (700, 735), "particle -> mesh query", small_font, (22, 101, 52))

    _draw_arrow(draw, (1440, 630), (1330, 730), (217, 119, 6), width=7, head=18)
    _draw_centered_label(draw, (1560, 610), "mesh supplies normal", small_font, (180, 83, 9))

    _draw_arrow(draw, (1100, 660), (760, 520), (147, 51, 234), width=7, head=18)
    _draw_centered_label(draw, (900, 585), "reaction -> rigid motion", small_font, (126, 34, 206))

    takeaway_box = (1470, 690, 2130, 930)
    _draw_boxed_text(
        draw,
        takeaway_box,
        "How to say it",
        [
            "Particles touch the mesh.",
            "Rigid state receives the reaction.",
        ],
        _load_sans_font(24, bold=True),
        _load_sans_font(19, bold=False),
        fill=(254, 242, 242),
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)
    return out_path


def _make_force_diagnostic_mechanism_diagram(
    out_path: Path,
    snapshot_path: Path,
    summary: dict[str, Any],
) -> Path:
    W, H = 2200, 1180
    img = Image.new("RGB", (W, H), (248, 250, 252))
    draw = ImageDraw.Draw(img)
    title_font = _load_sans_font(36, bold=True)
    sec_font = _load_sans_font(26, bold=True)
    body_font = _load_sans_font(20, bold=False)
    small_font = _load_sans_font(18, bold=False)

    navy = (24, 58, 96)
    slate = (71, 85, 105)
    line = (181, 194, 209)
    white = (255, 255, 255)
    blue = (235, 244, 255)
    green = (236, 253, 245)
    amber = (255, 247, 237)

    draw.text((60, 34), "Trigger-Substep Force Diagnostic", font=title_font, fill=navy)
    draw.text(
        (60, 84),
        "Clean box control first: verify direction and magnitude before blaming bunny geometry.",
        font=small_font,
        fill=slate,
    )

    snap_panel = (60, 150, 1310, 1080)
    _draw_round_box(draw, snap_panel, white, line, radius=20, width=3)
    draw.text((84, 170), "Trigger snapshot: box sanity check", font=sec_font, fill=navy)

    if snapshot_path.exists():
        snap = Image.open(snapshot_path).convert("RGB")
        box_w = snap_panel[2] - snap_panel[0] - 40
        box_h = snap_panel[3] - snap_panel[1] - 80
        scale = min(box_w / snap.width, box_h / snap.height)
        new_size = (max(1, int(snap.width * scale)), max(1, int(snap.height * scale)))
        snap = snap.resize(new_size)
        paste_x = snap_panel[0] + 20 + (box_w - new_size[0]) // 2
        paste_y = snap_panel[1] + 52 + (box_h - new_size[1]) // 2
        img.paste(snap, (paste_x, paste_y))

    topk_records = summary.get("topk_particle_records", [])
    top = topk_records[0] if topk_records else {}
    right_x0 = 1380
    right_x1 = 2140
    box1 = (right_x0, 150, right_x1, 410)
    box2 = (right_x0, 450, right_x1, 720)
    box3 = (right_x0, 760, right_x1, 1080)

    _draw_boxed_text(
        draw,
        box1,
        "What the colors mean",
        [
            "white = outward surface normal",
            "red = external contact force",
            "blue = spring-maintenance force",
            "green = resulting total acceleration",
        ],
        sec_font,
        body_font,
        fill=blue,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    _draw_boxed_text(
        draw,
        box2,
        "Box control sanity check",
        [
            f"contact nodes = {int(summary.get('contact_node_count', 0))}",
            f"wrong-direction nodes = {int(summary.get('contact_nodes_with_wrong_direction', 0))}",
            f"inward-accel nodes = {int(summary.get('contact_nodes_with_inward_acceleration', 0))}",
            f"force-below-stop nodes = {int(summary.get('contact_nodes_force_below_stop', 0))}",
        ],
        sec_font,
        body_font,
        fill=green,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    interpretation_lines = [
        "This clean case says the diagnostic is working: contact direction is outward in the box control.",
        f"Top contact node: ext_n = {float(top.get('external_force_normal', 0.0)):.3f} N, "
        f"stop_force = {float(top.get('stop_force_required_normal', 0.0)):.3f} N.",
        "So the remaining question for bunny is not whether we can measure it; the question is whether the ear flips the normal or just starves the contact magnitude.",
    ]
    _draw_boxed_text(
        draw,
        box3,
        "How to use this on bunny",
        interpretation_lines,
        sec_font,
        body_font,
        fill=amber,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)
    return out_path


def _make_rope_profile_summary_diagram(out_path: Path, profile: dict[str, Any]) -> Path:
    W, H = 2200, 1180
    img = Image.new("RGB", (W, H), (248, 250, 252))
    draw = ImageDraw.Draw(img)
    title_font = _load_sans_font(36, bold=True)
    sec_font = _load_sans_font(26, bold=True)
    body_font = _load_sans_font(20, bold=False)
    small_font = _load_sans_font(18, bold=False)

    navy = (24, 58, 96)
    slate = (71, 85, 105)
    line = (181, 194, 209)
    white = (255, 255, 255)
    blue = (235, 244, 255)
    green = (236, 253, 245)
    amber = (255, 247, 237)
    red = (239, 68, 68)

    draw.text((60, 34), "Newton Playground Profiling", font=title_font, fill=navy)
    draw.text(
        (60, 84),
        "Full rope-control replay | viewer = null | 3 profiled runs + 1 warmup | sim_dt fixed",
        font=small_font,
        fill=slate,
    )

    runs = profile.get("runs", [])
    wall_ms = statistics.mean(float(run["wall_ms"]) for run in runs) if runs else 0.0
    sim_time = statistics.mean(float(run["sim_time_sec"]) for run in runs) if runs else 0.0
    realtime_slowdown = (wall_ms / 1000.0) / sim_time if sim_time > 1.0e-12 else 0.0
    aggregate = profile.get("aggregate", {})
    total_sub = float(aggregate.get("total_substep", {}).get("mean_of_run_means_ms", 0.0))
    ranking = list(profile.get("bottleneck_ranked_ops", []))[:6]

    chart_box = (60, 150, 1320, 1080)
    _draw_round_box(draw, chart_box, white, line, radius=20, width=3)
    draw.text((84, 170), "Top costs per substep [ms]", font=sec_font, fill=navy)

    label_x = 90
    bar_x0 = 430
    bar_x1 = 1230
    top_y = 255
    row_gap = 115
    max_ms = max((float(item.get("mean_of_run_means_ms", 0.0)) for item in ranking), default=1.0)
    max_ms = max(max_ms, 1.0e-6)
    palette = [
        (52, 211, 153),
        (59, 130, 246),
        (249, 115, 22),
        (168, 85, 247),
        (236, 72, 153),
        (234, 179, 8),
    ]
    for idx, item in enumerate(ranking):
        y = top_y + idx * row_gap
        name = str(item.get("op_name", "unknown"))
        ms = float(item.get("mean_of_run_means_ms", 0.0))
        width = int((bar_x1 - bar_x0) * ms / max_ms)
        draw.text((label_x, y + 12), name, font=body_font, fill=(31, 41, 55))
        draw.rounded_rectangle((bar_x0, y, bar_x0 + width, y + 48), radius=14, fill=palette[idx % len(palette)])
        draw.text((bar_x0 + width + 18, y + 12), f"{ms:.3f} ms", font=body_font, fill=navy)

    right_x0 = 1390
    right_x1 = 2140
    box1 = (right_x0, 150, right_x1, 400)
    box2 = (right_x0, 440, right_x1, 710)
    box3 = (right_x0, 750, right_x1, 1080)
    _draw_boxed_text(
        draw,
        box1,
        "Full replay",
        [
            f"trajectory frames = {int(profile.get('trajectory_frames', 0))}",
            f"total substeps = {int(profile.get('total_substeps', 0))}",
            f"simulated time = {sim_time:.2f} s",
            f"wall time = {wall_ms / 1000.0:.2f} s (~{realtime_slowdown:.1f}x slower than real time)",
        ],
        sec_font,
        body_font,
        fill=blue,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    collide_ms = float(aggregate.get("model_collide", {}).get("mean_of_run_means_ms", 0.0))
    grid_ms = float(aggregate.get("particle_grid_build", {}).get("mean_of_run_means_ms", 0.0))
    _draw_boxed_text(
        draw,
        box2,
        "What is not expensive",
        [
            f"particle_grid_build = {grid_ms:.3f} ms",
            f"model_collide = {collide_ms:.3f} ms",
            "This replay is not collision-bound.",
            "The main cost stays inside the rope spring-mass update itself.",
        ],
        sec_font,
        body_font,
        fill=green,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )
    share = (float(aggregate.get("solver_step", {}).get("mean_of_run_means_ms", 0.0)) / total_sub * 100.0) if total_sub > 1.0e-12 else 0.0
    _draw_boxed_text(
        draw,
        box3,
        "Conclusion",
        [
            f"`solver_step` still takes about {share:.1f}% of the substep wall time.",
            "Top sub-costs are integrate_particles, spring_forces, write_kinematic_state, and drag_correction.",
            "So the next optimization target is the internal rope update path, not viewer/render and not broad-phase collision.",
        ],
        sec_font,
        body_font,
        fill=amber,
        outline=line,
        title_fill=navy,
        body_fill=(31, 41, 55),
    )

    draw.text((1440, 1110), "red bars = real bottlenecks in the current replay case", font=small_font, fill=red)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)
    return out_path


def _render_code_png(
    out_path: Path,
    header: str,
    lines: list[str],
    font_size: int = 20,
    max_chars: int = 88,
) -> Path:
    mx, my = 32, 24
    font = _load_mono_font(font_size)
    font_bold = _load_mono_font(font_size, bold=True)
    bg, panel = (246, 248, 251), (255, 255, 255)
    txt_c, hdr_c = (34, 34, 34), (48, 87, 138)
    gutter_c = (120, 128, 138)
    line_hl_fill = (255, 219, 77, 92)
    line_hl_border = (240, 182, 86)
    token_colors = {
        Comment: txt_c,
        Keyword: txt_c,
        Keyword.Constant: txt_c,
        Name.Builtin: txt_c,
        Name.Function: txt_c,
        Name.Class: txt_c,
        String: txt_c,
        Number: txt_c,
        Operator: txt_c,
        Punctuation: txt_c,
        Token.Error: (176, 32, 32),
        Text: txt_c,
    }

    all_lines = [header, ""] + [ln[:max_chars] for ln in lines]
    probe = Image.new("RGB", (100, 100))
    d = ImageDraw.Draw(probe)
    bbox = d.textbbox((0, 0), "M", font=font)
    cw = max(8, bbox[2] - bbox[0])
    lh = max(14, bbox[3] - bbox[1] + 5)
    W = min(max(mx * 2 + max_chars * cw, 1100), 1700)
    H = min(max(my * 2 + len(all_lines) * lh, 360), 1100)
    img = Image.new("RGBA", (W, H), bg + (255,))
    draw = ImageDraw.Draw(img, "RGBA")
    draw.rectangle(
        [8, 8, W - 8, H - 8],
        fill=panel + (255,),
        outline=(204, 212, 220, 255),
        width=2,
    )
    y = my + 6

    def _color_for_token(tok):
        for ttype, color in token_colors.items():
            if tok in ttype:
                return color
        return txt_c

    lexer = PythonLexer(stripnl=False)
    line_infos: list[tuple[int, str, bool]] = []
    for i, raw_line in enumerate(all_lines):
        highlight = raw_line.lstrip().startswith(">>>")
        ln = raw_line.replace(">>>", "", 1).lstrip() if highlight else raw_line
        line_infos.append((i, ln, highlight))

    for i, ln, highlight in line_infos:
        if i == 0:
            draw.text((mx, y), ln, font=font_bold, fill=hdr_c)
            y += lh
            continue
        if ln.startswith("SIG "):
            draw.text((mx, y), ln, font=font_bold, fill=hdr_c)
            y += lh
            continue

        line_no = ""
        code = ln
        stripped = ln.lstrip()
        leading = len(ln) - len(stripped)
        if stripped and stripped[0].isdigit():
            parts = stripped.split(" ", 1)
            if len(parts) == 2 and parts[0].isdigit():
                line_no = parts[0]
                code = " " * leading + parts[1]

        x = mx
        x_code_start = x
        if line_no:
            draw.text(
                (x, y),
                f"{line_no:>4}",
                font=font_bold if highlight else font,
                fill=line_hl_border if highlight else gutter_c,
            )
            x += int(cw * 6.5)
            x_code_start = x

        rendered_parts: list[tuple[str, tuple[int, int, int]]] = []
        x_probe = x
        for tok, value in lex(code, lexer):
            if not value:
                continue
            parts = value.splitlines(keepends=True)
            for part in parts:
                clean = part.replace("\n", "")
                if not clean:
                    continue
                rendered_parts.append((clean, _color_for_token(tok)))
                x_probe += int(draw.textlength(clean, font=font))

        if highlight and x_probe > x_code_start:
            stroke_y = y + int(lh * 0.63)
            draw.line(
                [(x_code_start - 2, stroke_y), (x_probe + 4, stroke_y)],
                fill=line_hl_fill,
                width=11,
                joint="curve",
            )
            draw.line(
                [(x_code_start - 1, stroke_y + 1), (x_probe + 3, stroke_y + 1)],
                fill=(245, 199, 59, 110),
                width=7,
                joint="curve",
            )

        for clean, color in rendered_parts:
            draw.text((x, y), clean, font=font, fill=color)
            x += int(draw.textlength(clean, font=font))
        y += lh
    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)
    return out_path


def _pic_twocol(prs: Presentation, title: str, ll: str, lp: Path, rl: str, rp: Path):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)
    m = int(width * 0.02)
    gap = int(width * 0.03)
    cw = int((width - m * 2 - gap) / 2)
    lh = Inches(0.35)
    ch = int(height - m * 2 - lh)
    y0 = int(top + m)
    x0 = int(left + m)
    x1 = x0 + cw + gap
    _add_label(s, x0, y0, cw, lh, ll)
    _add_pic(s, lp, x0, int(y0 + lh), cw, ch)
    _add_label(s, x1, y0, cw, lh, rl)
    _add_pic(s, rp, x1, int(y0 + lh), cw, ch)
    return s


def _gif_single(prs: Presentation, title: str, gif: Path, subtitle: str | None = None):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)
    caption_h = int(Inches(0.38)) if subtitle else 0
    if subtitle:
        _add_label(s, left, top, width, caption_h, subtitle, font_size=12, bold=False)
    _add_pic(s, gif, left, top + caption_h, width, height - caption_h)
    return s


def _card_name_slide(prs: Presentation, name: str) -> None:
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    card = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        int(Inches(0.70)),
        int(Inches(0.50)),
        int(Inches(8.60)),
        int(Inches(4.55)),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = RGBColor(0x3E, 0x63, 0xD3)
    card.line.width = Pt(2.6)

    box = s.shapes.add_textbox(
        int(Inches(1.15)),
        int(Inches(2.00)),
        int(Inches(7.70)),
        int(Inches(0.70)),
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_marked_paragraph(
        p,
        name,
        font_name="Arial",
        font_size=20,
        color=RGBColor(0x11, 0x11, 0x11),
    )


def _card_media_slide(
    prs: Presentation,
    title: str,
    media_path: Path,
    *,
    media_top_offset: int | None = None,
) -> None:
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    card_left = int(Inches(0.70))
    card_top = int(Inches(0.42))
    card_w = int(Inches(8.60))
    card_h = int(Inches(4.75))

    card = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        card_left,
        card_top,
        card_w,
        card_h,
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = RGBColor(0xC7, 0xCD, 0xD6)
    card.line.width = Pt(1.2)

    title_box = s.shapes.add_textbox(
        card_left + int(Inches(0.18)),
        card_top + int(Inches(0.10)),
        card_w - int(Inches(0.36)),
        int(Inches(0.38)),
    )
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_marked_paragraph(
        p,
        title,
        font_name="Arial",
        font_size=10,
        color=RGBColor(0x33, 0x33, 0x33),
    )

    media_top = media_top_offset if media_top_offset is not None else card_top + int(Inches(0.55))
    _add_pic(
        s,
        media_path,
        card_left + int(Inches(0.16)),
        media_top,
        card_w - int(Inches(0.32)),
        card_h - (media_top - card_top) - int(Inches(0.16)),
    )


def _make_recall_bridge_diagram(out_path: Path) -> Path:
    W, H = 1800, 860
    img = Image.new("RGB", (W, H), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    title_font = _load_sans_font(26, bold=True)
    body_font = _load_sans_font(20, bold=False)

    boxes = [
        ((120, 260, 360, 510), (201, 214, 241), "PhysTwin"),
        ((470, 220, 870, 550), (219, 240, 206), "Bridge\nIntermediate\nRepresentation"),
        ((980, 260, 1220, 510), (243, 237, 205), "Newton\nRebuild"),
        ((1330, 260, 1640, 510), (224, 232, 215), "Newton\nInteraction"),
    ]

    for xy, fill, text in boxes:
        _draw_round_box(draw, xy, fill=fill, outline=(176, 186, 198), radius=24, width=3)
        x0, y0, x1, y1 = xy
        lines = text.split("\n")
        total_h = len(lines) * 34
        y = int((y0 + y1 - total_h) / 2)
        for line in lines:
            _draw_centered_label(draw, ((x0 + x1) / 2, y + 12), line, title_font if len(lines) == 1 else body_font, (60, 60, 60))
            y += 34

    arrow_y = 385
    _draw_arrow(draw, (360, arrow_y), (470, arrow_y), (123, 133, 148), width=8, head=18)
    _draw_arrow(draw, (870, arrow_y), (980, arrow_y), (123, 133, 148), width=8, head=18)
    _draw_arrow(draw, (1220, arrow_y), (1330, arrow_y), (123, 133, 148), width=8, head=18)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)
    return out_path


def _gif_compare_large(
    prs: Presentation,
    title: str,
    left_caption: str,
    left_gif: Path,
    right_caption: str,
    right_gif: Path,
    footer: str | None = None,
):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    title_box = s.shapes.add_textbox(
        REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H
    )
    _set_title_textbox(title_box, title, size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)

    gap = int(Inches(0.22))
    footer_h = int(Inches(0.28)) if footer else 0
    caption_h = int(Inches(0.42))
    pic_h = height - caption_h - footer_h - int(Inches(0.08))
    pic_w = int((width - gap) / 2)
    x0 = left
    x1 = left + pic_w + gap
    y0 = top
    caption_y = y0 + pic_h + int(Inches(0.04))

    _add_pic(s, left_gif, x0, y0, pic_w, pic_h)
    _add_pic(s, right_gif, x1, y0, pic_w, pic_h)
    _add_label(s, x0, caption_y, pic_w, caption_h, left_caption, font_size=12, bold=False)
    _add_label(s, x1, caption_y, pic_w, caption_h, right_caption, font_size=12, bold=False)
    if footer:
        _add_label(
            s,
            left,
            caption_y + caption_h + int(Inches(0.02)),
            width,
            footer_h,
            footer,
            font_size=11,
            bold=False,
        )
    return s


def _picture_with_footer(prs: Presentation, title: str, image: Path, footer_lines: list[str]):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)
    _add_pic(s, image, REF_MEDIA_PIC_LEFT, REF_MEDIA_PIC_TOP, REF_MEDIA_PIC_W, REF_MEDIA_PIC_H)
    footer = s.shapes.add_textbox(REF_BODY_LEFT, REF_MEDIA_FOOTER_TOP, REF_BODY_W, REF_MEDIA_FOOTER_H)
    _set_lines(footer, footer_lines)
    return s


def _pipeline_flow_slide(
    prs: Presentation,
    title: str,
    steps: list[tuple[str, str]],
    takeaway: str,
):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=24)

    left, top, width, height = _detach_body_placeholder(s)
    margin_x = int(Inches(0.30))
    gap = int(Inches(0.20))
    box_h = int(Inches(1.65))
    y = int(top + Inches(0.55))
    footer_h = int(Inches(0.42))
    n = len(steps)
    box_w = int((width - 2 * margin_x - (n - 1) * gap) / n)

    for idx, (heading, detail) in enumerate(steps):
        x = left + margin_x + idx * (box_w + gap)
        shape = s.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            y,
            box_w,
            box_h,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFB)
        shape.line.color.rgb = RGBColor(0x7D, 0x98, 0xB6)
        shape.line.width = Pt(2.0)
        _set_lines(shape, [heading, f"  {detail}"])
        tf = shape.text_frame
        if tf.paragraphs:
            for r in tf.paragraphs[0].runs:
                r.font.size = Pt(20)
                r.font.bold = True
            if len(tf.paragraphs) > 1:
                for r in tf.paragraphs[1].runs:
                    r.font.size = Pt(16)
                    r.font.bold = False
        if idx < n - 1:
            arrow_w = int(gap)
            arrow_h = int(Inches(0.48))
            arrow = s.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                x + box_w + int(Inches(0.02)),
                y + int((box_h - arrow_h) / 2),
                arrow_w - int(Inches(0.04)),
                arrow_h,
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0x7D, 0x98, 0xB6)
            arrow.line.color.rgb = RGBColor(0x7D, 0x98, 0xB6)

    _add_label(
        s,
        REF_BODY_LEFT + int(Inches(0.30)),
        y + box_h + int(Inches(0.18)),
        REF_BODY_W - int(Inches(0.60)),
        footer_h,
        takeaway,
        font_size=16,
        bold=False,
    )
    return s


def _make_factor_chart(rows: list[dict], out_path: Path) -> Path:
    if not rows:
        return _write_text_image(
            out_path,
            "64-run OFF Grid: Reconstructed Aggregate Summary",
            [
                "Contact scale is the strongest lever:",
                "  scale 0.5 -> about -2.90",
                "  scale 2.0 -> about -0.14",
                "  scale 3.0 -> about -0.13",
                "",
                "Bunny mass helps, but is not the main lever:",
                "  mass 1 kg -> about -2.65",
                "  mass 1000 kg -> about -0.10",
                "  mass 5000 kg -> about -0.088",
                "",
                "Drop height hurts, but is not a standalone fix:",
                "  height 0.1 m -> about -0.216",
                "  height 1.0 m -> about -1.98",
            ],
            width=1500,
            line_h=42,
        )

    def group_mean(key: str):
        uniq = sorted({r[key] for r in rows})
        vals = []
        for u in uniq:
            subset = [r["min_clearance_max_to_bunny_top"] for r in rows if r[key] == u]
            vals.append(statistics.mean(subset))
        return uniq, vals

    fig, axes = plt.subplots(1, 3, figsize=(14, 4.5))
    for ax, key, title in zip(
        axes,
        ("contact_dist_scale", "rigid_mass", "drop_height_m"),
        ("Contact Scale", "Bunny Mass", "Drop Height"),
    ):
        xs, ys = group_mean(key)
        ax.bar(range(len(xs)), ys, color="#4C78A8")
        ax.axhline(0.0, color="#999999", lw=1)
        ax.set_title(title, fontsize=12)
        ax.set_xticks(range(len(xs)))
        ax.set_xticklabels([str(x).rstrip("0").rstrip(".") for x in xs], fontsize=10)
        ax.set_ylabel(
            "Avg max-clearance to bunny top\n(closer to 0 is better)", fontsize=9
        )
        ax.grid(axis="y", alpha=0.25)
    fig.suptitle("64-run OFF grid: average penetration trend", fontsize=14)
    fig.tight_layout()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return out_path


def _make_bunny_sweep_chart(csv_path: Path, out_path: Path) -> Path:
    rows = list(csv.DictReader(csv_path.open(encoding="utf-8")))
    labels = [
        r["case_dir"]
        .replace("center_", "center ")
        .replace("earridge_", "ear ridge ")
        .replace("bodycenter_", "body center ")
        for r in rows
    ]
    max_pen = [float(r["max_penetration_depth_bunny_mesh_m"]) * 1000.0 for r in rows]
    p99_pen = [float(r["final_penetration_p99_bunny_mesh_m"]) * 1000.0 for r in rows]

    fig, ax = plt.subplots(figsize=(12.5, 5.4))
    x = list(range(len(labels)))
    w = 0.36
    ax.bar([i - w / 2 for i in x], max_pen, width=w, label="Max penetration (mm)", color="#d95f02")
    ax.bar([i + w / 2 for i in x], p99_pen, width=w, label="Final p99 penetration (mm)", color="#1b9e77")
    ax.axhline(5.0, color="#b22222", linestyle="--", linewidth=1.2, label="5 mm threshold")
    ax.axhline(2.0, color="#1e5aa8", linestyle="--", linewidth=1.2, label="2 mm threshold")
    ax.set_ylabel("Penetration (mm)")
    ax.set_title("Bunny mesh: local geometry and contact-scale sweep")
    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=25, ha="right")
    ax.legend(fontsize=8, ncol=2, loc="upper right")
    fig.tight_layout()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=180, bbox_inches="tight")
    plt.close(fig)
    return out_path

MAX_CODE_EXCERPT_LINES = 20
MAX_CODE_HIGHLIGHT_LINES = 5


def _select_excerpt_line_numbers(
    candidate_numbers: list[int],
    highlight_numbers: set[int],
    *,
    max_lines: int = MAX_CODE_EXCERPT_LINES,
) -> list[int]:
    ordered = sorted(dict.fromkeys(int(v) for v in candidate_numbers))
    if len(ordered) <= max_lines:
        return ordered

    clamped_highlights = [ln for ln in ordered if ln in highlight_numbers][:MAX_CODE_HIGHLIGHT_LINES]
    if not clamped_highlights:
        head = max_lines // 2
        tail = max_lines - head
        keep = ordered[:head] + ordered[-tail:]
        return sorted(dict.fromkeys(keep))

    selected: set[int] = set(clamped_highlights)
    radius = 1
    while len(selected) < min(max_lines, len(ordered)):
        added = False
        for h in clamped_highlights:
            for ln in (h - radius, h + radius):
                if ln in ordered and ln not in selected:
                    selected.add(ln)
                    added = True
                    if len(selected) >= max_lines:
                        break
            if len(selected) >= max_lines:
                break
        if len(selected) >= max_lines:
            break
        if not added:
            break
        radius += 1

    for ln in ordered:
        if len(selected) >= max_lines:
            break
        if ln not in selected:
            selected.add(ln)

    return sorted(selected)


def _extract_code_segments(
    path: Path,
    segments: list[tuple[int, int]],
    *,
    highlight_lines: set[int] | None = None,
) -> list[str]:
    if path.exists():
        source_lines = path.read_text(encoding="utf-8").splitlines()
    else:
        max_lineno = max(end for _, end in segments)
        fallback_lines = ["# fallback excerpt: original Newton source is not vendored in this checkout"] * max_lineno
        for lineno, text in FALLBACK_CODE_LINE_MAP.get(path, {}).items():
            if 1 <= lineno <= max_lineno:
                fallback_lines[lineno - 1] = text
        source_lines = fallback_lines

    clamped_highlights = set(sorted(highlight_lines or set())[:MAX_CODE_HIGHLIGHT_LINES])
    candidate_numbers: list[int] = []
    for start, end in segments:
        candidate_numbers.extend(range(start, end + 1))
    selected_numbers = _select_excerpt_line_numbers(
        candidate_numbers,
        clamped_highlights,
        max_lines=MAX_CODE_EXCERPT_LINES,
    )

    out: list[str] = []
    prev_lineno: int | None = None
    for lineno in selected_numbers:
        if prev_lineno is not None and lineno != prev_lineno + 1:
            out.append("...")
        prefix = ">>> " if lineno in clamped_highlights else ""
        out.append(f"{prefix}{lineno:4d} {source_lines[lineno - 1]}")
        prev_lineno = lineno
    return out


def _split_top_level_commas(text: str) -> list[str]:
    parts: list[str] = []
    buf: list[str] = []
    depth = 0
    for ch in text:
        if ch in "([{":
            depth += 1
        elif ch in ")]}":
            depth = max(0, depth - 1)
        if ch == "," and depth == 0:
            piece = "".join(buf).strip()
            if piece:
                parts.append(piece)
            buf = []
            continue
        buf.append(ch)
    piece = "".join(buf).strip()
    if piece:
        parts.append(piece)
    return parts


def _compact_signature_lines(
    path: Path, start: int, end: int, *, max_len: int = 108
) -> list[str]:
    src = path.read_text(encoding="utf-8").splitlines()
    raw = " ".join(src[start - 1 : end])
    raw = re.sub(r"#.*", "", raw)
    raw = re.sub(r"\s+", " ", raw).strip()
    m = re.search(r"def\s+([A-Za-z_][A-Za-z0-9_]*)\((.*)\)\s*:", raw)
    if not m:
        return []
    fn_name = m.group(1)
    params_block = m.group(2).strip()
    params: list[str] = []
    for part in _split_top_level_commas(params_block):
        part = part.strip()
        if not part:
            continue
        if ":" in part:
            part = part.split(":", 1)[0].strip()
        if "=" in part:
            part = part.split("=", 1)[0].strip()
        if part:
            params.append(part)
    sig = f"SIG def {fn_name}(" + ", ".join(params) + ")"
    wrapped = textwrap.wrap(sig, width=max_len, break_long_words=False, break_on_hyphens=False)
    return wrapped[:2]


def _code_excerpt_image(
    path: Path, title: str, lines: list[str], out_dir: Path
) -> Path:
    safe = "".join(c if c.isalnum() or c == "_" else "_" for c in title)[:50]
    try:
        rel_path = path.relative_to(ROOT)
    except ValueError:
        rel_path = path
    header = f"# {rel_path}"
    if not path.exists():
        header += " (fallback excerpt)"
    return _render_code_png(out_dir / f"{safe}.png", header, lines)


def _write_text_image(
    out_path: Path, title: str, lines: list[str], width: int = 1200, line_h: int = 56
) -> Path:
    bg = (247, 249, 252)
    panel = (255, 255, 255)
    border = (205, 214, 224)
    text = (34, 34, 34)
    title_c = (40, 74, 120)
    font = _load_mono_font(24)
    title_font = _load_mono_font(30)
    height = 120 + len(lines) * line_h
    img = Image.new("RGB", (width, height), bg)
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle(
        [20, 20, width - 20, height - 20],
        radius=18,
        fill=panel,
        outline=border,
        width=2,
    )
    draw.text((50, 45), title, font=title_font, fill=title_c)
    y = 105
    for line in lines:
        draw.text((60, y), line, font=font, fill=text)
        y += line_h
    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)
    return out_path


def _make_pptx_slide_preview(pptx_path: Path, out_path: Path, slide_no: int = 1) -> Path:
    if out_path.exists() and pptx_path.exists() and out_path.stat().st_mtime >= pptx_path.stat().st_mtime:
        return out_path
    if not pptx_path.exists():
        return _write_text_image(
            out_path,
            "Missing PPTX Recall Asset",
            [f"Could not find: {pptx_path}"],
            width=1400,
            line_h=48,
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    pdf_path = out_path.parent / f"{pptx_path.stem}.pdf"
    pdf_is_fresh = pdf_path.exists() and pdf_path.stat().st_mtime >= pptx_path.stat().st_mtime
    try:
        if not pdf_is_fresh:
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(out_path.parent),
                    str(pptx_path),
                ],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
        prefix = out_path.with_suffix("")
        subprocess.run(
            [
                "pdftoppm",
                "-png",
                "-f",
                str(slide_no),
                "-singlefile",
                str(pdf_path),
                str(prefix),
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        return out_path
    except subprocess.CalledProcessError:
        return _write_text_image(
            out_path,
            "Earlier Deck Preview Unavailable",
            [
                f"Could not render: {pptx_path.name}",
                "Fallback: continue the meeting with the regenerated bridge deck below.",
            ],
            width=1400,
            line_h=48,
        )


def _boxed_bullets(prs: Presentation, title: str, heading: str, bullets: list[str]):
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, title, size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)

    box = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        REF_SUMMARY_BOX_LEFT,
        REF_SUMMARY_BOX_TOP,
        REF_SUMMARY_BOX_W,
        REF_SUMMARY_BOX_H,
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFD)
    box.line.color.rgb = RGBColor(0xCC, 0xD5, 0xE0)
    box.line.width = Pt(1.5)

    text_box = s.shapes.add_textbox(
        REF_SUMMARY_TEXT_LEFT,
        REF_SUMMARY_TEXT_TOP,
        REF_SUMMARY_TEXT_W,
        REF_SUMMARY_TEXT_H,
    )
    tf = text_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_marked_paragraph(
        p,
        heading,
        font_name="Arial",
        font_size=24,
        color=RGBColor(0x2C, 0x4A, 0x78),
        accent_color=RGBColor(0x1F, 0x4E, 0x79),
        bold_default=True,
    )

    for line in bullets:
        p = tf.add_paragraph()
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        _set_marked_paragraph(
            p,
            line,
            font_name="Arial",
            font_size=22,
            color=RGBColor(0x22, 0x22, 0x22),
            accent_color=RGBColor(0x1F, 0x4E, 0x79),
        )
    return s


def build_deck(
    prs: Presentation, rows: list[dict], gif_dir: Path, image_dir: Path
) -> None:
    _delete_all_slides(prs)

    factor_chart = _make_factor_chart(rows, image_dir / "factor_chart.png")
    phys_png = _code_excerpt_image(
        PHYSTWIN_SIM,
        "phystwin_self_collision",
        _extract_code_segments(
            PHYSTWIN_SIM,
            [(196, 205), (220, 220), (294, 296)],
            highlight_lines={196, 198, 199, 205, 294},
        ),
        image_dir,
    )
    newton_self_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_self_collision",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(79, 96)],
            highlight_lines={80, 88, 90, 94},
        ),
        image_dir,
    )
    newton_contact_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_particle_body_contact",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(216, 222), (245, 248), (263, 267)],
            highlight_lines={216, 218, 245, 248, 263},
        ),
        image_dir,
    )
    mesh_query_png = _code_excerpt_image(
        NEWTON_GEOM_KERNELS,
        "newton_mesh_contact_detection",
        _extract_code_segments(
            NEWTON_GEOM_KERNELS,
            [(1112, 1115), (1132, 1132)],
            highlight_lines={1113, 1114, 1115, 1132},
        ),
        image_dir,
    )
    dist_map_png = _code_excerpt_image(
        NEWTON_IMPORT_IR,
        "contact_dist_mapping",
        _extract_code_segments(
            NEWTON_IMPORT_IR,
            [(828, 833)],
            highlight_lines={829, 830, 831, 832},
        ),
        image_dir,
    )

    baseline_gif = _ensure_gif(
        BASELINE_OFF_MP4, gif_dir, stem_override="baseline_off_ccd1_h0p1_bm5000"
    )
    best_gif = _ensure_gif(
        BEST_OFF_MP4, gif_dir, stem_override="best_off_ccd3_h0p1_bm5000"
    )
    mass_light_gif = _ensure_gif(
        MASS_LIGHT_MP4, gif_dir, stem_override="mass_bm1_ccd1_h0p1"
    )
    mass_mid1_gif = _ensure_gif(
        MASS_MID1_MP4, gif_dir, stem_override="mass_bm10_ccd1_h0p1"
    )
    mass_mid2_gif = _ensure_gif(
        MASS_MID2_MP4, gif_dir, stem_override="mass_bm1000_ccd1_h0p1"
    )
    mass_heavy_gif = _ensure_gif(
        MASS_HEAVY_MP4, gif_dir, stem_override="mass_bm5000_ccd1_h0p1"
    )
    height_low_gif = _ensure_gif(
        HEIGHT_LOW_MP4, gif_dir, stem_override="height_h0p1_ccd1_bm5000"
    )
    height_mid1_gif = _ensure_gif(
        HEIGHT_MID1_MP4, gif_dir, stem_override="height_h0p25_ccd1_bm5000"
    )
    height_mid2_gif = _ensure_gif(
        HEIGHT_MID2_MP4, gif_dir, stem_override="height_h0p5_ccd1_bm5000"
    )
    height_high_gif = _ensure_gif(
        HEIGHT_HIGH_MP4, gif_dir, stem_override="height_h1p0_ccd1_bm5000"
    )
    ccd_small_gif = _ensure_gif(
        CCD_SMALL_MP4, gif_dir, stem_override="ccd_s0p5_h0p1_bm5000"
    )
    ccd_mid1_gif = _ensure_gif(
        CCD_MID1_MP4, gif_dir, stem_override="ccd_s1p0_h0p1_bm5000"
    )
    ccd_mid2_gif = _ensure_gif(
        CCD_MID2_MP4, gif_dir, stem_override="ccd_s2p0_h0p1_bm5000"
    )
    ccd_large_gif = _ensure_gif(
        CCD_LARGE_MP4, gif_dir, stem_override="ccd_s3p0_h0p1_bm5000"
    )

    _title_slide(
        prs,
        "PhysTwin → Newton Bridge",
        [
            "Meeting Update — OFF cloth-bunny penetration",
            f"{date.today().isoformat()}",
        ],
    )

    _body(
        prs,
        "3 Questions For This Meeting",
        [
            "Q1. Can ON directly match PhysTwin self-collision?",
            "Q2. In OFF, why does cloth still penetrate bunny?",
            "Q3. What progress did 64 runs actually buy us?",
        ],
    )

    _pic_twocol(
        prs,
        "Hypothesis: Can We Directly Migrate Self-Collision?",
        "PhysTwin: approaching-only impulse",
        phys_png,
        "Newton: always-on penalty force",
        newton_self_png,
    )

    _body(
        prs,
        "Takeaway: ON Is Not A Parity Target",
        [
            "Self-collision models are mathematically different",
            "ON = Newton-native ablation, not PhysTwin parity",
            "Main target = OFF cloth-bunny external contact",
        ],
    )

    _pic_twocol(
        prs,
        "OFF Contact: Why Particle Radius Matters",
        "Detection: mesh closest-point query uses margin + radius",
        mesh_query_png,
        "Force: penetration depth subtracts particle radius",
        newton_contact_png,
    )

    _twocol(
        prs,
        "Discussion: Why Increasing Particle Radius Helps",
        [
            "Rigid bunny is NOT a particle cloud",
            "  Newton queries the mesh surface",
            "  via closest-point / sign-normal search",
            "",
            "Detection test uses:",
            "  d < margin + particle radius",
        ],
        [
            "So larger radius means:",
            "  contact is detected earlier",
            "  penetration depth c becomes larger sooner",
            "",
            "Same bunny mesh, different effective particle radius",
            "  more support before deep penetration",
        ],
    )

    _twocol(
        prs,
        "Why Rope Looked Easier Than Cloth",
        [
            "Rope:",
            "  1D bead-chain contact",
            "  radius ≈ 0.013 m",
            "  1753 particles",
        ],
        [
            "Cloth:",
            "  2D sheet support problem",
            "  radius ≈ 0.0051 m",
            "  4409 particles",
        ],
    )

    _body(
        prs,
        "OFF Experiment Matrix",
        [
            "OFF only, object_mass = 1 kg / particle",
            "`contact_collision_dist` scale = 0.5 / 1.0 / 2.0 / 3.0",
            "Drop height = 0.1 / 0.25 / 0.5 / 1.0 m",
            "Bunny mass = 1 / 10 / 1000 / 5000 kg",
            "64 runs: video + scene + summary + csv",
        ],
    )

    _gif_grid_2x2(
        prs,
        "Hypothesis H1: Bunny Mass Matters",
        [
            ("Bunny mass = 1 kg", mass_light_gif),
            ("Bunny mass = 10 kg", mass_mid1_gif),
            ("Bunny mass = 1000 kg", mass_mid2_gif),
            ("Bunny mass = 5000 kg", mass_heavy_gif),
        ],
        common_settings="self-collision OFF | object mass = 1 kg / particle | `contact_collision_dist` scale = 1.0 | drop height = 0.1 m",
    )

    _gif_grid_2x2(
        prs,
        "Hypothesis H2: Height Matters, But Not Alone",
        [
            ("Drop height = 0.1 m", height_low_gif),
            ("Drop height = 0.25 m", height_mid1_gif),
            ("Drop height = 0.5 m", height_mid2_gif),
            ("Drop height = 1.0 m", height_high_gif),
        ],
        common_settings="self-collision OFF | object mass = 1 kg / particle | bunny mass = 5000 kg | `contact_collision_dist` scale = 1.0",
    )

    _gif_grid_2x2(
        prs,
        "Hypothesis H3: Contact Shell Thickness Matters Most",
        [
            ("`contact_collision_dist` scale = 0.5", ccd_small_gif),
            ("`contact_collision_dist` scale = 1.0", ccd_mid1_gif),
            ("`contact_collision_dist` scale = 2.0", ccd_mid2_gif),
            ("`contact_collision_dist` scale = 3.0", ccd_large_gif),
        ],
        common_settings="self-collision OFF | object mass = 1 kg / particle | bunny mass = 5000 kg | drop height = 0.1 m",
    )

    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, "Result: `contact_collision_dist` Is The Strongest Lever", size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)
    _add_pic(s, factor_chart, left, top, width, height)
    _add_label(
        s,
        int(left + 20),
        int(top + height - Inches(0.45)),
        int(width - 40),
        Inches(0.35),
        "Larger `contact_collision_dist` (larger Newton `particle_radius`) helps most. Bunny mass helps. Height hurts, but does not fix the issue alone.",
        font_size=12,
        bold=False,
    )

    _gif_twocol(
        prs,
        "Best OFF Example: Baseline vs Best",
        "Baseline\n`contact_collision_dist` scale = 1.0",
        baseline_gif,
        "Best so far\n`contact_collision_dist` scale = 3.0",
        best_gif,
        common_settings="self-collision OFF | object mass = 1 kg / particle | drop height = 0.1 m | bunny mass = 5000 kg",
    )

    _boxed_bullets(
        prs,
        "What This Week Actually Proved",
        "Validated / Falsified Hypotheses",
        [
            "Validated: self-collision cannot be directly ported",
            "Validated: bunny mass matters, but is not sufficient",
            "Validated: larger `contact_collision_dist` (larger Newton `particle_radius`) strongly reduces penetration",
            "Falsified: lower drop height alone is enough",
            "Falsified: 10x / 100x stiffness-damping scaling is a robust fix",
        ],
    )

    _body(
        prs,
        "Project Progress",
        [
            "Separated parity from ablation",
            "Localized failure to OFF external contact",
            "Built a reproducible 64-run grid + metrics",
            "Identified `contact_collision_dist` / Newton `particle_radius` as the main knob",
        ],
    )

    _body(
        prs,
        "Next Step",
        [
            "Narrow sweep: `contact_collision_dist` scale = 2.0 / 2.5 / 3.0 / 3.5",
            "Keep drop height = 0.1 m, bunny mass = 5000 kg",
            "Then try moderate ke / kd around the stable region",
        ],
    )

    _title_slide(prs, "Thank You", ["Questions?"])


TRANSCRIPT = textwrap.dedent(
    f"""\
    # Meeting Transcript — PhysTwin → Newton Bridge

    语言：中文主讲 + English terminology  
    主题：OFF cloth-bunny penetration after isolating self-collision mismatch

    ---

    ## Slide 1 — Title

    今天这次汇报，我想把叙事彻底改清楚。  
    我们不再把 `ON` case 讲成 “PhysTwin self-collision 成功迁移到 Newton”。  
    现在更准确的说法是：

    1. **PhysTwin self-collision 和 Newton native self-collision 不是一个模型。**
    2. 所以 `ON` case 只是 Newton-native ablation，不是 parity target。
    3. 我们的主问题是：**`OFF` case 里，cloth 为什么还会穿过 bunny？**

    ---

    ## Slide 2 — 3 Questions

    这一页我会明确三个问题。

    第一，`ON` 能不能直接对齐 PhysTwin self-collision？  
    第二，`OFF` 里 cloth-bunny penetration 的主因是什么？  
    第三，这周我们到底拿到了什么 progress，不只是做了很多实验而已。

    ---

    ## Slide 3 — Self-Collision Mismatch: Source Proof

    这里是第一条关键结论，用源码直接证明。

    PhysTwin 在 [spring_mass_warp.py]({PHYSTWIN_SIM}#L196) 里，只有在：
    - `dis_len < collision_dist`
    - 且 `dot(dis, relative_v) < -1e-4`

    时才触发 self-collision。  
    也就是它要求**距离够近，而且还在相向靠近**。

    然后在 [spring_mass_warp.py]({PHYSTWIN_SIM}#L281) 到 [spring_mass_warp.py]({PHYSTWIN_SIM}#L296)，它不是加 penalty force，而是：
    - 先累计 `J_sum`
    - 再做 `J_average`
    - 最后 `v_new = v_old - J_average / m`

    这是一条 **impulse / velocity correction** 路径。

    再看 Newton。  
    在 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L47) 到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L96)，native self-collision 是：
    - 查 hash grid 近邻
    - 计算 `err = d - r_i - r_j`
    - 只要 `err <= k_cohesion` 就直接加 `particle_force(...)`

    它没有 PhysTwin 那个 “must be approaching” 的 gate。  
    它也不是 impulse correction，而是 penalty-force accumulation。

    所以这页的结论非常明确：

    **ON case 不能再被表述成 PhysTwin self-collision parity。**

    ---

    ## Slide 4 — Scope Shift

    既然 self-collision 机制本身不同，我们就不能继续把主要精力放在 “为什么 ON 不像 PhysTwin” 上。  
    那样会把问题叙述错位。

    所以这周我们把 scope 缩到一个更干净的问题：

    **即使把 self-collision 关掉，cloth 还是会穿过 bunny。**

    这说明主问题不是 self-collision，而是 **external particle-body contact**。

    ---

    ## Slide 5 — OFF Contact Model: Source Proof

    这页是第二条核心源码证据。

    在 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L828) 到 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L833)，bridge 会把：

    - `contact_collision_dist`

    映射成：

    - `particle_radius = contact_collision_dist * 0.5`

    也就是说，IR 里的距离参数，最后变成了 Newton 里的 `particle_radius`。

    然后在 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L216) 到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L267)，particle-body contact 的核心是：

    - `c = dot(n, px - bx) - particle_radius`
    - `fn = n * c * ke`
    - `fd = n * min(vn,0) * kd`

    这里第一次出现的 `ke` 和 `kd` 先定义清楚：
    - `ke` = contact stiffness coefficient，决定 penetration 被转成多大的法向弹性力
    - `kd` = contact damping coefficient，决定法向相对速度被耗散得有多强

    所以这是一个非常标准的 **penalty contact**：

    - 先发生 penetration `c`
    - 再用 penalty spring / damping 往回推

    在 [kernels.py]({NEWTON_GEOM_KERNELS}#L1104) 到 [kernels.py]({NEWTON_GEOM_KERNELS}#L1132)，对 mesh shape 的处理是：
    - 先 `mesh_query_point_sign_normal(...)`
    - 再拿最近点和法向
    - 然后判断 `d < margin + radius`

    所以 radius 在这里直接决定：
    - 查询范围有多大
    - 多早能进入接触检测

    再到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L216)，penetration depth 写成：
    - `c = dot(n, px - bx) - particle_radius`

    也就是说 radius 不只是“画大一点”，它直接进入了接触深度。

    下面统一不用 “shell scale” 这种口语。  
    这里更严谨的写法是：
    - 调大 bridge 输入里的 `contact_collision_dist`
    - 于是 Newton 里的 `particle_radius` 也会一起变大

    所以调大 `contact_collision_dist` / `particle_radius` 的物理含义是：
    - 更早 detect
    - 更早产生 penalty support

    这就是为什么它对 rigid mesh contact 真的有用。

    ---

    ## Slide 6 — Discussion: Why Increasing Particle Radius Helps

    这页是对上一页源码的直接解释，不再只停留在代码层面。

    左边先回答第一个问题：

    **bunny 是不是被实现成 particle cloud？**

    不是。  
    它在 Newton 里是 `mesh shape`，通过最近点和法向查询来处理，不是把 bunny 离散成一堆 particle 再和 cloth 互撞。

    右边回答第二个问题：

    **为什么调大 particle radius 会有用？**

    因为 radius 在这条链路里进入了两个关键位置：

    1. 检测阶段：`d < margin + radius`
       - radius 越大，越早触发 contact detection

    2. 力阶段：`c = dot(n, px - bx) - particle_radius`
       - radius 越大，越早得到更大的有效 penetration depth

    所以就算 bunny mesh 本身没有变，particle 实际上也在带着一个更大的 Newton `particle_radius` 去撞这个 mesh。  
    这就是为什么调大 `contact_collision_dist`（也就是调大 Newton `particle_radius`）会更不容易穿。

    ---

    ## Slide 7 — Why Rope Looked Easier

    这一页解释为什么 rope 的问题相对不明显，而 cloth 会更快暴露出 penetration。

    不是因为 rope 走了另一套碰撞代码。  
    rope 和 cloth 在 Newton 里本质上都走 particle-vs-shape penalty contact。

    关键差别在几何和任务本身：

    - rope 更像一串 bead-chain，1D 接触带
    - cloth 是整片 2D sheet，需要 bunny 提供面积支撑

    另外数值上：

    - rope 的 Newton `particle_radius` 约 `0.013m`
    - cloth 的 Newton `particle_radius` 只有约 `0.0051m`

    所以 rope 更像“粗珠子撞 bunny”，cloth 更像“薄片压 bunny”。  
    同一个 penalty model，在 cloth 上更容易表现成 penetration。

    ---

    ## Slide 8 — 64-run OFF Grid

    为了不再凭直觉猜，我们把 OFF case 做成完整参数网格。

    固定条件：
    - `self_collision = off`
    - `object_mass = 1 kg / particle`
    - `frames = 180`
    - 渲染半径 `0.01m`，保证 cloth 后期仍可见

    Sweep 三个轴：
    - ``contact_collision_dist`` scale = `0.5 / 1.0 / 2.0 / 3.0`
    - `drop_height = 0.1 / 0.25 / 0.5 / 1.0`
    - `bunny_mass = 1 / 10 / 1000 / 5000`

    总计 64 个 run。  
    每个 run 都有：
    - video
    - scene
    - summary
    - aggregate CSV

    ---

    ## Slide 9 — Hypothesis H1: Bunny Mass Matters

    这页只回答一个问题：

    **如果 bunny 更重，penetration 会不会减轻？**

    我们控制其它量不变：
    - ``contact_collision_dist`` scale = `1.0`
    - `drop height = 0.1`

    这一页不是只放两个例子，而是同时放四个 GIF：
    - `bunny mass = 1`
    - `bunny mass = 10`
    - `bunny mass = 1000`
    - `bunny mass = 5000`

    视频上能直接看到，轻 bunny 更容易被 cloth 顶走。  
    这支持了 `H1`：

    **bunny 太轻，确实会明显加重 penetration。**

    但这页也同时说明：
    - 就算 bunny 很重，penetration 也没有完全消失

    所以 `H1` 成立，但不是唯一主因。

    ---

    ## Slide 10 — Hypothesis H2: Height Matters, But Not Alone

    这页回答第二个问题：

    **是不是因为 drop 太高，所以撞击速度太大，才导致穿模？**

    我们同样控制：
    - ``contact_collision_dist`` scale = `1.0`
    - `bunny mass = 5000`

    这一页放四个高度：
    - `0.1`
    - `0.25`
    - `0.5`
    - `1.0`

    结果很清楚：
    - 高度高的时候 penetration 更明显
    - 但高度低的时候并没有自动消失

    所以 `H2` 的精确表述应该是：

    **drop height 是 aggravating factor，不是 root fix。**

    ---

    ## Slide 11 — OFF Hypothesis 3: `contact_collision_dist` Is The Strongest Lever

    这页是最关键的对照。

    问题是：

    **`contact_collision_dist` 会不会太小，也就是 Newton `particle_radius` 会不会太小？**

    我们控制：
    - `drop height = 0.1`
    - `bunny mass = 5000`

    这一页放四个 ``contact_collision_dist`` scale：
    - `0.5`
    - `1.0`
    - `2.0`
    - `3.0`

    这页视频给出的结论最直接：

    **`contact_collision_dist` 变大之后，cloth 明显更接近挂在 bunny 上，而不是整体压过去。**

    这就是为什么我说：
    - 现在最强的主导变量不是质量
    - 而是 ``contact_collision_dist`` scale

    ---

    ## Slide 12 — Aggregate Trend

    这页把单个视频观察，变成整个 64-run grid 的 aggregate evidence。

    我们用 `min_clearance_max_to_bunny_top` 作为 penetration proxy。  
    它越接近 `0`，说明 cloth 的最高点越不容易整体落到 bunny 顶面以下。

    结果非常清楚：

    ### 1. `contact_collision_dist scale` 是最强主效应
    平均值从：
    - `scale=0.5` 的 `-2.90`
    提升到
    - `scale=2.0` 的 `-0.14`
    - `scale=3.0` 的 `-0.13`

    这说明：
    **当前 Newton `particle_radius` 太小** 是当前最强的 penetration 原因之一。

    ### 2. `bunny mass` 有帮助，但不是唯一根因
    - `bm=1` 平均大约 `-2.65`
    - `bm=1000` 大约 `-0.10`
    - `bm=5000` 大约 `-0.088`

    ### 3. `drop_height` 会恶化 penetration，但不是单独解
    - `h=0.1` 平均约 `-0.216`
    - `h=1.0` 平均约 `-1.98`

    ---

    ## Slide 13 — Best OFF Example

    这页是 baseline 和当前 best 的直接对照。

    baseline：
    - `ccd=1.0, h=0.1, bm=5000`

    current best：
    - `ccd=3.0, h=0.1, bm=5000`

    best 仍然没有完全消灭 penetration，  
    但已经把 `min_clearance_max_to_bunny_top` 拉到了大约 `-0.026`。

    这个数值已经非常接近“不穿”的边界了。

    所以这页的结论是：

    **现在最值得继续精修的是 `contact_collision_dist` / Newton `particle_radius`，而不是继续大范围扫质量。**

    ---

    ## Slide 14 — Which Hypotheses Survived

    这页把假设分成两类。

    ### 已验证 / 部分验证
    - `H1 bunny too light`：成立，但不是唯一主因
    - `H2 contact_collision_dist too small / particle_radius too small`：强成立
    - `H3 higher drop height worsens penetration`：成立

    ### 被弱化 / 被否定
    - “只要把 drop height 降低就能解决” → 不成立
    - “只要把 stiffness / damping 系数提很高就能 robust solve” → 不成立

    这里我补一条单独验证过的 negative result：
    - `body_ke/body_kd` 提高 `10x` 没稳健解决 penetration
    - `100x` 直接数值炸掉

    所以现在不应该继续往极硬 penalty 方向硬推。

    ---

    ## Slide 15 — What Progress We Actually Made

    我们这周真正的 progress 不是“又新增了很多视频”，而是三件更扎实的事：

    1. **把 ON 从 parity target 降级成 ablation**
       这让整个叙事不再自相矛盾。

    2. **把失败定位到 OFF external contact**
       说明我们现在要解决的是 particle-body contact，不是 self-collision。

    3. **用 64-run grid 把主导变量找出来**
       目前最强的杠杆是 ``contact_collision_dist`` scale。

    这是项目层面的实质性进展，因为我们已经从“问题很多”收敛到了“哪一个 knob 最值得继续推”。

    ---

    ## Slide 16 — Next Step

    下一步不需要再做大范围 brute-force 了。

    我建议直接收敛到：
    - `OFF`
    - `drop_height = 0.1`
    - `bunny_mass = 5000`

    然后只细扫：
    - `contact_collision_dist scale = 2.0 / 2.5 / 3.0 / 3.5`

    如果这一步找到一个几乎不穿的窗口，再做一小圈：
    - moderate `body_ke/body_kd`

    也就是说，我们现在已经不在大范围盲扫阶段，而是在做 **targeted refinement**。

    ---

    ## Slide 17 — Thank You

    结论再收一次：

    - `ON` 不能再被讲成 PhysTwin parity
    - `OFF` penetration 的主因已经被收敛到 external `contact_collision_dist` / Newton `particle_radius`
    - 我们已经拿到了下一步最值得精修的参数区间

    这就是本周最重要的 progress。
    """
)

TRANSCRIPT_MEETING = textwrap.dedent(
    f"""\
    # Meeting Transcript — PhysTwin → Newton Bridge

    语言：中文主讲 + English terminology  
    形式：按 slide 编号逐页口播  
    目标：对齐 `PhysTwin-Full Meetings.pptx` 的会议节奏，但内容更新为本周真实实验结果

    ---

    ## Slide 1 — PhysTwin-Full Meeting
    今天我汇报 PhysTwin 到 Newton bridge 的本周进展。  
    这次报告我会用 hypothesis-driven 的结构来讲，不只是展示结果，而是讲清楚我们验证了什么、证伪了什么。

    ## Slide 2 — Xinjie
    这一页只是 speaker page，我会直接进入 agenda。

    ## Slide 3 — This Week: 6 Hypotheses
    这周我们围绕 6 个 hypothesis 组织：
    1. H1A：PhysTwin deformables 能否和 novel rigid bodies 稳定交互。  
    2. H1B：Why cloth still penetrates bunny even when rigid mesh contact already works.  
    3. H1C：What do the latest low-mass cloth + bunny OFF runs actually show?  
    4. H2：为什么 Newton self-collision 不能被讲成 PhysTwin self-collision。  
    5. H3：多个 deformables 之间能否稳定互动。  
    6. H4：为什么 MPM 路线做不了真正 two-way coupling。

    ## Slide 4 — H1A: Deformable ↔ Rigid Body Interaction
    先讲 H1A。  
    H1A 是最基础的一层：如果连 deformable-rigid 都做不稳，后面的 cloth、multi-rope 和 MPM 比较都没有意义。

    ## Slide 5 — Bridge Pipeline: PhysTwin IR → Newton
    整个 pipeline 是：
    PhysTwin IR 导出 spring-mass object，然后在 Newton 里重建 particle、springs、ground 和 bunny mesh。  
    这里的重点不是“能不能读进去”，而是 physics semantics 有没有保住：
    topology、spring stiffness、mass、collision radius、drag 和 contact path 是否还对应原模型。

    ## Slide 6 — Newton Detects Particle-Mesh Contact
    这页用源码证明 Newton 是如何做 particle-shape detection 的。  
    关键点是：它按 `(particle, shape)` 对逐个查询，而且对 mesh 走的是 closest-point / sign-normal 路线。  
    这说明 bunny 在 Newton 里是一个 mesh surface，不是一团 particle cloud。

    ## Slide 7 — Newton Applies Force To The Rigid Body
    这页证明力是怎么施加的。  
    particle-body contact 不是一个简单的“位置投影”，而是 elastic + damping + friction 的 penalty force。  
    同时 rigid body 也会收到 reaction wrench，所以 interaction 是 bidirectional 的。

    ## Slide 8 — H1A Result: ✅ Validated
    结论是 H1A 已经成立。  
    rope / deformable 这条路说明：PhysTwin 的 deformable object 放进 Newton 后，已经能和 novel rigid bunny 做真实接触与载荷交换。

    ## Slide 9 — H1A Addendum: Strengthen Particle-Shape Contact
    这里我补一个非常直接的 engineering observation。

    在完全相同的 bunny center-drop 设置下，只改一个量：
    - `shape_contact_scale = 1x alpha`
    - 对比
    - `shape_contact_scale = 2x alpha`

    视频上可以直接看到：
    - `1x alpha` 时，bunny 顶部的局部 penetration 仍然更明显
    - `2x alpha` 时，external support 提前建立，cloth 更容易挂在 bunny 顶部

    这页的作用不是说 `2x alpha` 就是最终答案。  
    它只是先把一个事实讲清楚：
    - **particle-shape contact strength 是 deformable-rigid support 的直接强旋钮**

    ## Slide 10 — H1B: Why Cloth Still Penetrates Bunny
    这一段专门收口教授最关心的问题：
    - continuous mesh 为什么还会穿模？
    - bunny 和 thick box 为什么表现不同？
    - 能不能给出一个量化定义下的不穿模 demo？

    这一节不是继续铺现象，而是按 hypothesis 去做：
    1. 先从公式上说明 continuous mesh contact 为什么仍然依赖 particle radius  
    2. 再用 bunny vs thick box 验证局部几何是不是主因  
    3. 最后给出 bunny 自己的 quantitative working point

    ## Slide 11 — H1B Source Proof: Continuous Mesh Still Depends On Particle Radius
    这一页先不给结果，先把模型讲清楚。

    它专门回答一个被反复追问的问题：

    > rigid mesh 明明是 continuous surface，为什么 cloth-rigid 的结果还会依赖 particle spacing / particle radius？

    关键答案是：
    - Newton 的确把 mesh 当作 continuous surface 做 closest-point / sign-normal 查询；
    - 但参与接触判定的不是“点对曲面”，而是“**有半径的粒子球** 对曲面”。

    代码证据在 [kernels.py]({NEWTON_GEOM_KERNELS}#L1104-L1133)：
    - 对 mesh 先做 `mesh_query_point_sign_normal(...)`
    - 查询半径里直接有 `margin + radius`
    - 最终只有当 `d < margin + radius` 才会创建 soft contact

    也就是说：
    - surface 是 continuous 的
    - 但接触几何的有效厚度仍然由 `particle_radius` 决定

    所以后面如果粒子半径更小，cloth 中心可以离 mesh 更近才触发接触，实际就更容易出现可见 penetration。

    ## Slide 12 — H1B Baseline Comparison: Bunny Mesh vs Thick Box
    这一页先给老师看最直接的视频对照，而不是先堆指标。

    这里两个视频的设置完全一样：
    - 同一个 low-mass cloth
    - 同一个 rigid mass
    - 同一个 OFF contact pipeline
    - 只把 rigid geometry 从 bunny mesh 换成 thick box

    结果非常直观：
    - bunny baseline 里，局部 penetration 仍然明显可见
    - thick box baseline 里，overlap 会小很多，支撑 patch 也更稳定

    这页不是在说 thick box 已经完全解决了问题。  
    它要说明的是：
    - **局部几何形状确实会显著影响 penetration severity**
    - 所以后面继续讨论 bunny ear / thin patch 就有了明确的实验对照依据

    ## Slide 13 — H1B Thick-Box Working Point: Baseline vs 4x `shape_contact_scale`
    这一页开始给真正的 working point。

    thick box 的好处是几何简单，所以更适合先定义：
    - 什么叫 “解决了”
    - 什么叫 “仍然只是定性改善”

    这里我给的标准是：
    - `max penetration depth < 5 mm`
    - rigid body 不被异常打飞

    这里我先定义一下 `shape_contact_scale`。  
    它不是 Newton 原生变量，而是我们 bridge/demo 层加的实验参数，用来同时缩放 particle-shape 路径里的：
    - `soft_contact_ke/kd/kf`
    - `shape_material_ke/kd/kf`

    这里第一次出现的 `kf` 也说明一下：
    - `kf` = contact friction / tangential damping coefficient
    - 它控制切向接触响应，不是法向 stiffness

    这里的 `4x shape_contact_scale` 指的是：
    - particle-shape contact stiffness 相对于当前 low-mass 基线再提高 4 倍

    在 `4x shape_contact_scale` 下，当前 thick-box 结果已经达到：
    - `max penetration ≈ 3.8 mm`
    - `body_speed_max ≈ 0.0028 m/s`

    所以：
    - **thick box 已经可以作为量化定义下的 non-penetration cloth-rigid demo**

    ## Slide 14 — H1B Bunny Working Point: Baseline vs Best Current Setup
    box 解决以后，下一步就看 bunny 自己能不能进入同样的阈值。

    这里我选的 best case 不是继续把 cloth 投向耳朵区域，而是：
    - 把 cloth 中心移到更厚实的 `body-center`
    - `shape_contact_scale = 8x` relative to the current low-mass particle-shape contact baseline
    - 保留 `decoupled shape materials`

    这个 case 的结果是：
    - `max penetration ≈ 4.99 mm`
    - `body_speed_max ≈ 0.0156 m/s`

    也就是说，按刚才公开定义的标准：
    - bunny 这条线也已经进入 quantitative non-penetration 区间

    但这一页最重要的不是单个数字，而是让老师直接看到：
    - bunny 确实比 thick box 更难
    - 但并不是完全做不到

    ## Slide 15 — H1B Result: What We Can Claim Today
    到这里，H1B 可以明确收口成四句话：

    1. Newton 的 rigid mesh contact 的确是 continuous surface query  
    2. 但 contact force 仍然显式依赖 `particle_radius`，所以 particle spacing / `particle_radius` 会直接影响 penetration  
    3. thick box 更容易进入 non-penetration 区间，因为支撑 patch 更厚、更宽、更稳定  
    4. bunny 本身现在也已经有一个满足阈值的 best working point

    这就是这一节真正的交付，不再只是“现象有所改善”。

    ## Slide 16 — H1C: Latest Low-Mass Cloth + Bunny OFF Runs
    这页现在切到**最新重新补跑**的 `bm0p1 / bm0p5` 视频，不再引用之前坏掉的旧结果。

    这组实验是严格按统一规则重跑的：
    - self-collision OFF
    - `alpha = 0.1 / 4409`
    - `mass + spring_ke + spring_kd` 同时乘 `alpha`
    - `particle-shape contact` 也按同一个 `alpha` 缩
    - `shape_contact_damping_multiplier = 1`
    - 另外只改一个量：
      - `contact_dist_scale = 5.0`
    - 同时显式打开：
      - `decoupled shape materials`
      也就是保留 rigid-rigid 的 baseline `shape_material_*`，只缩 particle-shape 这条链

    这里 `contact_dist_scale = 5.0` 的物理意义是：
    - 更大的 `contact_collision_dist`
    - 经过 bridge 映射后也就是更大的 Newton `particle_radius`
    - 所以 cloth 会更早和 bunny mesh 建立接触壳

    新结果说明：
    - bunny 不再像之前错误结果那样沉到地面以下
    - `bm0p1` 和 `bm0p5` 的 bunny 顶部最终高度都基本保持在初始高度附近
    - 同时 cloth 对 bunny 的最终最大 clearance 也明显改善

    所以 `H1C` 这页现在的 takeaway 是：
    - **这里放的是最新、正确、可引用的 `bm0p1 / bm0p5` 结果**
    - **它们证明在 decoupled alpha-scaling 下，rigid-rigid support 可以保住**
    - **增大 `particle_radius` / `contact_collision_dist` 会进一步减轻 bunny mesh penetration**

    ## Slide 17 — H1B Source Proof: Why Particle Radius Matters in OFF
    这页把 H1B 的关键机理钉死。

    我们证明了两件事：
    1. detection 里 `margin + particle_radius` 会直接决定何时触发 mesh contact；  
    2. force 里 `c = dot(n, px - bx) - particle_radius` 会直接决定 penalty depth。  

    所以 `particle_radius` 不是视觉参数，而是 Newton 接触模型里的真实几何参数。  
    这也就是为什么同样是 continuous mesh，particle spacing / `particle_radius` 仍然会显著影响 penetration。

    ## Slide 18 — H1B OFF Experiment Matrix
    这一页讲 H1B 的实验设计。  
    OFF case 里我们固定 self-collision 关闭，只扫三个轴：
    - ``contact_collision_dist`` scale
    - `drop height`
    - `bunny mass`

    这样做的目的，是把 external particle-body contact 的主导因素单独抽出来，而不再和 self-collision 混在一起。

    ## Slide 19 — H1B Hypothesis 1: Bunny Mass Matters
    这页只回答一个问题：bunny mass 会不会影响 penetration。  
    结果是会。  
    bunny 越轻，越容易被 cloth 顶走；但把 bunny 变重之后，penetration 并不会自动消失。

    所以 bunny mass 是重要变量，但它不是唯一主因。

    ## Slide 20 — H1B Hypothesis 2: Height Matters, But Not Alone
    这页回答第二个问题：是不是只要 drop height 降低就能解决。  
    结果是不能。  
    高度越高 penetration 越重，但低高度仍然会穿，所以它是 aggravating factor，不是 root fix。

    ## Slide 21 — H1B Hypothesis 3: `contact_collision_dist` Is The Strongest Lever
    这页是 H1B 里最关键的实验对照。  
    在固定低高度和大 bunny mass 下，增大 ``contact_collision_dist`` scale 的效果最明显。  
    这说明当前外部接触里，最强的主导因素就是 `contact_collision_dist`，也就是 bridge 映射到 Newton 的 `particle_radius`。

    ## Slide 22 — H1B Result: `contact_collision_dist` Is The Strongest Lever
    这里把单个视频观察汇总成 aggregate trend。  
    64-run 结果显示：
    ``contact_collision_dist`` scale 的主效应比 bunny mass 和 height 都更强。  
    这就是 H1B 这一节最重要的 quantitative conclusion。

    ## Slide 23 — H1B Best OFF Example
    这一页是 H1B 的可视化收口。  
    baseline 和 current best 的区别很直观：最佳区间已经更接近“被 bunny 稳定支撑”，虽然还没有完全消灭 penetration。  
    这说明问题已经从“广泛失败”推进到了“接近可用区间”。

    ## Slide 24 — H1B What This Week Actually Proved
    这里总结 H1B 已验证和已证伪的假设：
    - bunny mass 有用但不够  
    - 单独把 cloth 总质量降到 `0.1 kg` 也不够  
    - 更大的 ``contact_collision_dist`` / Newton `particle_radius` 明显减轻 penetration  
    - 只降高度不够  
    - 单纯暴力增加 stiffness / damping 系数不是稳健解

    到这里，H1B 负责的问题已经非常清楚：
    - 为什么 cloth 仍然会穿 bunny
    - 在 OFF case 里，真正主导 penetration 的变量是什么

    ## Slide 25 — H2: Why Newton Self-Collision ≠ PhysTwin
    进入 H2。  
    这一节不再承担 OFF penetration 的归因工作。  
    它只回答一个问题：
    **为什么 Newton native self-collision 不能被讲成 PhysTwin self-collision，以及为什么不能简单把 self-collision 打开当作解决方案。**

    ## Slide 26 — H2 Terms: Gate, Update, Force
    这一页先把三个 terminology 定死。

    第一，`approach-gated`。  
    它的意思是：pair 不仅要够近，而且 relative velocity 还必须沿连线方向继续相向。

    第二，`impulse update`。  
    它的意思是：源码直接写 velocity，例如 `v_new = ...`，而不是先累积 force。

    第三，`penalty force`。  
    它的意思是：先把 overlap 转成 force，再由 solver integration 把 force 积成新速度。

    这三个定义后面都会对应到源码原文，不是口头概念。

    ## Slide 27 — H2 Newton Source: eval_particle_contact (Self-Collision)
    这一页先单独给 Newton native self-collision 的入口源码。  
    它本质上是 generic particle-particle penalty contact，不是 cloth-topology-aware self-collision。  
    这点非常重要，因为它决定了后面为什么 ON case 不能简单等同于 PhysTwin。

    ## Slide 28 — H2 Source Proof: PhysTwin vs Newton Self-Collision
    这页只引用几段最关键的原文，不铺长代码。

    PhysTwin 侧，trigger gate 是：
    - `dis_len < collision_dist`
    - `wp.dot(dis, relative_v) < -1e-4`

    也就是：只有接近中的 pair 才进 collision。  
    update 则是：
    - `J_average = J_sum / valid_count`
    - `v_new[tid] = v1 - J_average / m1`

    也就是说，PhysTwin 是 **approach-gated + impulse-style velocity correction**。

    Newton 侧，trigger 条件是：
    - `err = d - radius - particle_radius[index]`
    - `if err <= k_cohesion:`

    它只看 overlap / cohesion，不检查是否仍在 approaching。  
    update 则是：
    - `f += particle_force(...)`

    所以 Newton 是 **overlap-triggered + penalty-force accumulation**。

    这页的结论就是：
    - PhysTwin 和 Newton 这里是 model difference，不是 tuning difference。

    ## Slide 29 — H2 Result: Different Model, Different Claim
    所以 H2 这里我们只做一个 claim：

    - Newton native self-collision 和 PhysTwin self-collision 不是同一个模型
    - 因此 `self-collision ON` 不能作为 parity evidence
    - 在这套 slides 里，`self-collision ON` 只保留为 Newton-native ablation

    ## Slide 30 — H2 ON Ablation Demo: Why We Cannot Simply Turn Self-Collision ON
    这一页直接放 ON case 的视频，不再只讲抽象结论。

    目的就是给老师看：
    - 即使已经把 H1B 的 external contact 问题分离出来
    - 只要简单打开 Newton native self-collision
    - 结果仍然会出现明显的 particle-cloud instability

    这页的作用非常直接：
    - 不是证明 ON 完全不可研究
    - 而是证明 **当前不能把 “打开 self-collision ON” 当作一个简单的工程解**

    ## Slide 31 — H3: Multi-Deformable Interaction
    进入 H3。  
    这个问题是：多个 deformables 能不能彼此稳定互动，而不是只和 rigid 互动。

    ## Slide 29 — One Rope + Grounded Sloth: Design Decision
    我们的设计选择是：  
    不直接信任 native all-particle contact，因为那会让 same-object neighbors 和 spring 打架。  
    所以我们统一采用 selective pair filtering，然后分别展示：
    - `rope + rope`
    - `rope + sloth`
    两个最小 demo。

    ## Slide 30 — Custom Kernel: _eval_cross_object_contact
    这页用源码证明 H3 的实现方式。  
    我们没有换物理模型，而是在 Newton 的 HashGrid + penalty 框架里，只筛出 rope-vs-sloth pairs。  
    这样就能保住 Newton native 风格，同时避免 same-object artifact。

    ## Slide 31 — H3 Demos: Rope + Rope and Rope + Sloth
    这页视频现在同时放两个 demo：
    - 左边：`rope + rope on ground`
    - 右边：`rope + grounded sloth`

    这样老师可以直接看到：
    - 同类 slender deformables 之间能互动
    - rope 和更复杂的 deformable object 之间也能互动
    而且两条线都不再依赖 rigid bunny。

    ## Slide 32 — H3 Result: ✅ Validated
    所以 H3 现在是成立的。  
    Multi-deformable interaction 是可做的，只是需要 selective pair filtering，而不是全开 native particle contact。  
    现在 slides 里同时有：
    - `rope + rope`
    - `rope + sloth`
    两条最小展示线。

    ## Slide 33 — H4: MPM Two-Way Coupling
    进入 H4。  
    这部分不是说 MPM 不可用，而是问：它适不适合我们这条 spring-mass two-way 路线。

    ## Slide 34 — MPM: Kinematic Collider Works
    这页先讲 positive result。  
    MPM one-way kinematic collider 是能工作的，sloth / zebra 推沙子都没有问题。  
    所以我们不是说 MPM 全部失败，而是说 failure 出在 two-way support 上。

    ## Slide 35 — Why MPM ≠ Two-Way with Spring-Mass
    这页是理论解释。  
    关键问题有三个：
    1. mass ratio  
    2. ground-plane SDF competition  
    3. patch mesh quality  
    所以这个问题更像 architecture mismatch，而不是小 bug。

    ## Slide 36 — H4 Result: ❌ Infeasible (by design)
    这页给最终判断。  
    在当前项目目标下，MPM two-way 不是值得继续主推的方向。  
    SemiImplicit penalty contact 才是主线。

    ## Slide 37 — Next Step
    现在下一步很明确：
    不再做大范围 brute-force，而是缩到最 promising 的区域，细扫 ``contact_collision_dist`` scale，再做温和的 stiffness tuning。

    ## Slide 38 — Summary
    最后收口：
    H1 deformable-rigid interaction 已验证，  
    H1B 已经把 bunny penetration 主因收敛到 external `contact_collision_dist` / Newton `particle_radius`，  
    H1C 已经给出最新的 decoupled alpha-scaled light-bunny reference runs，  
    H2 已经用源码证明 Newton self-collision 不是 PhysTwin self-collision，  
    H3 multi-deformable interaction 已验证，  
    H4 MPM two-way 这条路被证伪。

    这就是本周 meeting 最想讲清楚的 progress。
    """
)

TRANSCRIPT_MEETING += textwrap.dedent(
    f"""

    ---

    # Additional Background Notes For Q&A

    这一部分不是要全部逐字讲出来，而是给我在 meeting 里防追问用的。  
    原则是：
    - Slide 上只保留 15 秒内能读完的核心句子
    - Transcript 里保留更多背景、定义、以及“如果教授追问，我该怎么解释”

    ---

    ## H1 Background — Why Start From Deformable ↔ Rigid

    为什么我们把 H1 放在最前面？  
    因为它是整个 bridge 的最小可验证闭环。

    如果一个 PhysTwin spring-mass object 导入 Newton 之后，连最基础的 rigid body interaction 都做不稳，那么：
    - 后面的 cloth self-collision 分析没有意义
    - multi-deformable interaction 也没有意义
    - 更不要说 MPM two-way 这种更重的链路

    所以 H1 的角色不是“又一个 demo”，而是 bridge correctness 的第一层门槛。

    这里教授可能会问：
    - 为什么不直接拿 cloth 开始？

    我的回答应该是：
    - cloth 引入了 self-collision、面状支撑、局部折叠等额外因素
    - rope / generic deformable ↔ rigid 更适合先验证基础 contact pipeline
    - 先证明 particle-shape interaction 没有语义错误，再去处理 cloth 的额外复杂性

    ---

    ## H1 Background — What Counts As “Validated”

    当我说 H1 validated，不是说 Newton 和 PhysTwin 在 every-frame 轨迹上完全一致。  
    我说的 validated 是：

    1. IR 里的 object topology、mass、collision radius、spring parameters 已经成功转到 Newton
    2. Newton 的 particle-shape contact pipeline 能稳定地让 deformable 和 novel rigid bunny 发生双向作用
    3. 这个 interaction 在视频上是可解释的，不是数值爆炸偶然呈现出的视觉假象

    这里要特别强调：
    - 我们做的是 **bridge validation**
    - 不是宣称 “PhysTwin and Newton are the same simulator”

    这两件事必须分开。

    ---

    ## Source Note — Why `create_soft_contacts` Matters

    [create_soft_contacts]({NEWTON_GEOM_KERNELS}#L1013) 这一页的价值不是“秀代码”，而是明确三个物理事实：

    1. contact detection 是 per `(particle, shape)` pair 做的  
    2. rigid bunny 作为 mesh shape 被单独查询  
    3. particle 和 shape 的 world assignment 是明确区分的

    也就是说，deformable-rigid interaction 在 Newton 里不是黑盒。  
    我们知道：
    - 谁在检测
    - 检测条件是什么
    - mesh 是如何被访问的

    如果教授追问：
    - 这是 broad phase 还是 narrow phase？

    我应该回答：
    - 对 particle-shape 这条 soft contact 路径来说，核心 narrow-phase 语义就在 `create_soft_contacts`
    - 它把 particle 变换到 shape local space，然后根据 shape type 分支处理
    - 对 mesh 就是 nearest-point / signed-normal query

    ---

    ## Source Note — Why `eval_particle_body_contact` Matters

    [eval_particle_body_contact]({NEWTON_KERNELS_CONTACT}#L156) 这一页的价值，是把“接触存在”变成“接触力怎么施加”。

    这里教授可能会问：
    - Newton 的 deformable-rigid contact 是 constraint-based 还是 penalty-based？

    我应该非常明确地回答：
    - 在这条 SemiImplicit 路线上，它是 **penalty-based**
    - 法向弹性项 `fn = n * c * ke`
    - 阻尼项 `fd = n * min(vn,0) * kd`
    - 切向摩擦项 `ft`

    所以这个系统天然允许：
    - 先发生一定 penetration
    - 再用力往回推

    这也是为什么 penetration 不能简单用“有没有 bug”来理解，很多时候它是 penalty stiffness / `particle_radius` / timestep 三者共同决定的。

    ---

    ## H2 Background — Why Self-Collision Became The Main Story

    H2 之所以成为主线，不是因为 self-collision 比 rigid interaction 更重要，而是因为：

    - H1 已经说明 bridge 能做最基本 interaction
    - 真正还没解释清楚的 failure，集中在 cloth
    - cloth 把两个问题叠在一起：
      1. self-collision
      2. external support on bunny

    所以如果我不先把 self-collision 的叙事边界说清楚，后面所有 cloth 结果都会被讲混。

    这就是为什么我们必须先下一个方法论判断：
    - `ON` 不是 PhysTwin parity
    - `OFF` 才是当前 mainline engineering target

    ---

    ## H2 Background — Why ON Cannot Be A Direct Parity Claim

    这里我需要准备一个更强的口头版本，因为教授很可能会问：
    - 你为什么说不能直接迁移？
    - 是参数没调好，还是模型真的不同？

    我的回答应该分三层：

    ### 第一层：触发条件不同
    PhysTwin 要求：
    - `dis_len < collision_dist`
    - 且 `dot(dis, relative_v) < -1e-4`
    见 [spring_mass_warp.py]({PHYSTWIN_SIM}#L196)

    Newton native self-contact 则是：
    - 只要 `err <= k_cohesion` 就加 penalty force
    见 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L88)

    也就是说：
    - PhysTwin 是 approaching-only
    - Newton native 不是

    ### 第二层：求解变量不同
    PhysTwin 修正 velocity  
    Newton 积累 force

    这意味着即使数值上调到“表面上接近”，底层物理解释仍然不同。

    ### 第三层：cloth topology awareness 不同
    Newton native particle contact 不知道 cloth 的结构邻居、1-ring、2-ring、surface adjacency。  
    它把问题当 generic particle soup。

    所以我不是说：
    - “这条路完全不能用”
    而是说：
    - “这条路不能再被宣称成 PhysTwin self-collision parity”

    ---

    ## H2 Background — Why OFF Became The Main Engineering Target

    OFF case 的价值在于：
    - 它去掉了 self-collision 这个 confounder
    - 只看 cloth ↔ bunny external contact

    这样我们就能把问题写成一个更干净的形式：

    > 给定一张没有自碰撞的 cloth，为什么它仍然会穿过 bunny？

    这时候若还出现穿模，就说明：
    - 主问题已经不在 self-collision
    - 而在 external penalty contact

    这一步对于项目很关键，因为它把“布料系统很复杂”拆成了：
    - self-contact 问题
    - external contact 问题

    两条线可以分别定位。

    ---

    ## Source Note — Why Particle Radius Matters Against A Mesh

    这个点老师很可能追问，我需要讲得比 slide 上更细。

    常见误解是：
    - rigid bunny 是 mesh，所以 particle radius 不应该重要

    但源码已经清楚说明，particle radius 重要，而且是两次重要。

    ### 第一次重要：检测
    在 [create_soft_contacts]({NEWTON_GEOM_KERNELS}#L1113) 里，mesh query 的搜索半径直接用了：
    - `margin + radius / min_scale`

    然后在 [create_soft_contacts]({NEWTON_GEOM_KERNELS}#L1132) 又判断：
    - `if d < margin + radius`

    所以 radius 会决定：
    - 最近点查询能搜多远
    - contact 何时被认定为 active

    ### 第二次重要：受力
    在 [eval_particle_body_contact]({NEWTON_KERNELS_CONTACT}#L216)：
    - `c = dot(n, px - bx) - particle_radius`

    所以在相同 particle center 位置下，radius 越大：
    - `c` 越早进入有效接触区
    - `fn = n * c * ke` 越早变大

    这就是为什么 `particle_radius` 对 rigid mesh contact 绝不是一个渲染参数，而是 Newton 接触模型里的真实几何参数。

    ---

    ## H2 Background — Why Rope Looked Easier

    这一点也很容易被问：
    - 为什么 rope 的现象相对没那么严重？

    更完整的回答是：

    1. rope 是 1D bead-chain，cloth 是 2D sheet  
    2. rope 不需要 bunny 提供大面积支撑  
    3. rope 的有效 Newton `particle_radius` 更大  
    4. rope 的视觉判据也更宽松

    举例说：
    - rope 只要能挂、绕、滑，视觉上就像“工作了”
    - cloth 则要求 bunny 真的把整片布托住

    所以同样的 penalty softness，在 rope 上可能只表现成“有点软”，在 cloth 上就直接暴露成 penetration。

    ---

    ## H2 Experiment Design — Why The 64-run Grid Was Necessary

    如果只跑一两个视频，很容易陷入叙事偏见。  
    所以这周我们做 64-run full grid，不是为了“多跑一些”，而是为了区分主效应。

    三个轴的设计逻辑是：

    ### ``contact_collision_dist`` scale
    测试 `contact_collision_dist` / `particle_radius` 假说  
    如果这个轴最强，就说明 detection/support 本身太弱

    ### `drop height`
    测试 kinetic energy 假说  
    如果高度一降就好，那问题更像 tunneling / impact velocity

    ### `bunny mass`
    测试 rigid support 假说  
    如果 bunny 太轻会明显恶化，说明 rigid recoil 是 penetration 的一部分来源

    这三个轴是 orthogonal enough 的，所以可以帮助我们做 first-order attribution。

    ---

    ## H2 Results — What The 64-run Grid Really Tells Us

    这里要比 slide 上更细一点，因为教授可能会问：
    - 为什么你说 ``contact_collision_dist`` / `particle_radius` 是最强主效应？

    我的回答应该是：

    - 当 ``contact_collision_dist`` scale 从 `0.5` 提到 `2.0 / 3.0` 时，
      `min_clearance_max_to_bunny_top` 的平均改善幅度远大于 mass 和 height
    - `bunny mass` 的确有帮助，但从 `1000 -> 5000` 的收益已经在变小
    - `drop height` 明显恶化 penetration，但不能单独修复问题

    这三条一起说明：

    > 当前 penetration 的第一主导项是 external `contact_collision_dist` / Newton `particle_radius`，  
    > 不是 bunny mass，也不是 drop height 单独决定。

    ---

    ## H2 Results — Why “Lower Height Alone” Was Falsified

    这是一个很容易被误会的点。  
    我们不是说 height 不重要，而是说：

    - **只调低 height 不能把问题从 fail 变成 pass**

    理由是：
    - 低高度时确实 penetration 变轻
    - 但 cloth 仍然会整体低于 bunny 顶面

    所以正确表述应该是：
    - `drop height` 是 aggravating factor
    - 不是 root fix

    这个 phrasing 在答辩时很重要，因为它避免把 negative result 讲成“height 没用”这种过头的结论。

    ---

    ## H2 Results — Why “Just Increase Stiffness / Damping” Was Falsified

    这个也是常见建议，所以我们必须提前准备。

    我们单独做过：
    - `10x` stiffness/damping
    - `100x` stiffness/damping

    结果是：
    - `10x` 没有稳健消除 penetration
    - `100x` 直接数值爆炸

    所以对这个系统来说，问题不是简单 “spring 不够硬”，而是：
    - penalty support 本身的触发时机
    - `particle_radius`
    - 和系统能量尺度的相对关系

    换句话说：
    - 太软会穿
    - 太硬会炸
    - 所以外部 contact 不是一个单旋钮问题

    ---

    ## H2 Progress Statement — What We Actually Achieved

    这一段在 meeting 里很重要，因为它决定老师最后怎么评价“这周有没有实质进展”。

    我们真正达成的 progress 是：

    1. **把 ON 从 parity story 中剥离出来**  
       这让叙事不再自相矛盾。

    2. **把失败具体定位到 OFF external `contact_collision_dist` / Newton `particle_radius`**  
       不是泛泛地说 “cloth 很难”。

    3. **用 64-run grid 证明了主导因子排序**  
       `contact_collision_dist` / `particle_radius` > bunny mass > drop height

    这三点合起来，说明项目已经从“现象观察”推进到了“机制归因”。

    ---

    ## H3 Background — Why Multi-Deformable Needed A Custom Path

    H3 里最值得强调的是：
    - 我们不是因为写 custom kernel 就背离 Newton
    - 而是为了在 Newton 现有 penalty 框架下，修正一个 pair-selection 问题

    native all-particle contact 对 `rope + sloth` 这种 multi-object case 的问题在于：
    - same-object 邻居和 cross-object 邻居被一视同仁
    - 这会让 rope 和 sloth 的内部结构也被错误地强排斥

    所以 custom kernel 的本质是：
    - 保持 Newton penalty 模型
    - 只改 pair filtering

    这是一个比“重写物理”温和得多的 intervention。

    ---

    ## H3 Result — Why It Counts As Validation

    H3 成立，不是因为“视频好看”，而是因为：

    1. rope 和 sloth 确实产生了 cross-object interaction  
    2. same-object artifact 被抑制住了  
    3. 结果在更干净的 ground-only 场景里仍然可读、稳定

    也就是说，这不是一个纯 render trick，而是 pair filtering 后的真实 contact behavior。

    ---

    ## H4 Background — Why MPM Looked Tempting

    H4 之所以值得单列，是因为它一开始在直觉上很有吸引力：
    - MPM 对 sand / granular interaction 很自然
    - one-way kinematic collider 已经工作

    所以很容易产生一个直觉：
    - 那是不是 two-way 也只差一点点？

    这周最重要的 H4 贡献，就是证明：
    - 这个直觉是错的

    ---

    ## H4 Background — Why Two-Way Fails By Design

    更完整的解释是：

    1. **mass ratio**  
       spring-mass body 太重，sand 给不出对应 support

    2. **ground-plane SDF competition**  
       动态 patch collider 还没接管 contact，ground 就先赢了

    3. **patch mesh quality / locality**  
       patch collider 太局部，不像真正的 rigid support surface

    所以 H4 的 failure 更像是：
    - architectural mismatch
    而不是：
    - 参数没调好

    ---

    ## Closing Notes For Q&A

    如果老师最后问：
    - 你们这一周最关键的 insight 是什么？

    我应该回答：

    > 最关键的 insight 是我们把“cloth 会穿 bunny”这件事，从一个笼统现象，收敛成了一个更具体的机制判断：  
    > 当前主问题是 OFF external penalty contact 里的 `contact_collision_dist` / `particle_radius` 太小；  
    > 同时 ON 不能再被讲成 PhysTwin parity，因为两边 self-collision 模型根本不同。

    如果老师问：
    - 下一周最值得做的事是什么？

    我应该回答：

    > 在当前最 promising 区域，继续精扫 `contact_collision_dist`；  
    > 而不是再去大范围 brute-force 其它参数。  
    > 这会让我们更快判断：这个方向到底能不能达到一个“几乎不穿”的工作点。
    """
)


def build_deck_reference(
    prs: Presentation, rows: list[dict], gif_dir: Path, image_dir: Path
) -> None:
    _delete_all_slides(prs)

    factor_chart = _make_factor_chart(rows, image_dir / "factor_chart.png")
    phys_png = _code_excerpt_image(
        PHYSTWIN_SIM,
        "phystwin_self_collision_ref",
        _compact_signature_lines(PHYSTWIN_SIM, 261, 272)
        + [""]
        + _extract_code_segments(
            PHYSTWIN_SIM,
            [(196, 200), (205, 206), (220, 220), (294, 296)],
            highlight_lines={196, 198, 199, 205, 220, 294, 295},
        ),
        image_dir,
    )
    create_soft_png = _code_excerpt_image(
        NEWTON_GEOM_KERNELS,
        "newton_create_soft_contacts_ref",
        _compact_signature_lines(NEWTON_GEOM_KERNELS, 1014, 1038)
        + [""]
        + _extract_code_segments(
            NEWTON_GEOM_KERNELS,
            [(1013, 1014), (1112, 1115), (1132, 1134)],
            highlight_lines={1014, 1113, 1114, 1132, 1133},
        ),
        image_dir,
    )
    particle_body_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_eval_particle_body_contact_ref",
        _compact_signature_lines(NEWTON_KERNELS_CONTACT, 156, 186)
        + [""]
        + _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(215, 223), (244, 248), (269, 273)],
            highlight_lines={216, 222, 223, 245, 248, 271, 273},
        ),
        image_dir,
    )
    newton_self_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_self_collision_ref",
        _compact_signature_lines(NEWTON_KERNELS_CONTACT, 48, 62)
        + [""]
        + _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(79, 80), (83, 94)],
            highlight_lines={80, 88, 90, 94},
        ),
        image_dir,
    )
    mesh_query_png = _code_excerpt_image(
        NEWTON_GEOM_KERNELS,
        "newton_mesh_query_ref",
        _extract_code_segments(
            NEWTON_GEOM_KERNELS,
            [(1112, 1115), (1132, 1132)],
            highlight_lines={1113, 1114, 1115, 1132},
        ),
        image_dir,
    )
    newton_contact_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_contact_force_ref",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(215, 223), (244, 248), (263, 267)],
            highlight_lines={216, 222, 223, 245, 248, 263, 265},
        ),
        image_dir,
    )
    cross_rope_png = _code_excerpt_image(
        ROOT / "Newton" / "phystwin_bridge" / "demos" / "demo_rope_sloth_ground_contact.py",
        "cross_object_contact_ref",
        _compact_signature_lines(
            ROOT / "Newton" / "phystwin_bridge" / "demos" / "demo_rope_sloth_ground_contact.py",
            45,
            59,
        )
        + [""]
        + _extract_code_segments(
            ROOT / "Newton" / "phystwin_bridge" / "demos" / "demo_rope_sloth_ground_contact.py",
            [(45, 53), (67, 68), (80, 81), (87, 93)],
            highlight_lines={45, 51, 52, 67, 80, 87, 91},
        ),
        image_dir,
    )

    rigid_h1_gif = _ensure_gif(CURRENT_ROPE_BUNNY_MP4, gif_dir, stem_override="rope_bunny_current")
    rope_ground_gif = _ensure_gif(
        ROOT / "tmp" / "rope_ground_drop_full_close_20260322" / "rope_ground_drop.mp4",
        gif_dir,
        stem_override="rope_ground_drop",
    )
    multi_rope_gif = _ensure_gif(
        ROOT / "tmp" / "rope_sloth_dualsource_drop1p0_r2_20260322" / "rope_sloth_ground_drop.mp4",
        gif_dir,
        stem_override="rope_sloth_ground_drop",
    )
    mpm_sloth_gif = _ensure_gif(MPM_SLOTH_MP4, gif_dir, stem_override="sloth_mpm_sand_one_way_wide")

    baseline_gif = _ensure_gif(BASELINE_OFF_MP4, gif_dir, stem_override="baseline_off_ccd1_h0p1_bm5000")
    best_gif = _ensure_gif(BEST_OFF_MP4, gif_dir, stem_override="best_off_ccd3_h0p1_bm5000")
    mass_light_gif = _ensure_gif(MASS_LIGHT_MP4, gif_dir, stem_override="mass_bm1_ccd1_h0p1")
    mass_mid1_gif = _ensure_gif(MASS_MID1_MP4, gif_dir, stem_override="mass_bm10_ccd1_h0p1")
    mass_mid2_gif = _ensure_gif(MASS_MID2_MP4, gif_dir, stem_override="mass_bm1000_ccd1_h0p1")
    mass_heavy_gif = _ensure_gif(MASS_HEAVY_MP4, gif_dir, stem_override="mass_bm5000_ccd1_h0p1")
    height_low_gif = _ensure_gif(HEIGHT_LOW_MP4, gif_dir, stem_override="height_h0p1_ccd1_bm5000")
    height_mid1_gif = _ensure_gif(HEIGHT_MID1_MP4, gif_dir, stem_override="height_h0p25_ccd1_bm5000")
    height_mid2_gif = _ensure_gif(HEIGHT_MID2_MP4, gif_dir, stem_override="height_h0p5_ccd1_bm5000")
    height_high_gif = _ensure_gif(HEIGHT_HIGH_MP4, gif_dir, stem_override="height_h1p0_ccd1_bm5000")
    ccd_small_gif = _ensure_gif(CCD_SMALL_MP4, gif_dir, stem_override="ccd_s0p5_h0p1_bm5000")
    ccd_mid1_gif = _ensure_gif(CCD_MID1_MP4, gif_dir, stem_override="ccd_s1p0_h0p1_bm5000")
    ccd_mid2_gif = _ensure_gif(CCD_MID2_MP4, gif_dir, stem_override="ccd_s2p0_h0p1_bm5000")
    ccd_large_gif = _ensure_gif(CCD_LARGE_MP4, gif_dir, stem_override="ccd_s3p0_h0p1_bm5000")
    advisor_bm0p1_gif = _ensure_gif(ADVISOR_BM0P1_MP4, gif_dir, stem_override="advisor_bm0p1")
    advisor_bm0p5_gif = _ensure_gif(ADVISOR_BM0P5_MP4, gif_dir, stem_override="advisor_bm0p5")
    h2_on_failure_gif = _ensure_gif(H2_ON_FAILURE_MP4, gif_dir, stem_override="h2_on_failure_s0p5")
    h1_bunny_radius_1x_gif = _ensure_gif(P0_BUNNY_OFF_MP4, gif_dir, stem_override="h1_bunny_radius_1x")
    h1_bunny_radius_2x_gif = _ensure_gif(H1_BUNNY_RADIUS_2X_MP4, gif_dir, stem_override="h1_bunny_radius_2x")
    p0_bunny_gif = _ensure_gif(P0_BUNNY_OFF_MP4, gif_dir, stem_override="p0_bunny_off_decoupled")
    p0_bunny_final_gif = _ensure_gif(P0_BUNNY_FINAL_MP4, gif_dir, stem_override="p0_bunny_bodycenter8x_nopen")
    p0_box_baseline_gif = _ensure_gif(P0_BOX_BASELINE_MP4, gif_dir, stem_override="p0_box_off_baseline")
    p0_box_gif = _ensure_gif(P0_BOX_OFF_MP4, gif_dir, stem_override="p0_box_off_nopen")
    p0_box_summary = _load_json_or_fallback(
        P0_BOX_OFF_SUMMARY,
        {
            "max_penetration_depth_rigid_m": 0.0038,
            "body_speed_max": 0.0028,
        },
    )
    p0_bunny_final_summary = _load_json_or_fallback(
        P0_BUNNY_FINAL_SUMMARY,
        {
            "max_penetration_depth_bunny_mesh_m": 0.00499,
            "body_speed_max": 0.0156,
        },
    )
    p0_bunny_sweep_png = _make_bunny_sweep_chart(
        P0_BUNNY_SWEEP_CSV, image_dir / "bunny_mesh_sweep_metrics.png"
    )
    bridge_flow_png = _make_bridge_pipeline_diagram(image_dir / "bridge_pipeline_diagram.png")

    _title_slide(prs, "PhysTwin-Full Meeting", [])
    _title_slide(prs, "Xinjie", [date.today().isoformat()])
    _body(prs, "This Week: 6 Hypotheses", [
        "H1A  Can PhysTwin deformables interact with novel rigid bodies?",
        "H1B  Why cloth still penetrates bunny even when mesh contact works",
        "H1C  What do the latest low-mass cloth + bunny OFF runs actually show?",
        "H2   Why Newton self-collision is not PhysTwin self-collision",
        "H3   Can multiple deformables interact with each other?",
        "H4   Why can't MPM do two-way coupling with spring-mass?",
    ])

    _section(prs, "H1A: Deformable ↔ Rigid Body Interaction")
    _picture_with_footer(prs, "Bridge Pipeline: PhysTwin IR → Newton", bridge_flow_png, [
        "Read left to right: PhysTwin outputs → IR v2 → Newton reconstruction → interaction layer.",
    ])
    _picture_with_footer(prs, "Newton Source: create_soft_contacts", create_soft_png, [
        "Proves: particle-shape contacts are detected per (particle, shape) pair",
        "This is how a deformable object sees an arbitrary rigid bunny mesh",
    ])
    _picture_with_footer(prs, "Newton Source: eval_particle_body_contact", particle_body_png, [
        "Proves: contact force is elastic + damping + friction",
        "The rigid body receives the equal-and-opposite reaction wrench",
    ])
    _gif_single(prs, "H1A Result: ✅ Validated", rigid_h1_gif, "Deformable-rigid interaction is already working in Newton.")
    _gif_compare_large(
        prs,
        "H1A Addendum: Stronger Particle-Shape Contact Helps",
        "Baseline `shape_contact_scale` = 1x alpha",
        h1_bunny_radius_1x_gif,
        "Scaled `shape_contact_scale` = 2x alpha",
        h1_bunny_radius_2x_gif,
        footer="self-collision OFF | cloth total mass = 0.1 kg | bunny mass = 0.5 kg | center drop",
    )
    _section(prs, "H1B: Why Cloth Still Penetrates Bunny")
    _pic_twocol(
        prs,
        "H1B Source Proof: Continuous Mesh Still Depends On Particle Radius",
        "Detection: continuous mesh query still starts at d < margin + radius",
        mesh_query_png,
        "Force: penalty penetration still uses c = signed distance - particle radius",
        newton_contact_png,
    )
    _gif_compare_large(
        prs,
        "H1B Baseline Comparison: Bunny Mesh vs Thick Box",
        "Bunny baseline: same cloth, local penetration remains visible on the curved top patch",
        p0_bunny_gif,
        "Thick-box baseline: same cloth, broader support patch already reduces penetration",
        p0_box_baseline_gif,
        footer="Same low-mass cloth | self-collision OFF | rigid mass = 0.5 kg | decoupled shape materials",
    )
    _gif_compare_large(
        prs,
        "H1B Thick-Box Working Point: Baseline vs 4x `shape_contact_scale`",
        "Baseline thick box: visible overlap remains because penalty contact is still soft",
        p0_box_baseline_gif,
        "4x `shape_contact_scale`: the cloth stays on top and rigid recoil remains negligible",
        p0_box_gif,
        footer="Here `shape_contact_scale` means the bridge-side scale applied to Newton particle-shape contact stiffness.",
    )
    _gif_compare_large(
        prs,
        "H1B Bunny Working Point: Baseline vs Best Current Setup",
        "Bunny baseline: center drop, baseline `shape_contact_scale`, local penetration still visible",
        p0_bunny_gif,
        "Best current bunny case: body-center target plus 8x `shape_contact_scale` reduces penetration below the threshold",
        p0_bunny_final_gif,
        footer="The bunny case is harder than the box case because the local support patch is thinner and more curved",
    )
    _boxed_bullets(
        prs,
        "H1B Result: What We Can Claim Today",
        "Success Definition",
        [
            "Define solved as: max penetration depth < 5 mm and rigid body stays stable",
            f"Current thick-box demo: {p0_box_summary['max_penetration_depth_rigid_m']*1000:.1f} mm max penetration, {p0_box_summary['body_speed_max']:.4f} m/s rigid speed",
            f"Current bunny best case: {p0_bunny_final_summary['max_penetration_depth_bunny_mesh_m']*1000:.2f} mm max penetration, {p0_bunny_final_summary['body_speed_max']:.4f} m/s rigid speed",
            "Takeaway: continuous mesh contact is already working, but bunny geometry is a harder support target than a thick box",
        ],
    )
    _gif_twocol(
        prs,
        "H1C: Latest Low-Mass Cloth + Bunny OFF Runs",
        "Bunny mass = 0.1 kg",
        advisor_bm0p1_gif,
        "Bunny mass = 0.5 kg",
        advisor_bm0p5_gif,
        common_settings="self-collision OFF | `mass-spring-scale` = alpha | `shape_contact_scale` = alpha | `shape_contact_damping_multiplier` = 1 | `contact_dist_scale` = 5 | decoupled shape materials | total cloth mass = 0.1 kg | drop height = 0.5 m",
    )

    _pic_twocol(prs, "H1B Source Proof: Why Particle Radius Matters in OFF", "Detection: mesh closest-point query uses margin + radius", mesh_query_png, "Force: penetration depth subtracts particle radius", newton_contact_png)
    _body(prs, "H1B OFF Experiment Matrix", [
        "OFF only, object mass = 1 kg / particle",
        "`contact_collision_dist` scale = 0.5 / 1.0 / 2.0 / 3.0",
        "Drop height = 0.1 / 0.25 / 0.5 / 1.0 m",
        "Bunny mass = 1 / 10 / 1000 / 5000 kg",
        "64 runs: video + scene + summary + aggregate CSV",
    ])
    _gif_grid_2x2(prs, "H1B Hypothesis 1: Bunny Mass Matters", [
        ("Bunny mass = 1 kg", mass_light_gif),
        ("Bunny mass = 10 kg", mass_mid1_gif),
        ("Bunny mass = 1000 kg", mass_mid2_gif),
        ("Bunny mass = 5000 kg", mass_heavy_gif),
    ], common_settings="self-collision OFF | object mass = 1 kg / particle | `contact_collision_dist` scale = 1.0 | drop height = 0.1 m")
    _gif_grid_2x2(prs, "H1B Hypothesis 2: Height Matters, But Not Alone", [
        ("Drop height = 0.1 m", height_low_gif),
        ("Drop height = 0.25 m", height_mid1_gif),
        ("Drop height = 0.5 m", height_mid2_gif),
        ("Drop height = 1.0 m", height_high_gif),
    ], common_settings="self-collision OFF | object mass = 1 kg / particle | bunny mass = 5000 kg | `contact_collision_dist` scale = 1.0")
    _gif_grid_2x2(prs, "H1B Hypothesis 3: `contact_collision_dist` Is The Strongest Lever", [
        ("`contact_collision_dist` scale = 0.5", ccd_small_gif),
        ("`contact_collision_dist` scale = 1.0", ccd_mid1_gif),
        ("`contact_collision_dist` scale = 2.0", ccd_mid2_gif),
        ("`contact_collision_dist` scale = 3.0", ccd_large_gif),
    ], common_settings="self-collision OFF | object mass = 1 kg / particle | bunny mass = 5000 kg | drop height = 0.1 m")
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, "H1B Result: `contact_collision_dist` Is The Strongest Lever", size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)
    _add_pic(s, factor_chart, left, top, width, height)
    _add_label(s, int(left + 20), int(top + height - Inches(0.45)), int(width - 40), Inches(0.35), "Across 64 runs, larger `contact_collision_dist` (larger Newton `particle_radius`) helps most. Bunny mass helps. Height hurts. Height alone does not solve the issue.", font_size=12, bold=False)
    _gif_compare_large(
        prs,
        "H1B Best OFF Example",
        "Baseline\n`contact_collision_dist` = 1.0x",
        baseline_gif,
        "Best current\n`contact_collision_dist` = 3.0x",
        best_gif,
        footer="OFF | mass = 1 kg / particle | drop = 0.1 m | bunny = 5000 kg",
    )
    _boxed_bullets(prs, "H1B What This Week Actually Proved", "Validated / Falsified Hypotheses", [
        "Validated: self-collision cannot be directly ported",
        "Validated: bunny mass matters, but is not sufficient",
        "Validated: even unified mass+ke+kd scaling does not make light bunny masses stable",
        "Validated: larger `contact_collision_dist` (larger Newton `particle_radius`) strongly reduces penetration",
        "Falsified: lower drop height alone is enough",
        "Falsified: 10x / 100x stiffness-damping scaling is a robust fix",
    ])

    _section(prs, "H2: Why Newton Self-Collision ≠ PhysTwin")
    _body(prs, "H2 Terms: Gate, Update, Force", [
        "`approach-gated`: pair must still move toward each other before collision activates",
        "`impulse update`: source code writes velocity directly with `v_new = ...`",
        "`penalty force`: source code first accumulates force, then the solver integrates it",
    ])
    _picture_with_footer(prs, "H2 Newton Source: eval_particle_contact (Self-Collision)", newton_self_png, [
        "Proves: Newton has a native particle-particle self-contact kernel",
        "But it is a generic penalty model, not a cloth-topology-aware self-collision model",
    ])
    _pic_twocol(
        prs,
        "H2 Source Proof: PhysTwin vs Newton Self-Collision",
        "PhysTwin\n`dis_len < collision_dist`\n`wp.dot(dis, relative_v) < -1e-4`\n`v_new[tid] = v1 - J_average / m1`",
        phys_png,
        "Newton\n`err = d - radius - particle_radius[index]`\n`if err <= k_cohesion:`\n`f += particle_force(...)`",
        newton_self_png,
    )
    _boxed_bullets(
        prs,
        "H2 Result: Different Model, Different Claim",
        "What We Can Say Precisely",
        [
            "PhysTwin self-collision is approach-gated and impulse-style",
            "Newton native self-collision is overlap-triggered and penalty-force based",
            "So `self-collision ON` cannot be used as parity evidence",
            "We therefore keep `self-collision ON` only as a Newton-native ablation",
        ],
    )
    _gif_single(
        prs,
        "H2 ON Ablation Demo: Why We Cannot Simply Turn Self-Collision ON",
        h2_on_failure_gif,
        "Current ON case still develops unstable particle-cloud behavior, so enabling Newton native self-collision is not a valid replacement for PhysTwin self-collision.",
    )

    _section(prs, "H3: Multi-Deformable Interaction")
    _body(prs, "One Rope + Grounded Sloth: Design Decision", [
        "Native all-particle contact makes same-object neighbors fight springs",
        "We keep one Newton particle system, but filter to rope-vs-sloth pairs only",
        "Goal: preserve real rope-sloth interaction without internal blow-up",
    ])
    _picture_with_footer(prs, "Custom Kernel: _eval_cross_object_contact", cross_rope_png, [
        "Proves: we keep Newton's HashGrid and penalty model while filtering only cross-object pairs",
        "This makes multi-rope interaction stable and readable",
    ])
    _gif_twocol(
        prs,
        "H3 Demos: Rope + Rope and Rope + Sloth",
        "Two ropes on ground\ncross-rope particle contact only",
        rope_ground_gif,
        "One rope + grounded sloth\ncross-object particle contact only",
        multi_rope_gif,
        common_settings="No rigid bunny | Newton particle contact is filtered to cross-object pairs only",
    )
    _body(prs, "H3 Result: ✅ Validated", [
        "Both rope+rope and rope+sloth interactions are feasible without a rigid probe",
        "Selective pair filtering solved the same-object artifact",
        "This is a cleaner minimal test for future cloth-cloth interaction",
    ])

    _section(prs, "H4: MPM Two-Way Coupling")
    _gif_single(prs, "MPM: Kinematic Collider Works", mpm_sloth_gif, "One-way kinematic sand interaction works. Two-way support is the real failure case.")
    _body(prs, "Why MPM ≠ Two-Way with Spring-Mass", [
        "Mass ratio is extreme: deformable body is too heavy relative to sand response",
        "Ground-plane SDF competes against patch colliders and wins too often",
        "Patch meshes are too weak and too local for stable support forces",
    ])
    _body(prs, "H4 Result: ❌ Infeasible (by design)", [
        "Current MPM route is fine for kinematic colliders",
        "It is not a reliable two-way support model for spring-mass cloth or animals",
        "So the main research path stays on SemiImplicit contact, not MPM two-way",
    ])
    _body(prs, "Next Step", [
        "Narrow sweep: `contact_collision_dist` scale = 2.0 / 2.5 / 3.0 / 3.5",
        "Keep drop height = 0.1 m and bunny mass = 5000 kg",
        "Then test moderate ke / kd only inside that stable region",
    ])

    _body(prs, "Summary", [
        "✅ H1 deformable-rigid interaction is validated",
        "✅ H1B localizes bunny penetration to external `contact_collision_dist` / Newton `particle_radius`",
        "✅ H1C provides the latest decoupled alpha-scaled light-bunny reference runs",
        "✅ H2 shows Newton self-collision is not the PhysTwin self-collision model",
        "✅ H3 multi-deformable interaction works with selective filtering",
        "❌ H4 MPM two-way is not a viable path for this project",
    ])


TRANSCRIPT_MEETING = textwrap.dedent(
    f"""\
    # Meeting Transcript — PhysTwin → Newton Bridge

    语言：中文主讲 + English terminology  
    形式：short-form meeting script  
    目标：每页只讲 hypothesis、evidence、以及这周真正的 progress

    ---

    ## Slide 1 — Xinjie Zhang
    这一页不讲技术内容。  
    它只是 speaker page，后面马上进入 old baseline recall 和这周的新结论。

    ## Slide 2 — Recall 1: PhysTwin → Newton Bridge: Working
    这一页只做一个简单 recall。  
    这里放的是旧的 cloth baseline，意思不是再讲一遍细节，而是提醒老师：
    **在这次 meeting 之前，cloth 这条 bridge baseline 已经是 working 的。**

    ## Slide 3 — Recall 1: PhysTwin → Newton Bridge: Working
    这一页换成 zebra。  
    takeaway 和上一页一样：
    **zebra 这条 bridge baseline 以前也已经做成了。**

    ## Slide 4 — Recall 1: PhysTwin → Newton Bridge: Working
    这一页换成 sloth。  
    我要传达的信息还是一样：
    **sloth 这条 bridge baseline 也已经做成了。**

    ## Slide 5 — Recall 1: PhysTwin → Newton Bridge: Working
    这一页换成 rope。  
    还是同一句话：
    **rope 这条 bridge baseline 也已经做成了。**

    ## Slide 6 — Recall 1: PhysTwin → Newton Bridge: Working
    这一页不是视频，而是最短流程图。  
    它只想提醒老师一个结构性结论：
    PhysTwin 先输出 object state，bridge 写成中间表示层，Newton 再重建并继续做 interaction。  
    所以前 2 到 5 页那些 old baseline，不是偶然跑通，而是走同一条 bridge pipeline 跑通的。

    ## Slide 7 — Five Hypotheses, One Story
    五个 hypotheses 很简单。  
    Already done：bridge 已经把多个 PhysTwin object 重建进 Newton。  
    New this week：bunny penetration 现在被收敛到 external support quality，而不是 missing bridge step。  
    Today's claims：ON 不是 PhysTwin parity，filtered multi-deformable interaction works，MPM two-way 不是主线。

    ## Slide 8 — H1: Particle-Mesh Contact Is Already Bidirectional
    第 8 页先把 two-way deformable-rigid interaction 讲清楚。  
    左边 `create_soft_contacts` 说明 Newton 确实在对 bunny mesh 做 closest-point / signed-normal query。  
    右边 `eval_particle_body_contact` 说明接触力不是单向投影，而是 elastic、damping、friction 组成的 bidirectional penalty force。  
    所以 rope+bunny 这里不是看起来碰到了，而是真的有双向力交换。

    ## Slide 9 — H1 Method Note: Why We Do Not Lower Rope Mass Alone
    这里先解释一个方法问题：为什么我们不能只把 rope mass 调小。

    这里我只用 Newton core 证明，不再引用 bridge 层。

    第一段源码在 [kernels_particle.py]({NEWTON_KERNELS_PARTICLE}#L27) 到 [kernels_particle.py]({NEWTON_KERNELS_PARTICLE}#L50)。  
    Newton 的 spring force 直接写成：
    - `fs = dir * (ke * c + kd * dcdt)`

    也就是说，弹簧力这一层只看：
    - spring stiffness `ke`
    - spring damping `kd`
    - shape error `c`
    - relative velocity term `dcdt`

    这里没有任何“如果 mass 变小，就自动把 spring force 一起缩小”的补偿项。

    第二段源码在 [solver.py]({NEWTON_SOLVER_BASE}#L34) 到 [solver.py]({NEWTON_SOLVER_BASE}#L44)。  
    Newton 的粒子积分直接写成：
    - `v1 = v0 + (f0 * inv_mass + g) * dt`

    所以同一条 spring force，最后进入速度更新时还要再乘一次 `inv_mass = 1 / m`。  
    这就意味着：
    - 如果只把 mass 调小，而 `ke`、`kd` 不变，
    - 那么 spring 产生的 acceleration 会直接变大，
    - 有效的 `ke / m` 和 `kd / m` 都会一起变化。

    物理上，这不是“同一根 rope 只是更轻了”这么简单。  
    它同时改变了系统的 characteristic time scale。  
    更直白地说：
    - stiffness per unit mass 变大，
    - damping per unit mass 也变大，
    - rope 的响应速度、振动频率、以及接触时的回弹方式都会变。

    所以如果我们想比较“同一套 rope physics 在不同 rigid bunny 质量下会怎样”，  
    那就不能只单独改 deformable object 的 mass。  
    否则你已经不是在测“换 bunny mass”的 effect，  
    而是在同时改一整套 rope dynamics。

    ## Slide 10 — H1: Short IR Pipeline Reconstructs The Right Newton Objects
    第 10 页开始把 “rebuild” 这件事钉在源码上。  
    左边这段 importer 代码直接调用：
    - `builder.add_particles(...)`
    - `builder.add_spring(...)`

    这说明 IR 里的 object state 不是被当成渲染轨迹读进去，  
    而是被真正实例化成 Newton 的 native particles 和 native springs。

    右边这段 importer 代码继续做两件事：
    - 先 `builder.add_body(...)`
    - 再 `builder.add_shape_mesh(body=body, ...)`
    - 最后 `builder.finalize(device=device)`

    这三步的物理含义很明确：  
    bunny 不是贴图，不是离线 mesh animation，  
    而是一个真正的 Newton rigid body，上面再挂一个真正参与接触查询的 mesh shape。  
    `finalize()` 之后，整个系统才进入 Newton solver 的原生状态空间。

    所以这页真正想证明的是：  
    **bridge 在做 native object reconstruction，不是在 replay 旧轨迹。**

    ## Slide 11 — H1: Mesh Defines Contact; Rigid State Moves The Bunny
    第 11 页把 body / mesh / wrench 的物理角色彻底拆开。

    左边是 geometry kernel。  
    当 `d < margin + radius` 时，Newton 不只是说“碰到了”，  
    它会把：
    - `body_pos`
    - `body_vel`
    - `world_normal`
    - `soft_contact_particle`

    这些量写进 soft-contact buffer。  
    这一步的物理含义是：
    **mesh 决定 contact point 和 contact normal 在哪里。**

    右边是 contact kernel。  
    先算出：
    - friction `ft`
    - total contact `f_total`

    然后：
    - `wp.atomic_sub(particle_f, ..., f_total)`
    - 对 rigid body 的 `body_f` 累加 spatial force / torque

    这一步的物理含义是：
    **同一份 contact 不只推 cloth particle，也会写进 rigid bunny 的 wrench accumulator。**

    所以最干净的 Newton 解释不是泛泛地说 “cloth touches bunny”。  
    更准确的说法是：
    - mesh 定义 contact geometry
    - rigid state 接收 reaction wrench
    - solver 再根据这个 wrench 去更新 bunny motion

    ## Slide 12 — H1 Result: Deformable-Rigid Interaction Is Validated
    视频把 H1 收口。  
    现在 PhysTwin deformable object 放进 Newton 后，已经能和 novel rigid bunny 做真实力交换。  
    所以 H1 结论是成立的。

    ## Slide 13 — H1 Addendum: Rope+Bunny Still Works At 1 kg And 5 kg Total Rope Mass
    这一页补一个同场景的 rope-weight 对照。  
    我们固定：
    - bunny mass = `5 kg`
    - drop height = `1.0 m`
    - same close manual camera

    左边把 rope total mass 设成 `1 kg`，右边设成 `5 kg`。  
    这两组都通过同一个 `auto-set-weight` 自动换算出统一的 `weight_scale`，所以：
    - mass
    - spring
    - contact
    都按同一个比例联动。

    视频结论很直接：
    - 两边都还能稳定做 rope-bunny interaction
    - 更重的 rope 会让 bunny 的响应更强
    - 但 interaction 本身并没有因为换 rope total mass 而消失

    所以这页的作用不是讲 penetration，  
    而是把 H1 讲得更完整：  
    **在固定 bunny mass = 5 kg 时，rope total mass 从 1 kg 到 5 kg，deformable-rigid interaction 仍然成立。**

    ## Slide 14 — H1 Addendum: Cloth+Bunny OFF Still Works At 1 kg And 5 kg Total Cloth Mass
    这一页继续补 H1，但换成 cloth+bunny，而且明确是：
    - self-collision OFF
    - bunny mass = `5 kg`
    - drop height = `0.5 m`

    左边把 cloth total mass 设成 `1 kg`，右边设成 `5 kg`。  
    这两组同样通过 `auto-set-weight` 自动换算统一的 `weight_scale`，所以：
    - mass
    - spring
    - contact
    仍然按同一个比例联动。

    这里我只强调 H1 层面的结论：  
    即使换成 cloth，而且保持 OFF self-collision，deformable-rigid interaction 仍然存在。  
    更重的 cloth 会让 bunny 的响应更强，但 interaction 本身没有消失。

    这页不负责回答 penetration 为什么发生。  
    penetration 的归因还是放到后面的 H2 去讲。  
    这页只补一个更完整的 H1 statement：  
    **在固定 bunny mass = 5 kg 时，cloth total mass 从 1 kg 到 5 kg，OFF self-collision 的 cloth-bunny interaction 仍然成立。**

    ## Slide 15 — H1 Addendum: Cloth+Bunny OFF Still Works At 0.1 kg And 0.5 kg Total Cloth Mass
    这一页继续沿着同一条逻辑走，但把 cloth total mass 再往下压到：
    - `0.1 kg`
    - `0.5 kg`

    这里设置保持不变：
    - self-collision OFF
    - bunny mass = `5 kg`
    - drop height = `0.5 m`
    - 只保留第一次 cloth-bunny 接触后的大约 `1 second` 视频

    两组仍然通过同一个 `auto-set-weight` 自动换算统一的 `weight_scale`，  
    所以 mass、spring、contact 还是按同一个比例联动。

    这页的结论也只停留在 H1 层面：  
    即使 cloth total mass 压到 `0.1 kg` 和 `0.5 kg`，OFF self-collision 的 cloth-bunny interaction 仍然存在。  
    这里我们只是在补 interaction evidence，不是在提前下 penetration 的最终结论。

    ## Slide 16 — H2 Counterexample: 1x/2x Radius Still Penetrates Bunny
    从这一页开始，我正式转到 H2 的 counterexample。  
    这里固定：
    - self-collision OFF
    - cloth total mass = `0.1 kg`
    - bunny mass = `5 kg`
    - drop height = `0.5 m`
    - 视频在第一次 cloth-bunny 接触后 `1 second` 结束

    这一页只比较：
    - `1x radius`
    - `2x radius`

    我在标题括号里直接写绝对数值：
    - `1x` 对应 `particle_radius = 0.0051 m`
    - `2x` 对应 `particle_radius = 0.0102 m`

    这页要传达的是：  
    **即使把 radius 从 baseline 提到 2x，penetration 仍然存在。**

    ## Slide 17 — H2 Counterexample: 3x/4x Radius Still Penetrates Bunny
    这一页继续把 radius 往上推：
    - `3x radius = 0.0152 m`
    - `4x radius = 0.0203 m`

    目的不是证明 radius 没用，而是把结论讲准：  
    更大的 radius 会更早建立支撑，但 3x 和 4x 仍然没有把 bunny penetration 彻底消掉。

    ## Slide 18 — H2 Counterexample: 8x/10x Radius Still Penetrates Thin Ears
    这一页把 radius 再拉到更高：
    - `8x radius = 0.0406 m`
    - `10x radius = 0.0508 m`

    即使到这个量级，thin-ear geometry 仍然会留下 penetration。  
    所以更严谨的说法应该是：
    - larger radius helps
    - but radius alone is still not a complete fix for thin-ear support

    这三页合起来服务同一个 H2 conclusion：  
    **radius 是重要旋钮，但单靠 radius 并不能完全解决 bunny 几何上的支撑问题。**

    ## Slide 19 — H2 Counterexample: Bunny Is Harder Than Thick Box At Same Radius
    这一页把几何因素单独拉出来看。  
    我们保持：
    - self-collision OFF
    - cloth total mass = `0.1 kg`
    - rigid mass = `5 kg`
    - drop height = `0.5 m`
    - 视频在第一次 cloth-rigid 接触后 `1 second` 结束

    左边是 cloth + bunny，右边是 cloth + thick box。  
    这里 radius 完全一样，都是：
    - `particle_radius = 0.0051 m`

    所以这页不是在比较 radius，  
    而是在比较相同 radius 下，rigid geometry 本身会不会改变支撑质量。

    这页的 takeaway 很直接：  
    **同样的 cloth、同样的质量、同样的 radius，thick box 仍然比 bunny 更容易提供稳定支撑。**  
    所以后面的 penetration 问题，不能只归因到 radius，不然 box 和 bunny 不该有这么明显的差别。

    ## Slide 20 — H2 Counterexample: 1x/1.5x Bunny Size Still Penetrates
    这里我继续只改一个量：bunny size。  
    其它设置保持不变：
    - self-collision OFF
    - cloth total mass = `0.1 kg`
    - bunny mass = `5 kg`
    - `particle_radius = 0.0051 m`
    - drop height = `0.5 m`
    - 视频在第一次 cloth-bunny 接触后 `1 second` 结束

    这一页对比：
    - `1x bunny size (height = 0.194 m)`
    - `1.5x bunny size (height = 0.291 m)`

    这页要强调的是：  
    **即使把 bunny 几何整体放大到 1.5x，penetration 仍然存在。**  
    所以问题不能简单说成 “因为 bunny 太小”。

    ## Slide 21 — H2 Counterexample: 2x/3x Bunny Size Still Penetrates
    这一页继续把 bunny 再放大：
    - `2x bunny size (height = 0.389 m)`
    - `3x bunny size (height = 0.583 m)`

    更大的 bunny 的确会提供更宽、更高的支撑 patch，  
    但 penetration 仍然没有自然消失。

    所以这两页合起来服务同一个 H2 conclusion：  
    **更大的 rigid geometry 会改变支撑条件，但单纯放大 bunny size 也不是 penetration 的完整解。**

    ## Slide 22 — H2 Counterexample: Smaller `dt` Still Does Not Eliminate Penetration
    这里我单独测试一个数值问题：
    如果只把时间步继续缩小，会不会让 penetration 自然消失？

    我们保持：
    - self-collision OFF
    - cloth total mass = `0.1 kg`
    - bunny mass = `5 kg`
    - `particle_radius = 0.0051 m`
    - drop height = `0.5 m`
    - 视频在第一次 cloth-bunny 接触后 `1 second` 结束

    这一页对比的是：
    - case 1: `sim_dt = 1e-5`, `substeps = 3335`
    - case 2: `sim_dt = 1e-6`, `substeps = 33350`

    这页要讲清楚的是：  
    **更小的 `dt` 最多是 numerical sanity check，不会自动把 external support 问题变成 solved。**  
    如果 penetration 还在，就说明主问题仍然不是“步长太大”。

    ## Slide 23 — H2 Counterexample: Even 0.01 kg And 0.005 kg Still Penetrate Bunny
    这里我继续只改 cloth total mass，而且把它压到更小：
    - `0.01 kg`
    - `0.005 kg`

    其它设置保持：
    - self-collision OFF
    - bunny mass = `5 kg`
    - `particle_radius = 0.0051 m`
    - drop height = `0.5 m`
    - 视频在第一次 cloth-bunny 接触后 `1 second` 结束

    这页要说明的是：  
    **即使 cloth 总质量继续降到 `0.01 kg` 和 `0.005 kg`，penetration 仍然存在。**  
    所以更小的 cloth weight 也不是自动的完整解。

    ## Slide 13 — H2: Bunny Penetration Is Mainly A Support Problem
    H2 只问一个问题。  
    mesh contact 明明已经存在，为什么 cloth 还是会穿 bunny。  
    这周的判断是：主问题不是 missing contact，而是 external support quality 和 local geometry。

    ## Slide 15 — H2 Evidence: Continuous Mesh Contact Still Depends On `particle_radius`
    这里不能只说结论，要把源码和几何关系说清楚。

    第一段在 [kernels.py]({NEWTON_GEOM_KERNELS}#L1128) 到 [kernels.py]({NEWTON_GEOM_KERNELS}#L1142)。  
    只有当 `d < margin + radius` 时，Newton 才会把这对 `(particle, shape)` 写进 soft-contact buffer。  
    这说明即使 bunny 是 continuous mesh，  
    真正参与碰撞判定的几何对象也不是“点对曲面”，  
    而是“**带半径的粒子球** 对曲面”。

    第二段在 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L215) 到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L223)。  
    penalty depth 直接写成：
    - `c = dot(n, px - bx) - radius`
    - `fn = n * c * ke`

    这一步的物理含义是：
    - 同一个 particle center 位置下，`radius` 越大，
    - 有效 penetration `c` 越早进入作用区，
    - 法向恢复力 `fn` 也越早变大。

    所以 `particle_radius` 不是渲染参数。  
    它同时控制：
    - 何时开始认定接触存在
    - 以及 penalty support 何时开始真正把粒子往外推

    ## Slide 17 — H2: 64-Run OFF Grid Separates Three Main Effects
    这一页讲实验设计。  
    OFF case 里我们只扫三个轴：`contact_collision_dist`、`drop height`、`bunny mass`。  
    这样做的目的是把 penetration 的主导因素单独抽出来，而不是继续靠单个视频猜。

    ## Slide 18 — H2: `contact_collision_dist` Dominates Across 64 OFF Runs
    aggregate chart 给出最强的 quantitative conclusion。  
    `contact_collision_dist`，也就是 bridge 映射到 Newton 的 `particle_radius`，比 bunny mass 和 drop height 更能决定 penetration。  
    所以现在最该精修的旋钮是 `particle_radius`，不是继续大范围扫别的量。

    ## Slide 16 — H2: Bunny Now Has A Near-Working Point
    这页看 progress。  
    左边是 baseline，右边是 current best bunny case。  
    结论是问题已经从广泛失败推进到了 near-working point，而且最新 low-mass reruns 也保住了 rigid-rigid support。

    ## Slide 18 — H2 Result: Validated And Falsified Hypotheses
    到这里 H2 可以收成四句。  
    更大的 `contact_collision_dist / particle_radius` 明显减轻 penetration。  
    bunny mass 有帮助，但不是主导项。  
    thin-ear geometry 说明 radius 不是完整解。  
    只降 height 或暴力加 stiffness / damping 都不是稳健解。

    ## Slide 19 — H3: Newton Self-Collision Is Not PhysTwin Parity
    H3 不再处理 OFF penetration。  
    它只回答一个问题：为什么 Newton native self-collision 不能被讲成 PhysTwin self-collision parity。  
    这里我会顺带提一句：Newton 官方 examples 确实有 self-contact demo，但主线是 VBD，不是 SemiImplicit。

    ## Slide 20 — H3: PhysTwin And Newton Use Different Self-Collision Models
    这页要把模型差异讲成源码差异，而不是口号差异。

    PhysTwin 那边，条件是：
    - `dis_len < collision_dist`
    - `wp.dot(dis, relative_v) < -1e-4`

    然后它不是累 force，而是直接改 velocity。  
    也就是 pair 必须还在相向靠近，才会触发这条更新。

    Newton 这边，在 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L79) 到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L94)，  
    条件写成：
    - `err = d - radius_i - radius_j`
    - `if err <= k_cohesion:`
    - `particle_f[i] += ...`

    这意味着 Newton native self-collision 的判据是 overlap / cohesion，  
    更新变量是 force accumulator，不是 velocity overwrite。

    物理上这两条线的差别是：
    - PhysTwin 更像 **approach-gated impulse correction**
    - Newton 更像 **overlap-triggered penalty contact**

    所以这里不是“参数还没调到位”，  
    而是底层 contact model 已经不是同一个东西。

    ## Slide 21 — H3 Result: `self-collision ON` Is Only A Newton-Native Ablation
    所以 H3 的结论很直接。  
    `self-collision ON` 只能作为 Newton-native ablation，不能作为 parity evidence。  
    它可以展示现象，但不能承担 bridge correctness 的 claim。  
    官方 examples 里可以找到 self-contact，但那也不能直接当成 SemiImplicit parity 的证据。

    ## Slide 22 — H4: Pair-Filtered Multi-Deformable Interaction Is Viable
    H4 看的是 multi-deformable interaction。  
    这里的假设是：可以做，但前提是 pair filtering 先讲清楚。

    ## Slide 23 — H4: We Changed Pair Selection, Not The Physics Model
    这页源码的重点是：我们改的是 admissible pair set，不是 contact law。

    在 custom kernel 里，hash-grid neighbor query 还是照常做。  
    但真正决定 pair 能不能进 contact 的，是对象归属过滤：
    - same-object pair 直接跳过
    - cross-object pair 才继续算 penalty contact

    然后真正的力学形式没有换。  
    还是同一个 penalty contact：
    - 法向排斥
    - 阻尼
    - 摩擦

    也就是说，这里不是“我们重新发明了一套 rope-sloth physics”。  
    我们只是把错误的 same-object 排斥去掉，  
    让 Newton 原来的 contact law 只作用在真正应该互相碰撞的 cross-object pairs 上。

    物理原则可以概括成一句话：  
    **interaction law 不变，pair topology 变了。**  
    这就是为什么它能修掉 artifact，同时还保持 Newton penalty model 的解释一致性。

    ## Slide 24 — H4 Result: Filtered Cross-Object Contact Already Works In Two Demos
    视频给出两个最小正结果：`rope + rope` 和 `rope + sloth`。  
    两条线都已经工作，所以 H4 现在是 validated result，不是概念假设。

    ## Slide 25 — H4 Addendum: Two-Rope Weight Contrast Still Works
    这一页在 H4 里再补一个更直接的质量对照。  
    场景保持不变：
    - lower rope rests on ground
    - upper rope drops in a cross layout
    - 仍然只开 cross-rope contact

    这里比较两组总质量：
    - case 1: lower rope = `0.5 kg`, upper rope = `0.1 kg`
    - case 2: lower rope = `5 kg`, upper rope = `5 kg`

    这页要表达的不是 “质量完全不重要”，  
    而是：**即使把两根 rope 的总质量设成不一样，pair-filtered cross-rope interaction 仍然存在，而且视频上仍然可读。**

    ## Slide 26 — H4 Addendum: Two Crossed Ropes Also Work On A Rigid Box
    这一页再补一个 shared rigid support 的例子。  
    不再是 ground-only，而是：
    - 两根 crossed ropes 一起落到同一个 rigid box 上
    - rope-rope contact 仍然只保留 cross-object pairs

    这个例子的作用不是取代 ground-only minimal demo，  
    而是说明：
    **cross-rope interaction 不只在 ground 上能看见，在 shared rigid support 上也仍然可读。**

    ## Slide 27 — H4 Addendum: Native Franka Can Lift And Release The Bridge Rope
    这一页把 H4 再往 downstream interaction 推一步。  
    这里不再是 proxy sphere / capsule，  
    而是直接复用 Newton 自带的 `Franka Panda` robotics asset。

    这一页真正想表达的是：  
    **同一根 bridge rope 不只会和 rope / sloth / rigid box 互动，也已经能被 native Franka 做出可读的 `approach -> lift -> release`。**  
    这里我也会主动把 claim 说窄一点：  
    这是 `taskful native baseline`，不是 full force-feedback robot policy。

    ## Slide 25 — H5: MPM One-Way Is Useful; Two-Way Is Not The Main Path
    H5 是路线取舍。  
    我们不是说 MPM 没用，而是问它适不适合当前 spring-mass two-way 目标。  
    当前判断是：one-way 有价值，但 two-way 不该继续做主线。

    ## Slide 26 — H5: One-Way MPM Works; Two-Way Should Be Deprioritized
    这一页先承认 positive result。  
    one-way kinematic collider 是工作的。  
    但 two-way support 更像 architecture mismatch，不是小参数问题，所以不应该继续当主线。

    ## Slide 27 — Summary: Which Hypotheses Survived
    最后的总括很清楚。  
    H1 证明 bridge 已经工作。  
    H2 证明 penetration 主因是 external support quality，而且 even 10x radius 也不能单独消掉 thin-ear penetration。  
    H3 说明 self-collision 不是 parity target。  
    H4 给出 pair-filtered multi-deformable 的正结果，而且补出了 native Franka 对同一根 bridge rope 的 `lift and release` baseline。  
    H5 说明 MPM two-way 不是主线。

    ## Slide 28 — Next Step: Fewer Slides, Stronger Evidence
    下一步不再继续加页数。  
    每个 hypothesis 只保留一页 flow、一页 code proof、一页 best demo。  
    同时继续精扫 `particle_radius` 附近的稳定区间，把 H2 从 near-working point 推到更硬的结论。

    ---

    # Meeting Summary

    一分钟 summary 可以直接这样讲：

    1. H1：bridge 已经保住 deformable-rigid interaction。  
    2. H2：bunny penetration 的主因是 external support quality，而不是 missing mesh contact。  
    3. H3：Newton native self-collision 不是 PhysTwin parity。  
    4. H4：pair-filtered multi-deformable interaction 已经有正结果，而且 native Franka 已经能对同一根 bridge rope 做出可读的 `lift and release`。  
    5. H5：MPM one-way 有用，但 two-way 不是这条项目主线。
    """
)

TRANSCRIPT_MEETING = textwrap.dedent(
    """\
    # Meeting Transcript — PhysTwin → Newton Bridge

    语言：中文主讲 + English terminology  
    形式：按照当前手工调整后的 37 页 deck 逐页口播  
    目标：让 transcript 和 `My Adjust.pptx` 的页序、标题、叙事完全对齐

    ---

    ## Slide 1 — Xinjie Zhang
    这一页只是 opening speaker page。  
    我会直接从 recall 开始，不在这里展开技术内容。

    ## Slide 2 — Recall 1: PhysTwin → Newton Bridge: Working
    第一页 recall 先把旧 cloth PASS case 拉回来。  
    它的作用只是提醒老师：bridge baseline 不是今天新讲的东西。

    ## Slide 3 — Recall 1: PhysTwin → Newton Bridge: Working
    第二页 recall 对应 zebra。  
    作用同样是告诉老师：旧 PASS case 已经建立，我们今天不重复证明旧结论。

    ## Slide 4 — Recall 1: PhysTwin → Newton Bridge: Working
    第三页 recall 对应 sloth。  
    这里仍然只是视觉回顾，不承担新的 hypothesis 证明。

    ## Slide 5 — Recall 1: PhysTwin → Newton Bridge: Working
    第四页 recall 对应 rope。  
    它提醒老师：rope 这一条 bridge baseline 已经成立。

    ## Slide 6 — Recall 1: PhysTwin → Newton Bridge: Working
    最后一页 recall 回到 bridge flow 本身。  
    它的作用是把旧 deck 的主线重新唤起：PhysTwin outputs -> bridge IR -> Newton reconstruction。

## Slide 7 — Recall 2: Bunny+Rope
这一页 recall 的重点不再是旧 PASS case，而是 bunny+rope。  
它提醒老师：novel rigid bunny interaction 这条线早就已经开始工作，只是当时还没有展开 cloth 归因。

## Slide 8 — Verification: Particle-Mesh Contact Is Already Bidirectional
recall 到这里就收住。  
下面开始正式做 verification。  
这页不能只说 “particle-mesh contact 存在”，而是要把源码里的检测和受力分开讲。

左边源码在 [kernels.py]({NEWTON_GEOM_KERNELS}#L1104) 到 [kernels.py]({NEWTON_GEOM_KERNELS}#L1133)。  
Newton 先把 particle 位置变到 shape local frame，再做 mesh query。  
只有当：
- `d < margin + radius`

时，才会把这对 `(particle, shape)` 写进 soft-contact buffer。  
这一步的物理含义是：
- 参与判定的不是“点碰曲面”
- 而是“带半径的粒子球”去碰 bunny mesh

右边源码在 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L203) 到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L213)，以及 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L233) 到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L236)。  
接触深度直接写成：
- `c = dot(n, px - bx) - radius`

然后同一份 contact 会分成：
- 法向弹性 `fn`
- 阻尼 `fd`
- 摩擦 `ft`

最后：
- 对 particle 做 `particle_f -= total_contact`
- 对 rigid bunny 做 `body_f += total_contact` 和 `body_t += cross(...)`

这一步的物理含义是：  
**同一份 penalty contact 同时更新 deformable 和 rigid 两边。**

所以这页真正要证明的是：  
**rope-bunny interaction 不是看起来碰到了，而是 Newton 源码里已经存在真实的双向力交换。**

    ## Slide 9 — Verification: Method Note: Why We Do Not Lower Rope Mass Alone
    这一页也不能只讲方法论口号，而要直接从 Newton core 解释为什么。

    第一段源码在 [kernels_particle.py]({NEWTON_KERNELS_PARTICLE}#L27) 到 [kernels_particle.py]({NEWTON_KERNELS_PARTICLE}#L50)。  
    弹簧力直接写成：
    - `fs = dir * (ke * c + kd * dcdt)`

    也就是说，force law 这一层只看：
    - `ke`
    - `kd`
    - shape error `c`
    - relative velocity term `dcdt`

    这里没有“如果 mass 变小，就自动把 spring force 一起缩小”的补偿。

    第二段源码在 [solver.py]({NEWTON_SOLVER_BASE}#L34) 到 [solver.py]({NEWTON_SOLVER_BASE}#L44)。  
    Newton 的粒子积分直接写成：
    - `v1 = v0 + (f0 * inv_mass + g) * dt`

    所以同一条 spring force 真正进入动力学时，还要再乘一次 `inv_mass = 1 / m`。

    物理上，这意味着：
    - 如果只单独把 rope mass 调小
    - 但 `ke`、`kd` 不变
    - 那么系统的有效 `ke / m` 和 `kd / m` 都会一起变

    结果就不是“同一根 rope 只是更轻了”，  
    而是 characteristic time scale、振动频率和接触响应都被改掉了。

    所以这页真正要证明的是：  
    **只改 deformable mass 不是一个干净的 weight experiment；如果要比重量，mass、spring、contact 必须联动缩放。**

    ## Slide 10 — Verification: Short IR Pipeline Reconstructs The Right Newton Objects
    这页也按“源码 + 物理原则”来讲，不用抽象词。

    左边这段 importer 源码在 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L1651) 到 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L1658)，以及 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L1674) 到 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L1678)。  
    bridge 直接调用：
    - `builder.add_particles(...)`
    - `builder.add_spring(...)`

    这说明 IR 里的 object arrays 不是被当成渲染轨迹读进去，  
    而是被实例化成 Newton 的 native particles 和 native springs。

    右边这段 importer 源码在 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L1691) 到 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L1693)，[newton_import_ir.py]({NEWTON_IMPORT_IR}#L1705) 到 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L1706)，[newton_import_ir.py]({NEWTON_IMPORT_IR}#L1732) 到 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L1738)，以及 [newton_import_ir.py]({NEWTON_IMPORT_IR}#L1762)。  
    它依次做：
    - `builder.add_body(...)`
    - `builder.add_shape_mesh(body=body, ...)`
    - `builder.finalize(device=device)`

    这三步的物理含义是：
    - bunny 的 mass / inertia / motion state 进入 rigid body
    - bunny 的 contact geometry 进入 attached mesh shape
    - `finalize()` 之后整个系统才真正进入 Newton solver 的原生状态空间

    所以这页真正要证明的是：  
    **bridge 在做 native object reconstruction，不是在 replay 旧轨迹。**

    ## Slide 11 — Verification: Mesh Defines Contact; Rigid State Moves The Bunny
    第 11 页把 geometry 和 dynamics 在源码里真正拆开。

    左边是 geometry kernel，在 [kernels.py]({NEWTON_GEOM_KERNELS}#L1128) 到 [kernels.py]({NEWTON_GEOM_KERNELS}#L1142)。  
    一旦 `d < margin + radius`，Newton 不只是说“碰到了”，  
    它会把：
    - `body_pos`
    - `body_vel`
    - `world_normal`
    - `soft_contact_particle`

    这些量一起写进 contact buffer。  
    这一步的物理含义是：
    **mesh 决定 contact point / contact normal 在哪里，以及这份 contact 属于哪个 rigid body。**

    右边是 contact kernel，在 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L251) 到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L261)。  
    先算出总接触力 `total_contact`，然后：
    - `particle_f -= total_contact`
    - `body_f += total_contact`
    - `body_t += cross(r, total_contact)`

    这一步的物理含义是：
    **同一份 contact 最终会变成 rigid bunny 的平动力和转动力矩。**

    所以最准确的解释不是笼统地说 “cloth touches bunny”，  
    而是：
    - mesh 定义 contact geometry
    - rigid state 接收 reaction wrench
    - solver 再用这个 wrench 更新 bunny motion

    ## Slide 12 — H1: Can We Adjust the Weight of the Rope, Not Just the Bunny?
    H1 在这里被明确写成一个问题。  
    我们不是只改 bunny mass，也要看 deformable 自己的 total weight 能不能被统一控制，而且 interaction 还能不能保住。

    ## Slide 13 — H1 Method Note: Why We Do Not Lower Rope Mass Alone
    这一页虽然标题回到 H1，但源码标准和第 9 页一样，不能只说结论。

    还是同一组 Newton core 源码：
    - [kernels_particle.py]({NEWTON_KERNELS_PARTICLE}#L27) 到 [kernels_particle.py]({NEWTON_KERNELS_PARTICLE}#L50)：spring force 直接由 `ke`、`kd` 决定
    - [solver.py]({NEWTON_SOLVER_BASE}#L34) 到 [solver.py]({NEWTON_SOLVER_BASE}#L44)：速度更新再乘一次 `inv_mass`

    所以物理原则没有变：
    - 单独改 mass
    - 会连带改掉有效的 `ke / m`、`kd / m`
    - 也会一起改掉 rope 的动力学时间尺度

    这页在 H1 里的作用，就是把实验设计边界讲清楚：  
    **如果我们要比较重量效应，就必须用统一的 weight scale，让 mass、spring、contact 一起跟着变。**

    ## Slide 14 — H1: Rope+Bunny Still Works At 1 kg And 5 kg Total Rope Mass
    这里开始给 rope+bunny 的 weight 对照。  
    左边是 `1 kg` rope，右边是 `5 kg` rope。  
    takeaway 是：interaction 仍然存在，更重的 rope 只是让 bunny 响应更明显。

    ## Slide 15 — H1: Comparison
    这一页把 H1 的 rope 线和后面的 cloth 线连起来看。  
    它的作用是把“can interact”这件事和“why still penetrates”这件事分开。

## Slide 16 — H1: Cloth+Bunny OFF Still Works At 1 kg And 5 kg Total Cloth Mass
这里把 H1 扩展到 cloth+bunny，而且明确是 self-collision OFF。  
takeaway 仍然是 interaction 还在，不是说 weight 一变 interaction 就没了。

## Slide 17 — H2: Why Cloth Penetrates Bunny
到这里，H1 负责证明的 “can interact” 已经够了。  
下面开始进入 H2。  
H2 不再问 interaction 是否存在，而是问：mesh contact 明明已经存在，为什么 cloth 还是会穿 bunny。

    ## Slide 18 — H2 Counterexample: Smaller Weights Still Penetrates Bunny
    这一页先用 `1 kg` 和 `5 kg` cloth 总质量做第一层 counterexample。  
    目的是说明：仅仅改 cloth total mass，并不会自动让 penetration 消失。

    ## Slide 19 — H2 Counterexample: Smaller Weights Still Penetrates Bunny
    这里把质量继续降到 `0.1 kg` 和 `0.5 kg`。  
    即使更轻，penetration 仍然存在，所以“让 cloth 更轻”不是完整解。

    ## Slide 20 — H2 Counterexample: Even 0.01 kg And 0.005 kg Still Penetrate Bunny
    这一页把 cloth total mass 再压低一个量级，到 `0.01 kg` 和 `0.005 kg`。  
    结果还是 nonzero penetration，所以更小的 cloth weight 也不是自动解法。

    ## Slide 21 — H2 Evidence: Continuous Mesh Contact Still Depends On `particle_radius`
    这一页也按同一标准来讲：先讲源码，再讲物理原则。

    第一段在 [kernels.py]({NEWTON_GEOM_KERNELS}#L1128) 到 [kernels.py]({NEWTON_GEOM_KERNELS}#L1142)。  
    只有当：
    - `d < margin + radius`

    这对 `(particle, shape)` 才会进入 soft-contact buffer。  
    这说明 bunny 虽然是 continuous mesh，  
    但真正参与接触判定的几何对象并不是“点对曲面”，  
    而是“**带半径的粒子球** 对曲面”。

    第二段在 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L203) 到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L213)。  
    penalty depth 直接写成：
    - `c = dot(n, px - bx) - radius`
    - `fn = n * c * ke`

    物理含义就是：
    - 同一个 particle center 位置下
    - `radius` 越大
    - 有效 penetration `c` 越早进入作用区
    - 法向恢复力 `fn` 也越早变大

    所以 `particle_radius` 不是渲染参数。  
    它同时控制：
    - 何时开始认定接触存在
    - 以及 penalty support 何时开始真正把粒子往外推

    这页真正要证明的是：  
    **continuous mesh contact 仍然显式依赖 `particle_radius`，所以 particle spacing / radius 会直接影响 bunny penetration。**

    ## Slide 22 — H2 Counterexample: 1x/2x Radius Still Penetrates Bunny
    这页开始把 `particle_radius` 做成系统对照。  
    先看 `1x` 和 `2x`：radius 变大有帮助，但 penetration 没有消失。

    ## Slide 23 — H2 Counterexample: 3x/4x Radius Still Penetrates Bunny
    这里继续把 radius 往上推到 `3x` 和 `4x`。  
    这页服务的还是同一个结论：radius 是重要旋钮，但不是一开就 pass 的单旋钮。

    ## Slide 24 — H2 Counterexample: 8x/10x Radius Still Looks Weird
    到 `8x` 和 `10x` 时，视觉上已经开始变得奇怪。  
    所以这页除了说明 radius 不是完整解，还说明“暴力放大 radius”本身会带来新的 artifact。

    ## Slide 25 — H2 Counterexample: Bunny Is Harder Than Thick Box At Same Radius
    这一页把 geometry 单独拿出来看。  
    同样的 cloth、同样的 mass、同样的 radius，box 仍然比 bunny 更容易提供稳定支撑。

    ## Slide 26 — H2 Counterexample: 1x/1.5x Bunny Size Still Penetrates
    这里开始只改 bunny size。  
    即使把 bunny 放大到 `1.5x`，penetration 仍然存在，所以不能简单说“只是因为 bunny 太小”。

    ## Slide 27 — H2 Counterexample: 2x/3x Bunny Size Still Penetrates
    这一页继续把 bunny 放大到 `2x` 和 `3x`。  
    rigid geometry 的确变了，但 penetration 还是没有自动消失。

## Slide 28 — H2 Counterexample: Smaller `dt` Still Does Not Eliminate Penetration
这里单独测试 timestep。  
把 `dt` 再缩两个数量级以后，penetration 仍然在，所以它最多是 numerical sanity check，不是主解。

## Slide 29 — H3: Newton Self-Collision Is Not PhysTwin Parity
到这里，H2 的 external contact counterexample 已经够密了。  
H3 在这个版本里只负责给出方法论边界。  
我们不能把 Newton native self-collision 讲成 PhysTwin self-collision parity。

    ## Slide 30 — H3: PhysTwin And Newton Use Different Self-Collision Models
    这一页不能只说 “模型不同”，而要把差异落在源码里的 trigger 和 update 变量上。

    PhysTwin 这边，在 [spring_mass_warp.py]({PHYSTWIN_SIM}#L196) 到 [spring_mass_warp.py]({PHYSTWIN_SIM}#L206)，  
    触发条件写成：
    - `dis_len < collision_dist`
    - `wp.dot(dis, relative_v) < -1e-4`

    也就是说，pair 不仅要够近，还必须仍然在相向靠近。  
    然后在 [spring_mass_warp.py]({PHYSTWIN_SIM}#L294) 到 [spring_mass_warp.py]({PHYSTWIN_SIM}#L296)，  
    它不是累 force，而是直接写：
    - `v_new[tid] = ...`

    这条线的物理原则是：  
    **approach-gated impulse-style velocity correction。**

    Newton 这边，在 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L76) 到 [kernels_contact.py]({NEWTON_KERNELS_CONTACT}#L82)，  
    条件写成：
    - `err = d - radius_i - radius_j`
    - `if err <= k_cohesion:`

    一旦条件成立，它直接往 `particle_f` 里加 penalty contact force。  
    也就是说，Newton 的判据是 overlap / cohesion，  
    更新变量是 force accumulator，不是 velocity overwrite。

    这条线的物理原则是：  
    **overlap-triggered penalty-force accumulation。**

    所以这页真正要证明的是：  
    **这里不是 tuning difference，而是 contact model difference。**  
    也正因为如此，`self-collision ON` 不能被讲成 PhysTwin parity evidence。

## Slide 31 — H4: Multi-Deformable Interaction
H3 收口之后，下面转到 H4。  
这里的问题变成：多个 deformables 能不能彼此稳定互动，而不是只和 rigid 互动。

    ## Slide 32 — H4 Demos: Rope + Rope and Rope + Sloth
    这一页给出两个最小正结果：`rope + rope` 和 `rope + sloth`。  
    它们说明 pair-filtered multi-deformable interaction 已经是 working demo，不只是概念。

    ## Slide 33 — H4 Addendum: Two-Rope Weight Contrast Still Works
    这里再补一个 weight contrast。  
    即使 lower rope 和 upper rope 的 total mass 不一样，cross-rope interaction 仍然存在，而且视频上仍然可读。

    ## Slide 34 — H4 Addendum: Two Crossed Ropes Also Work On A Rigid Box
    这页把 shared rigid support 拉进 H4。  
    不是只有 ground-only 场景能看见 rope-rope interaction，换成 one rigid box 也还能看见。

    ## Slide 35 — H4 Addendum: Native Franka Can Lift And Release The Bridge Rope
    这页把 H4 再往 downstream interaction 推一步。  
    这里不再是 proxy sphere / capsule，而是直接复用 Newton 自带的 `Franka Panda` robotics asset。  
    它的作用是说明：同一根 bridge rope 不只会和 rope / sloth / rigid box 互动，也已经能被 native Franka 做出可读的 `approach -> lift -> release`。  
    这里我会主动把 claim 说窄一点：这是 `taskful native baseline`，不是 full force-feedback robot policy。

## Slide 36 — H5: MPM
H4 讲完 multi-deformable 之后，最后进入 H5。  
H5 不是说 MPM 没用，而是要把它从这条 spring-mass 主线里放回正确位置。

    ## Slide 37 — MPM: Kinematic Collider Works
    这页给 H5 的正结果。  
    one-way kinematic collider 是工作的；真正不该继续当主线的是 two-way spring-mass support。
    """
)


def build_deck_reference(
    prs: Presentation, rows: list[dict], gif_dir: Path, image_dir: Path
) -> None:
    _delete_all_slides(prs)

    factor_chart = _make_factor_chart(rows, image_dir / "factor_chart.png")
    phys_png = _code_excerpt_image(
        PHYSTWIN_SIM,
        "phystwin_self_collision_ref_v2",
        _extract_code_segments(
            PHYSTWIN_SIM,
            [(196, 200), (205, 206), (220, 220), (294, 296)],
            highlight_lines={196, 198, 199, 205, 220, 294, 295},
        ),
        image_dir,
    )
    create_soft_png = _code_excerpt_image(
        NEWTON_GEOM_KERNELS,
        "newton_create_soft_contacts_ref_v2",
        _extract_code_segments(
            NEWTON_GEOM_KERNELS,
            [(1013, 1014), (1112, 1115), (1132, 1134)],
            highlight_lines={1014, 1113, 1114, 1132, 1133},
        ),
        image_dir,
    )
    particle_body_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_eval_particle_body_contact_ref_v2",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(215, 223), (244, 248), (269, 273)],
            highlight_lines={216, 222, 223, 245, 248, 271, 273},
        ),
        image_dir,
    )
    newton_self_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_self_collision_ref_v2",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(79, 80), (83, 94)],
            highlight_lines={80, 88, 90, 94},
        ),
        image_dir,
    )
    mesh_query_png = _code_excerpt_image(
        NEWTON_GEOM_KERNELS,
        "newton_mesh_query_ref_v2",
        _extract_code_segments(
            NEWTON_GEOM_KERNELS,
            [(1112, 1115), (1132, 1132)],
            highlight_lines={1113, 1114, 1115, 1132},
        ),
        image_dir,
    )
    newton_contact_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_contact_force_ref_v2",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(215, 223), (244, 248), (263, 267)],
            highlight_lines={216, 222, 223, 245, 248, 263, 265},
        ),
        image_dir,
    )
    cross_rope_png = _code_excerpt_image(
        ROOT / "Newton" / "phystwin_bridge" / "demos" / "demo_rope_sloth_ground_contact.py",
        "cross_object_contact_ref_v2",
        _extract_code_segments(
            ROOT / "Newton" / "phystwin_bridge" / "demos" / "demo_rope_sloth_ground_contact.py",
            [(45, 53), (67, 68), (80, 81), (87, 93)],
            highlight_lines={45, 51, 52, 67, 80, 87, 91},
        ),
        image_dir,
    )

    rigid_h1_gif = _ensure_gif(CURRENT_ROPE_BUNNY_MP4, gif_dir, stem_override="rope_bunny_current_v2")
    baseline_gif = _ensure_gif(BASELINE_OFF_MP4, gif_dir, stem_override="baseline_off_ccd1_h0p1_bm5000_v2")
    best_gif = _ensure_gif(BEST_OFF_MP4, gif_dir, stem_override="best_off_ccd3_h0p1_bm5000_v2")
    mass_light_gif = _ensure_gif(MASS_LIGHT_MP4, gif_dir, stem_override="mass_bm1_ccd1_h0p1_v2")
    mass_mid1_gif = _ensure_gif(MASS_MID1_MP4, gif_dir, stem_override="mass_bm10_ccd1_h0p1_v2")
    mass_mid2_gif = _ensure_gif(MASS_MID2_MP4, gif_dir, stem_override="mass_bm1000_ccd1_h0p1_v2")
    mass_heavy_gif = _ensure_gif(MASS_HEAVY_MP4, gif_dir, stem_override="mass_bm5000_ccd1_h0p1_v2")
    height_low_gif = _ensure_gif(HEIGHT_LOW_MP4, gif_dir, stem_override="height_h0p1_ccd1_bm5000_v2")
    height_mid1_gif = _ensure_gif(HEIGHT_MID1_MP4, gif_dir, stem_override="height_h0p25_ccd1_bm5000_v2")
    height_mid2_gif = _ensure_gif(HEIGHT_MID2_MP4, gif_dir, stem_override="height_h0p5_ccd1_bm5000_v2")
    height_high_gif = _ensure_gif(HEIGHT_HIGH_MP4, gif_dir, stem_override="height_h1p0_ccd1_bm5000_v2")
    ccd_small_gif = _ensure_gif(CCD_SMALL_MP4, gif_dir, stem_override="ccd_s0p5_h0p1_bm5000_v2")
    ccd_mid1_gif = _ensure_gif(CCD_MID1_MP4, gif_dir, stem_override="ccd_s1p0_h0p1_bm5000_v2")
    ccd_mid2_gif = _ensure_gif(CCD_MID2_MP4, gif_dir, stem_override="ccd_s2p0_h0p1_bm5000_v2")
    ccd_large_gif = _ensure_gif(CCD_LARGE_MP4, gif_dir, stem_override="ccd_s3p0_h0p1_bm5000_v2")
    advisor_bm0p1_gif = _ensure_gif(ADVISOR_BM0P1_MP4, gif_dir, stem_override="advisor_bm0p1_v2")
    advisor_bm0p5_gif = _ensure_gif(ADVISOR_BM0P5_MP4, gif_dir, stem_override="advisor_bm0p5_v2")
    thin_ear_5x_gif = _ensure_gif(THIN_EAR_5X_MP4, gif_dir, stem_override="thin_ear_ccd5x_v2")
    thin_ear_10x_gif = _ensure_gif(THIN_EAR_10X_MP4, gif_dir, stem_override="thin_ear_ccd10x_v2")
    rope_ground_mp4 = ROOT / "tmp" / "rope_ground_drop_full_close_20260322" / "rope_ground_drop.mp4"
    rope_ground_gif = (
        _ensure_gif(rope_ground_mp4, gif_dir, stem_override="rope_ground_drop_v2")
        if rope_ground_mp4.exists()
        else old_assets_gif_dir / "rope_drag_on_overlay1x3.gif"
    )
    multi_rope_mp4 = ROOT / "tmp" / "rope_sloth_dualsource_drop1p0_r2_20260322" / "rope_sloth_ground_drop.mp4"
    multi_rope_gif = (
        _ensure_gif(multi_rope_mp4, gif_dir, stem_override="rope_sloth_ground_drop_v2")
        if multi_rope_mp4.exists()
        else old_assets_gif_dir / "sloth_cmp2x3.gif"
    )
    mpm_sloth_gif = _ensure_gif(MPM_SLOTH_MP4, gif_dir, stem_override="sloth_mpm_sand_one_way_wide_v2")
    baseline_gif = _ensure_gif(BASELINE_OFF_MP4, gif_dir, stem_override="baseline_off_ccd1_h0p1_bm5000_v2")
    best_gif = _ensure_gif(BEST_OFF_MP4, gif_dir, stem_override="best_off_ccd3_h0p1_bm5000_v2")
    mass_light_gif = _ensure_gif(MASS_LIGHT_MP4, gif_dir, stem_override="mass_bm1_ccd1_h0p1_v2")
    mass_mid1_gif = _ensure_gif(MASS_MID1_MP4, gif_dir, stem_override="mass_bm10_ccd1_h0p1_v2")
    mass_mid2_gif = _ensure_gif(MASS_MID2_MP4, gif_dir, stem_override="mass_bm1000_ccd1_h0p1_v2")
    mass_heavy_gif = _ensure_gif(MASS_HEAVY_MP4, gif_dir, stem_override="mass_bm5000_ccd1_h0p1_v2")
    height_low_gif = _ensure_gif(HEIGHT_LOW_MP4, gif_dir, stem_override="height_h0p1_ccd1_bm5000_v2")
    height_mid1_gif = _ensure_gif(HEIGHT_MID1_MP4, gif_dir, stem_override="height_h0p25_ccd1_bm5000_v2")
    height_mid2_gif = _ensure_gif(HEIGHT_MID2_MP4, gif_dir, stem_override="height_h0p5_ccd1_bm5000_v2")
    height_high_gif = _ensure_gif(HEIGHT_HIGH_MP4, gif_dir, stem_override="height_h1p0_ccd1_bm5000_v2")
    ccd_small_gif = _ensure_gif(CCD_SMALL_MP4, gif_dir, stem_override="ccd_s0p5_h0p1_bm5000_v2")
    ccd_mid1_gif = _ensure_gif(CCD_MID1_MP4, gif_dir, stem_override="ccd_s1p0_h0p1_bm5000_v2")
    ccd_mid2_gif = _ensure_gif(CCD_MID2_MP4, gif_dir, stem_override="ccd_s2p0_h0p1_bm5000_v2")
    ccd_large_gif = _ensure_gif(CCD_LARGE_MP4, gif_dir, stem_override="ccd_s3p0_h0p1_bm5000_v2")
    h1_bunny_radius_1x_gif = _ensure_gif(P0_BUNNY_OFF_MP4, gif_dir, stem_override="h1_bunny_radius_1x_v2")
    h1_bunny_radius_2x_gif = _ensure_gif(H1_BUNNY_RADIUS_2X_MP4, gif_dir, stem_override="h1_bunny_radius_2x_v2")
    p0_bunny_gif = _ensure_gif(P0_BUNNY_OFF_MP4, gif_dir, stem_override="p0_bunny_off_decoupled_v2")
    p0_bunny_final_gif = _ensure_gif(P0_BUNNY_FINAL_MP4, gif_dir, stem_override="p0_bunny_bodycenter8x_nopen_v2")
    p0_box_baseline_gif = _ensure_gif(P0_BOX_BASELINE_MP4, gif_dir, stem_override="p0_box_off_baseline_v2")
    p0_box_summary = _load_json_or_fallback(
        P0_BOX_OFF_SUMMARY,
        {
            "max_penetration_depth_rigid_m": 0.0038,
            "body_speed_max": 0.0028,
        },
    )
    p0_bunny_final_summary = _load_json_or_fallback(
        P0_BUNNY_FINAL_SUMMARY,
        {
            "max_penetration_depth_bunny_mesh_m": 0.00499,
            "body_speed_max": 0.0156,
        },
    )
    thin_ear_5x_summary = _load_json_or_fallback(
        THIN_EAR_5X_SUMMARY,
        {
            "contact_dist_scale": 5.0,
            "max_penetration_depth_bunny_mesh_m": 0.004997673444449902,
            "final_penetration_p99_bunny_mesh_m": 0.00040089187677949667,
            "rendered_frame_count": 40,
            "video_duration_target_sec": 1.334,
        },
    )
    thin_ear_10x_summary = _load_json_or_fallback(
        THIN_EAR_10X_SUMMARY,
        {
            "contact_dist_scale": 10.0,
            "max_penetration_depth_bunny_mesh_m": 0.005719680339097977,
            "final_penetration_p99_bunny_mesh_m": 0.0023300147056579596,
            "rendered_frame_count": 40,
            "video_duration_target_sec": 1.333333,
        },
    )
    bridge_flow_png = _make_bridge_pipeline_diagram(image_dir / "bridge_pipeline_diagram_v3.png")
    body_shape_png = _make_body_shape_mesh_diagram(image_dir / "body_shape_mesh_diagram_v1.png")

    _title_slide(prs, "PhysTwin-Full Meeting", ["Takeaway: this deck now follows four claims only."])
    _body(prs, "Four Claims, Not A Grab Bag", [
        "Takeaway: this meeting keeps only the pages that defend a claim.",
        "Claim 1-2: the bridge already works; bunny penetration is still an external contact + geometry problem.",
        "Claim 3-4: self-collision is not parity; keep multi-deformable and drop MPM two-way.",
    ])

    _section(prs, "The Bridge Already Works For Deformable-Rigid Interaction")
    _picture_with_footer(prs, "A Short IR Pipeline Reconstructs PhysTwin Objects In Newton", bridge_flow_png, [
        "Takeaway: PhysTwin exports a neutral IR, and Newton rebuilds spring-mass, rigid, and MPM in separate native worlds.",
        "Read left to right: PhysTwin outputs -> Bridge IR -> Newton reconstruction -> interaction layer.",
    ])
    _picture_with_footer(prs, "Mesh Defines Contact While The Rigid State Moves The Bunny", body_shape_png, [
        "Takeaway: explain this page in plain physics words first: cloth particles contact the bunny mesh surface.",
        "Only as a Newton-internal note: the contact geometry lives on a shape attached to the rigid body, which receives the reaction wrench.",
    ])
    _picture_with_footer(prs, "Newton Already Detects Particle-Mesh Contact On The Bunny", create_soft_png, [
        "Takeaway: Newton already detects `(particle, shape)` pairs directly, and here the bunny collision geometry is a mesh.",
        "This is the `closest-point / signed-normal` entry point for deformable-rigid interaction.",
    ])
    _picture_with_footer(prs, "Newton Uses **Bidirectional** Penalty Forces Between Cloth And Bunny", particle_body_png, [
        "Takeaway: Newton applies **Bidirectional** penalty forces, and the rigid bunny receives the equal-and-opposite reaction.",
        "**Elastic (弹性)**, **Damping (阻尼)**, and **Friction (摩擦力)** are all explicit parts of the same contact update.",
    ])
    _gif_single(
        prs,
        "Mainline 1 Result: The Bridge Already Works",
        rigid_h1_gif,
        "Takeaway: object-only rope | rope particle mass = 1.0 kg | bunny mass = 5.0 kg | drop height = 5.0 m | gravity on | deformable-rigid interaction is already working in Newton.",
    )
    _gif_compare_large(
        prs,
        "Stronger Particle-Shape Contact Builds Support Earlier",
        "Baseline `shape_contact_scale` = 1x alpha",
        h1_bunny_radius_1x_gif,
        "Scaled `shape_contact_scale` = 2x alpha",
        h1_bunny_radius_2x_gif,
        footer="self-collision OFF | cloth total mass = 0.1 kg | bunny mass = 0.5 kg | center drop",
    )

    _section(prs, "Bunny Penetration Is Still An External Contact Problem")
    _pic_twocol(
        prs,
        "Continuous Mesh Contact Still Depends On `particle_radius`",
        "Detection: continuous mesh query still starts at d < margin + radius",
        mesh_query_png,
        "Force: penalty penetration still uses c = signed distance - particle radius",
        newton_contact_png,
    )
    _gif_compare_large(
        prs,
        "The Bunny Is Harder Than A Thick Box Because Its Support Patch Is Thin",
        "Bunny baseline: same cloth, local penetration remains visible on the curved top patch",
        p0_bunny_gif,
        "Thick-box baseline: same cloth, broader support patch already reduces penetration",
        p0_box_baseline_gif,
        footer="self-collision OFF | cloth total mass = 0.1 kg | rigid mass = 0.5 kg | center drop | baseline particle-shape contact | decoupled shape materials",
    )
    _gif_compare_large(
        prs,
        "The Best Bunny Case Is Close To A Working Point, Not A Full Fix",
        "Bunny baseline: center drop, baseline `shape_contact_scale`, local penetration still visible",
        p0_bunny_gif,
        "Best current bunny case: body-center target plus 8x `shape_contact_scale` reduces penetration below the threshold",
        p0_bunny_final_gif,
        footer="self-collision OFF | cloth total mass = 0.1 kg | rigid mass = 0.5 kg | body-center target | `shape_contact_scale` = 8x | decoupled shape materials",
    )
    _body(prs, "The 64-Run OFF Grid Separates Main Effects Cleanly", [
        "Takeaway: this grid is for attribution, not brute-force search.",
        "We sweep only three axes: `contact_collision_dist`, drop height, and bunny mass.",
        "Everything else stays fixed in OFF mode so the main effects remain comparable.",
    ])
    _gif_grid_2x2(prs, "Heavier Bunny Helps, But Is Not The Main Lever", [
        ("Bunny mass = 1 kg", mass_light_gif),
        ("Bunny mass = 10 kg", mass_mid1_gif),
        ("Bunny mass = 1000 kg", mass_mid2_gif),
        ("Bunny mass = 5000 kg", mass_heavy_gif),
    ], common_settings="self-collision OFF | object mass = 1 kg / particle | `contact_collision_dist` scale = 1.0 | drop height = 0.1 m")
    _gif_grid_2x2(prs, "Lower Drop Height Helps, But Does Not Fix The Problem", [
        ("Drop height = 0.1 m", height_low_gif),
        ("Drop height = 0.25 m", height_mid1_gif),
        ("Drop height = 0.5 m", height_mid2_gif),
        ("Drop height = 1.0 m", height_high_gif),
    ], common_settings="self-collision OFF | object mass = 1 kg / particle | bunny mass = 5000 kg | `contact_collision_dist` scale = 1.0")
    _gif_grid_2x2(prs, "Larger `contact_collision_dist` Helps Most In OFF", [
        ("`contact_collision_dist` scale = 0.5", ccd_small_gif),
        ("`contact_collision_dist` scale = 1.0", ccd_mid1_gif),
        ("`contact_collision_dist` scale = 2.0", ccd_mid2_gif),
        ("`contact_collision_dist` scale = 3.0", ccd_large_gif),
    ], common_settings="self-collision OFF | object mass = 1 kg / particle | bunny mass = 5000 kg | drop height = 0.1 m")
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, "Across 64 OFF Runs, `contact_collision_dist` Dominates", size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)
    _add_pic(s, factor_chart, left, top, width, height)
    _add_label(
        s,
        int(left + 20),
        int(top + height - Inches(0.45)),
        int(width - 40),
        Inches(0.35),
        "Grid settings: self-collision OFF | object mass = 1 kg / particle | `contact_collision_dist` = 0.5 / 1.0 / 2.0 / 3.0 | bunny mass = 1 / 10 / 1000 / 5000 kg | drop height = 0.1 / 0.25 / 0.5 / 1.0 m",
        font_size=12,
        bold=False,
    )
    _gif_twocol(
        prs,
        "Latest Low-Mass OFF Runs Are Reference Cases, Not A New Claim",
        "Bunny mass = 0.1 kg",
        advisor_bm0p1_gif,
        "Bunny mass = 0.5 kg",
        advisor_bm0p5_gif,
        common_settings="self-collision OFF | `mass-spring-scale` = alpha | `shape_contact_scale` = alpha | `shape_contact_damping_multiplier` = 1 | `contact_dist_scale` = 5 | decoupled shape materials | total cloth mass = 0.1 kg | drop height = 0.5 m",
    )
    _gif_twocol(
        prs,
        "Thin-Ear Counterexample: Even 5x And 10x Radius Still Penetrate",
        (
            "5x radius close-up: non-zero penetration remains\n"
            f"max bunny-mesh penetration = {float(thin_ear_5x_summary['max_penetration_depth_bunny_mesh_m']) * 1000:.2f} mm"
        ),
        thin_ear_5x_gif,
        (
            "10x radius close-up: non-zero penetration still remains\n"
            f"max bunny-mesh penetration = {float(thin_ear_10x_summary['max_penetration_depth_bunny_mesh_m']) * 1000:.2f} mm"
        ),
        thin_ear_10x_gif,
        common_settings=(
            "close-up camera | no slowdown | video ends about 1s after first cloth-bunny contact | "
            "cloth total mass = 0.1 kg | bunny mass = 0.5 kg | "
            f"5x p99 = {float(thin_ear_5x_summary['final_penetration_p99_bunny_mesh_m']) * 1000:.2f} mm | "
            f"10x p99 = {float(thin_ear_10x_summary['final_penetration_p99_bunny_mesh_m']) * 1000:.2f} mm"
        ),
    )
    _boxed_bullets(prs, "Mainline 2 Result: Bunny Penetration Is A Support-Quality Problem", "Takeaway", [
        "Takeaway: the bridge already has continuous mesh contact, but bunny penetration is still an external support quality problem.",
        f"Best thick-box case: {p0_box_summary['max_penetration_depth_rigid_m']*1000:.1f} mm max penetration with stable rigid support.",
        f"Best bunny case: {p0_bunny_final_summary['max_penetration_depth_bunny_mesh_m']*1000:.2f} mm max penetration after moving to body-center and increasing contact strength; latest OFF reruns also preserve rigid-rigid support.",
        f"Thin-ear counterexample: even 10x radius still leaves {float(thin_ear_10x_summary['max_penetration_depth_bunny_mesh_m']) * 1000:.2f} mm bunny-mesh penetration, so radius helps but does not fully solve thin geometry.",
    ])

    _section(prs, "Self-Collision Is Not A PhysTwin Parity Target")
    _body(prs, "Three Collision Terms Explain Why ON Is Not Parity", [
        "**Gate**: PhysTwin only activates while pairs still approach each other.",
        "**Update**: PhysTwin writes velocity directly, while Newton first accumulates force.",
        "**Force**: Newton ON is overlap-triggered penalty contact, not impulse correction.",
    ])
    _pic_twocol(
        prs,
        "PhysTwin And Newton Use Different Self-Collision Models",
        "PhysTwin\n`dis_len < collision_dist`\n`wp.dot(dis, relative_v) < -1e-4`\n`v_new[tid] = v1 - J_average / m1`",
        phys_png,
        "Newton\n`err = d - radius - particle_radius[index]`\n`if err <= k_cohesion:`\n`f += particle_force(...)`",
        newton_self_png,
    )
    _boxed_bullets(prs, "Mainline 3 Result: ON Is Not Evidence Of Parity", "Takeaway", [
        "Takeaway: Newton native self-collision is not the PhysTwin self-collision model.",
        "PhysTwin is approach-gated and impulse-style, while Newton native self-collision is overlap-triggered and penalty-force based.",
        "So `self-collision ON` should not dominate the story and cannot be used as parity evidence.",
    ])

    _section(prs, "Keep Multi-Deformable Interaction; Drop MPM Two-Way")
    _body(prs, "Pair Filtering Makes Multi-Deformable Interaction Viable", [
        "Takeaway: multi-deformable interaction becomes viable only after explicit pair filtering.",
        "Native all-particle contact makes same-object neighbors fight springs.",
        "We therefore keep one Newton particle system, but filter to rope-vs-sloth pairs only.",
    ])
    _picture_with_footer(prs, "We Changed Pair Selection, Not The Physics Model", cross_rope_png, [
        "Takeaway: we keep Newton's HashGrid and penalty model, and only change pair selection.",
        "This is why rope-vs-sloth stays stable without rewriting the whole physics model.",
    ])
    _gif_twocol(
        prs,
        "Filtered Cross-Object Contact Already Works In Two Demos",
        "Two ropes on ground\nrope particle mass = 1.0 kg | vertical gap = 0.35 m",
        rope_ground_gif,
        "One rope + grounded sloth\nrope drop height = 0.35 m | sloth rests on ground",
        multi_rope_gif,
        common_settings="No rigid bunny | gravity = 9.8 m/s^2 | drag on | Newton particle contact filtered to cross-object pairs only",
    )
    _body(prs, "Mainline 4 Result: Multi-Deformable Interaction Is Validated", [
        "Takeaway: multi-deformable interaction is feasible once pair filtering is explicit.",
        "Both rope+rope and rope+sloth interactions work without relying on a rigid probe.",
        "This is the cleaner next step for future cloth-cloth interaction studies.",
    ])
    _gif_single(
        prs,
        "One-Way MPM Works, But It Solves A Narrower Problem",
        mpm_sloth_gif,
        "Takeaway: controller-driven sloth | one-way MPM sand pile | frames = 120 | sand radii = 0.60 / 0.80 m | sand height = 0.36 m | soft gravity = 0.0",
    )
    _body(prs, "MPM Two-Way Is Not The Main Path", [
        "Takeaway: we should not spend the main effort on MPM two-way coupling for this project.",
        "Mass ratio is too extreme, the ground SDF competes with patch colliders, and local patch support is too weak.",
        "So the main path stays on SemiImplicit external contact, not MPM two-way support.",
    ])
    _boxed_bullets(prs, "Summary: Four Claims, One Story", "Takeaway", [
        "Takeaway: the bridge already works; the remaining bottleneck is external contact quality on bunny geometry.",
        "Self-collision is not a PhysTwin parity target.",
        "Keep multi-deformable interaction, and deprioritize MPM two-way.",
    ])
    _body(prs, "Next Step: Fewer Slides, Stronger Evidence", [
        "Takeaway: next week should produce a smaller, cleaner, harder-to-argue-against deck.",
        "Keep the flow page, the body/shape explanation page, and only the strongest support demos.",
        "Replace weak demos with one clean bridge demo, one clean bunny support demo, and a 1-minute opening summary.",
    ])


def build_deck_reference(
    prs: Presentation, rows: list[dict], gif_dir: Path, image_dir: Path
) -> None:
    _delete_all_slides(prs)

    factor_chart = _make_factor_chart(rows, image_dir / "factor_chart.png")
    phys_png = _code_excerpt_image(
        PHYSTWIN_SIM,
        "phystwin_self_collision_ref_v3",
        _extract_code_segments(
            PHYSTWIN_SIM,
            [(196, 200), (205, 206), (220, 220), (294, 296)],
            highlight_lines={196, 199, 205, 220, 294},
        ),
        image_dir,
    )
    create_soft_png = _code_excerpt_image(
        NEWTON_GEOM_KERNELS,
        "newton_create_soft_contacts_ref_v3",
        _extract_code_segments(
            NEWTON_GEOM_KERNELS,
            [(1104, 1107), (1128, 1133)],
            highlight_lines={1104, 1105, 1106, 1128, 1132, 1133},
        ),
        image_dir,
    )
    particle_body_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_eval_particle_body_contact_ref_v3",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(203, 213), (233, 236)],
            highlight_lines={203, 204, 210, 211, 233, 236},
        ),
        image_dir,
    )
    newton_self_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_self_collision_ref_v3",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(76, 82)],
            highlight_lines={76, 78, 79, 82},
        ),
        image_dir,
    )
    mesh_query_png = _code_excerpt_image(
        NEWTON_GEOM_KERNELS,
        "newton_mesh_query_ref_v3",
        _extract_code_segments(
            NEWTON_GEOM_KERNELS,
            [(1104, 1107), (1128, 1133)],
            highlight_lines={1104, 1105, 1106, 1128, 1132, 1133},
        ),
        image_dir,
    )
    newton_contact_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "newton_contact_force_ref_v3",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(203, 213), (233, 236)],
            highlight_lines={203, 204, 210, 211, 233, 236},
        ),
        image_dir,
    )
    newton_spring_force_png = _code_excerpt_image(
        NEWTON_KERNELS_PARTICLE,
        "newton_spring_force_ref_v2",
        _extract_code_segments(
            NEWTON_KERNELS_PARTICLE,
            [(27, 29), (46, 50)],
            highlight_lines={27, 28, 46, 47, 50},
        ),
        image_dir,
    )
    newton_mass_integrate_png = _code_excerpt_image(
        NEWTON_SOLVER_BASE,
        "newton_particle_integrate_ref_v2",
        _extract_code_segments(
            NEWTON_SOLVER_BASE,
            [(34, 44)],
            highlight_lines={34, 38, 39, 44},
        ),
        image_dir,
    )
    import_particles_springs_png = _code_excerpt_image(
        NEWTON_IMPORT_IR,
        "import_particles_springs_ref_v1",
        _extract_code_segments(
            NEWTON_IMPORT_IR,
            [(1651, 1658), (1674, 1678)],
            highlight_lines={1652, 1655, 1656, 1677, 1678},
        ),
        image_dir,
    )
    import_body_mesh_png = _code_excerpt_image(
        NEWTON_IMPORT_IR,
        "import_body_mesh_ref_v1",
        _extract_code_segments(
            NEWTON_IMPORT_IR,
            [(1691, 1693), (1705, 1706), (1732, 1738), (1762, 1762)],
            highlight_lines={1691, 1693, 1733, 1736, 1762},
        ),
        image_dir,
    )
    mesh_contact_state_png = _code_excerpt_image(
        NEWTON_GEOM_KERNELS,
        "mesh_contact_state_ref_v1",
        _extract_code_segments(
            NEWTON_GEOM_KERNELS,
            [(1128, 1142)],
            highlight_lines={1128, 1133, 1136, 1138, 1142},
        ),
        image_dir,
    )
    body_wrench_png = _code_excerpt_image(
        NEWTON_KERNELS_CONTACT,
        "body_wrench_ref_v1",
        _extract_code_segments(
            NEWTON_KERNELS_CONTACT,
            [(251, 261)],
            highlight_lines={251, 253, 255, 259, 261},
        ),
        image_dir,
    )
    cross_rope_png = _code_excerpt_image(
        ROOT / "Newton" / "phystwin_bridge" / "demos" / "demo_rope_sloth_ground_contact.py",
        "cross_object_contact_ref_v3",
        _extract_code_segments(
            ROOT / "Newton" / "phystwin_bridge" / "demos" / "demo_rope_sloth_ground_contact.py",
            [(53, 67), (73, 79)],
            highlight_lines={53, 66, 67, 73, 74, 77},
        ),
        image_dir,
    )

    old_assets_gif_dir = FORMAL_SLIDE_DIR / "slides_assets" / "gif"
    rigid_h1_gif = (
        _ensure_gif(CURRENT_ROPE_BUNNY_MP4, gif_dir, stem_override="rope_bunny_current_v2")
        if CURRENT_ROPE_BUNNY_MP4.exists()
        else old_assets_gif_dir / "bunny_drop_m5.gif"
    )
    rope_ground_mp4 = ROOT / "tmp" / "rope_ground_drop_full_close_20260322" / "rope_ground_drop.mp4"
    rope_ground_gif = (
        _ensure_gif(rope_ground_mp4, gif_dir, stem_override="rope_ground_drop_v2")
        if rope_ground_mp4.exists()
        else old_assets_gif_dir / "rope_drag_on_overlay1x3.gif"
    )
    multi_rope_mp4 = ROOT / "tmp" / "rope_sloth_dualsource_drop1p0_r2_20260322" / "rope_sloth_ground_drop.mp4"
    multi_rope_gif = (
        _ensure_gif(multi_rope_mp4, gif_dir, stem_override="rope_sloth_ground_drop_v2")
        if multi_rope_mp4.exists()
        else old_assets_gif_dir / "sloth_overlay1x3.gif"
    )
    two_ropes_weight_case1_mp4 = (
        TWO_ROPES_WEIGHT_COMPARE_DIR
        / "lower0p5_upper0p1"
        / "two_ropes_lower0p5_upper0p1.mp4"
    )
    two_ropes_weight_case2_mp4 = (
        TWO_ROPES_WEIGHT_COMPARE_DIR
        / "lower5_upper5"
        / "two_ropes_lower5_upper5.mp4"
    )
    two_ropes_weight_case1_summary_path = (
        TWO_ROPES_WEIGHT_COMPARE_DIR
        / "lower0p5_upper0p1"
        / "two_ropes_lower0p5_upper0p1_summary.json"
    )
    two_ropes_weight_case2_summary_path = (
        TWO_ROPES_WEIGHT_COMPARE_DIR
        / "lower5_upper5"
        / "two_ropes_lower5_upper5_summary.json"
    )
    two_ropes_weight_case1_gif = _ensure_gif(
        two_ropes_weight_case1_mp4, gif_dir, stem_override="two_ropes_lower0p5_upper0p1_v1"
    )
    two_ropes_weight_case2_gif = _ensure_gif(
        two_ropes_weight_case2_mp4, gif_dir, stem_override="two_ropes_lower5_upper5_v1"
    )
    two_ropes_weight_case1_summary = _load_json_or_fallback(
        two_ropes_weight_case1_summary_path,
        {
            "rope_total_masses": [0.5, 0.1],
            "video_duration_target_sec": 1.6,
        },
    )
    two_ropes_weight_case2_summary = _load_json_or_fallback(
        two_ropes_weight_case2_summary_path,
        {
            "rope_total_masses": [5.0, 5.0],
            "video_duration_target_sec": 2.0,
        },
    )
    two_ropes_box_stack_gif = _ensure_gif(
        TWO_ROPES_BOX_STACK_MP4, gif_dir, stem_override="two_ropes_box_stack_v1"
    )
    two_ropes_box_stack_summary = _load_json_or_fallback(
        TWO_ROPES_BOX_STACK_SUMMARY,
        {
            "rope_total_masses": [0.1, 0.1],
            "box_total_mass_kg": 0.1,
            "box_half_extents": [0.05, 0.05, 0.05],
            "drop_height_m": 0.15,
            "body_speed_max": 0.032788608223199844,
        },
    )
    three_box_stack_gif = _ensure_gif(
        THREE_BOX_STACK_MP4, gif_dir, stem_override="three_box_stack_v1"
    )
    three_box_stack_summary = _load_json_or_fallback(
        THREE_BOX_STACK_SUMMARY,
        {
            "rope_total_masses": [0.1, 0.1],
            "box_count": 3,
            "box_total_mass_kg": 0.3,
            "box_half_extents": [0.05, 0.05, 0.05],
            "drop_height_m": 0.15,
            "body_speed_max": 0.0780295729637146,
        },
    )
    robot_rope_franka_gif = _ensure_gif(
        ROBOT_ROPE_FRANKA_MP4, gif_dir, stem_override="robot_rope_franka_lift_release_v1"
    )
    robot_rope_franka_summary = _load_json_or_fallback(
        ROBOT_ROPE_FRANKA_SUMMARY,
        {
            "robot_geometry": "native_franka",
            "rope_total_mass": 1753.0,
            "task": "lift_release",
            "render_mode": "presentation",
            "first_contact_time_s": 0.0884,
            "contact_duration_s": 0.0904,
            "release_frame": 1001,
            "contact_active_frames": 226,
            "min_clearance_min_m": -0.0027090779040008783,
            "rope_com_displacement_m": 0.15631510317325592,
            "rope_mid_segment_peak_z_delta_m": 0.009135246276855469,
        },
    )
    mpm_sloth_gif = (
        _ensure_gif(MPM_SLOTH_MP4, gif_dir, stem_override="sloth_mpm_sand_one_way_wide_v2")
        if MPM_SLOTH_MP4.exists()
        else old_assets_gif_dir / "sloth_overlay1x3.gif"
    )
    p0_bunny_gif = (
        _ensure_gif(P0_BUNNY_OFF_MP4, gif_dir, stem_override="p0_bunny_off_decoupled_v2")
        if P0_BUNNY_OFF_MP4.exists()
        else old_assets_gif_dir / "cloth_overlay1x3.gif"
    )
    p0_bunny_final_gif = (
        _ensure_gif(P0_BUNNY_FINAL_MP4, gif_dir, stem_override="p0_bunny_bodycenter8x_nopen_v2")
        if P0_BUNNY_FINAL_MP4.exists()
        else old_assets_gif_dir / "bunny_drop_m5.gif"
    )
    p0_box_baseline_gif = (
        _ensure_gif(P0_BOX_BASELINE_MP4, gif_dir, stem_override="p0_box_off_baseline_v2")
        if P0_BOX_BASELINE_MP4.exists()
        else old_assets_gif_dir / "rigid_probe_m5_scene_xz.gif"
    )
    thin_ear_5x_gif = _ensure_gif(THIN_EAR_5X_MP4, gif_dir, stem_override="thin_ear_ccd5x_v3")
    thin_ear_10x_gif = _ensure_gif(THIN_EAR_10X_MP4, gif_dir, stem_override="thin_ear_ccd10x_v3")
    thin_ear_5x_summary = _load_json_or_fallback(
        THIN_EAR_5X_SUMMARY,
        {
            "contact_dist_scale": 5.0,
            "max_penetration_depth_bunny_mesh_m": 0.004997673444449902,
            "final_penetration_p99_bunny_mesh_m": 0.00040089187677949667,
            "rendered_frame_count": 40,
            "video_duration_target_sec": 1.334,
        },
    )
    thin_ear_10x_summary = _load_json_or_fallback(
        THIN_EAR_10X_SUMMARY,
        {
            "contact_dist_scale": 10.0,
            "max_penetration_depth_bunny_mesh_m": 0.005719680339097977,
            "final_penetration_p99_bunny_mesh_m": 0.0023300147056579596,
            "rendered_frame_count": 40,
            "video_duration_target_sec": 1.333333,
        },
    )
    bridge_flow_png = _make_bridge_pipeline_diagram(image_dir / "bridge_pipeline_diagram_v4.png")
    recall_bridge_png = _make_recall_bridge_diagram(image_dir / "recall_bridge_diagram_v1.png")
    body_shape_png = _make_body_shape_mesh_diagram(image_dir / "body_shape_mesh_diagram_v4.png")
    force_diag_summary = _load_json_or_fallback(
        FORCE_DIAG_BOX_SUMMARY,
        {
            "contact_node_count": 2097,
            "contact_nodes_with_wrong_direction": 0,
            "contact_nodes_with_inward_acceleration": 0,
            "contact_nodes_force_below_stop": 2097,
            "contact_nodes_internal_dominates": 0,
            "dominant_issue_guess": "insufficient_contact_magnitude",
            "topk_particle_records": [
                {
                    "external_force_normal": 0.01804775558412075,
                    "stop_force_required_normal": 1.4134559631347656,
                }
            ],
        },
    )
    rope_profile_summary = _load_json_or_fallback(
        ROPE_PROFILE_FULL_JSON,
        {
            "trajectory_frames": 264,
            "total_substeps": 175421,
            "runs": [
                {"wall_ms": 26402.553208000427, "sim_time_sec": 8.77104977842464},
                {"wall_ms": 26523.329532999924, "sim_time_sec": 8.77104977842464},
                {"wall_ms": 26691.76469499962, "sim_time_sec": 8.77104977842464},
            ],
            "aggregate": {
                "total_substep": {"mean_of_run_means_ms": 0.13307293066524303},
                "solver_step": {"mean_of_run_means_ms": 0.06466363214421334},
                "model_collide": {"mean_of_run_means_ms": 0.0},
                "particle_grid_build": {"mean_of_run_means_ms": 0.0},
            },
            "bottleneck_ranked_ops": [
                {"op_name": "integrate_particles", "mean_of_run_means_ms": 0.020095355347191538, "call_count_total": 526263},
                {"op_name": "spring_forces", "mean_of_run_means_ms": 0.019438087868771287, "call_count_total": 526263},
                {"op_name": "write_kinematic_state", "mean_of_run_means_ms": 0.016849647324359788, "call_count_total": 526263},
                {"op_name": "drag_correction", "mean_of_run_means_ms": 0.015291624893795, "call_count_total": 526263},
                {"op_name": "integrate_bodies", "mean_of_run_means_ms": 0.001268380883015826, "call_count_total": 526263},
            ],
        },
    )
    force_diag_png = _make_force_diagnostic_mechanism_diagram(
        image_dir / "force_diagnostic_mechanism_v1.png",
        FORCE_DIAG_BOX_SNAPSHOT,
        force_diag_summary,
    )
    rope_profile_png = _make_rope_profile_summary_diagram(
        image_dir / "rope_profile_summary_v1.png",
        rope_profile_summary,
    )
    old_cloth_gif = old_assets_gif_dir / "cloth_overlay1x3.gif"
    old_zebra_gif = old_assets_gif_dir / "zebra_overlay1x3.gif"
    old_sloth_gif = old_assets_gif_dir / "sloth_overlay1x3.gif"
    old_rope_gif = old_assets_gif_dir / "rope_drag_on_overlay1x3.gif"
    recall_bunny_rope_m5_gif = (
        _ensure_gif(ROPE_BUNNY_FULL_M5_MP4, gif_dir, stem_override="rope_bunny_close_1m_m5_v2")
        if ROPE_BUNNY_FULL_M5_MP4.exists()
        else old_assets_gif_dir / "bunny_drop_m5.gif"
    )
    recall_bunny_rope_m500_gif = (
        _ensure_gif(ROPE_BUNNY_FULL_M500_MP4, gif_dir, stem_override="rope_bunny_close_1m_m500_v2")
        if ROPE_BUNNY_FULL_M500_MP4.exists()
        else old_assets_gif_dir / "bunny_drop_m500.gif"
    )
    rope_bunny_weight_1kg_gif = _ensure_gif(
        ROPE_BUNNY_WEIGHT_1KG_MP4, gif_dir, stem_override="rope_bunny_total1kg_m5_v1"
    )
    rope_bunny_weight_5kg_gif = _ensure_gif(
        ROPE_BUNNY_WEIGHT_5KG_MP4, gif_dir, stem_override="rope_bunny_total5kg_m5_v1"
    )
    rope_bunny_weight_1kg_summary = _load_json_or_fallback(
        ROPE_BUNNY_WEIGHT_1KG_SUMMARY,
        {
            "total_object_mass": 1.0,
            "body_speed_max": 0.04798571392893791,
        },
    )
    rope_bunny_weight_5kg_summary = _load_json_or_fallback(
        ROPE_BUNNY_WEIGHT_5KG_SUMMARY,
        {
            "total_object_mass": 5.0,
            "body_speed_max": 0.24445965886116028,
        },
    )
    cloth_bunny_weight_1kg_gif = _ensure_gif(
        CLOTH_BUNNY_WEIGHT_1KG_MP4, gif_dir, stem_override="cloth_bunny_total1kg_m5_v1"
    )
    cloth_bunny_weight_5kg_gif = _ensure_gif(
        CLOTH_BUNNY_WEIGHT_5KG_MP4, gif_dir, stem_override="cloth_bunny_total5kg_m5_v1"
    )
    cloth_bunny_weight_1kg_summary = _load_json_or_fallback(
        CLOTH_BUNNY_WEIGHT_1KG_SUMMARY,
        {
            "total_object_mass": 1.0,
            "body_speed_max": 0.016717961058020592,
        },
    )
    cloth_bunny_weight_5kg_summary = _load_json_or_fallback(
        CLOTH_BUNNY_WEIGHT_5KG_SUMMARY,
        {
            "total_object_mass": 5.0,
            "body_speed_max": 0.05247604846954346,
        },
    )
    cloth_bunny_weight_0p1kg_gif = _ensure_gif(
        CLOTH_BUNNY_WEIGHT_0P1KG_MP4, gif_dir, stem_override="cloth_bunny_total0p1kg_m5_v1"
    )
    cloth_bunny_weight_0p5kg_gif = _ensure_gif(
        CLOTH_BUNNY_WEIGHT_0P5KG_MP4, gif_dir, stem_override="cloth_bunny_total0p5kg_m5_v1"
    )
    cloth_bunny_weight_0p1kg_summary = _load_json_or_fallback(
        CLOTH_BUNNY_WEIGHT_0P1KG_SUMMARY,
        {
            "total_object_mass": 0.1,
            "body_speed_max": 0.013971414417028427,
            "video_duration_target_sec": 1.334,
        },
    )
    cloth_bunny_weight_0p5kg_summary = _load_json_or_fallback(
        CLOTH_BUNNY_WEIGHT_0P5KG_SUMMARY,
        {
            "total_object_mass": 0.5,
            "body_speed_max": 0.013971414417028427,
            "video_duration_target_sec": 1.334,
        },
    )
    cloth_bunny_radius_cases: list[tuple[str, Path, dict[str, float]]] = []
    for scale in (1, 2, 3, 4, 8, 10):
        mp4 = (
            CLOTH_BUNNY_RADIUS_COMPARE_DIR
            / f"r{scale}x"
            / "self_off"
            / f"cloth_bunny_total0p1kg_r{scale}x_off_m5.mp4"
        )
        summary_path = (
            CLOTH_BUNNY_RADIUS_COMPARE_DIR
            / f"r{scale}x"
            / "self_off"
            / f"cloth_bunny_total0p1kg_r{scale}x_off_m5_summary.json"
        )
        gif = _ensure_gif(mp4, gif_dir, stem_override=f"cloth_bunny_total0p1kg_r{scale}x_m5_v1")
        summary = _load_json_or_fallback(
            summary_path,
            {
                "contact_collision_dist_used": 0.01016022264957428 * float(scale),
                "body_speed_max": 0.013971414417028427,
                "video_duration_target_sec": 1.334,
            },
        )
        cloth_bunny_radius_cases.append((f"{scale}x", gif, summary))
    cloth_rigid_compare_bunny_gif = _ensure_gif(
        CLOTH_RIGID_COMPARE_BUNNY_MP4, gif_dir, stem_override="cloth_rigid_compare_bunny_m5_v1"
    )
    cloth_rigid_compare_box_gif = _ensure_gif(
        CLOTH_RIGID_COMPARE_BOX_MP4, gif_dir, stem_override="cloth_rigid_compare_box_m5_v1"
    )
    cloth_rigid_compare_bunny_summary = _load_json_or_fallback(
        CLOTH_RIGID_COMPARE_BUNNY_SUMMARY,
        {
            "contact_collision_dist_used": 0.01016022264957428,
            "body_speed_max": 0.013971414417028427,
        },
    )
    cloth_rigid_compare_box_summary = _load_json_or_fallback(
        CLOTH_RIGID_COMPARE_BOX_SUMMARY,
        {
            "contact_collision_dist_used": 0.01016022264957428,
            "body_speed_max": 0.0037922172341495752,
        },
    )
    cloth_bunny_size_cases: list[tuple[str, Path, dict[str, float]]] = []
    for dir_tag, label, scale in (
        ("x1", "1x", 0.12),
        ("x1p5", "1.5x", 0.18),
        ("x2", "2x", 0.24),
        ("x3", "3x", 0.36),
    ):
        mp4 = (
            CLOTH_BUNNY_SIZE_COMPARE_DIR
            / dir_tag
            / "self_off"
            / f"cloth_bunny_total0p1kg_{dir_tag}_off_m5.mp4"
        )
        summary_path = (
            CLOTH_BUNNY_SIZE_COMPARE_DIR
            / dir_tag
            / "self_off"
            / f"cloth_bunny_total0p1kg_{dir_tag}_off_m5_summary.json"
        )
        gif = _ensure_gif(
            mp4,
            gif_dir,
            stem_override=f"cloth_bunny_total0p1kg_bunnysize_{dir_tag}_m5_v1",
        )
        summary = _load_json_or_fallback(
            summary_path,
            {
                "bunny_scale": scale,
                "bunny_top_z": 1.61894823 * scale,
                "contact_collision_dist_used": 0.01016022264957428,
            },
        )
        cloth_bunny_size_cases.append((label, gif, summary))
    cloth_bunny_dt_cases: list[tuple[str, Path, dict[str, float]]] = []
    for dir_tag, label, sim_dt, substeps in (
        ("dt1e5", "case 1", 1.0e-5, 3335),
        ("dt1e6", "case 2", 1.0e-6, 33350),
    ):
        mp4 = (
            CLOTH_BUNNY_DT_COMPARE_DIR
            / dir_tag
            / "self_off"
            / f"cloth_bunny_total0p1kg_{dir_tag}_off_m5.mp4"
        )
        summary_path = (
            CLOTH_BUNNY_DT_COMPARE_DIR
            / dir_tag
            / "self_off"
            / f"cloth_bunny_total0p1kg_{dir_tag}_off_m5_summary.json"
        )
        gif = _ensure_gif(
            mp4,
            gif_dir,
            stem_override=f"cloth_bunny_total0p1kg_dtcompare_{dir_tag}_m5_v1",
        )
        summary = _load_json_or_fallback(
            summary_path,
            {
                "sim_dt": sim_dt,
                "substeps": substeps,
                "contact_collision_dist_used": 0.01016022264957428,
            },
        )
        cloth_bunny_dt_cases.append((label, gif, summary))
    cloth_bunny_ultralight_cases: list[tuple[str, Path, dict[str, float]]] = []
    for dir_tag, label, total_mass in (
        ("0p01kg", "0.01 kg", 0.01),
        ("0p005kg", "0.005 kg", 0.005),
    ):
        mp4 = (
            CLOTH_BUNNY_ULTRALIGHT_COMPARE_DIR
            / dir_tag
            / "self_off"
            / f"cloth_bunny_total{dir_tag}_off_m5.mp4"
        )
        summary_path = (
            CLOTH_BUNNY_ULTRALIGHT_COMPARE_DIR
            / dir_tag
            / "self_off"
            / f"cloth_bunny_total{dir_tag}_off_m5_summary.json"
        )
        gif = _ensure_gif(
            mp4,
            gif_dir,
            stem_override=f"cloth_bunny_total{dir_tag}_m5_v1",
        )
        summary = _load_json_or_fallback(
            summary_path,
            {
                "total_object_mass": total_mass,
                "contact_collision_dist_used": 0.01016022264957428,
                "video_duration_target_sec": 1.334,
            },
        )
        cloth_bunny_ultralight_cases.append((label, gif, summary))

    _card_name_slide(prs, "Xinjie Zhang")
    recall_title = "Recall 1: PhysTwin \u2192 Newton Bridge: Working"
    _card_media_slide(prs, recall_title, old_cloth_gif)
    _card_media_slide(prs, recall_title, old_zebra_gif)
    _card_media_slide(prs, recall_title, old_sloth_gif)
    _card_media_slide(prs, recall_title, old_rope_gif)
    _card_media_slide(prs, recall_title, recall_bridge_png, media_top_offset=int(Inches(1.05)))
    _gif_twocol(
        prs,
        "Recall 2: Bunny+Rope",
        "Bunny mass = 5 kg",
        recall_bunny_rope_m5_gif,
        "Bunny mass = 500 kg",
        recall_bunny_rope_m500_gif,
        common_settings="rope per-particle mass = 1.0 kg | 1753 particles -> rope total mass = 1753 kg | no scale | drop height = 1.0 m | slowdown = 1x | close manual camera",
    )

    _pic_twocol(
        prs,
        "H1: Particle-Mesh Contact Is Already **Bidirectional**",
        "Detection: `create_soft_contacts` writes contact point / normal from the bunny mesh query",
        create_soft_png,
        "Force: `fn + fd + ft` is subtracted from the particle and added as force + torque to the rigid bunny",
        particle_body_png,
    )
    _pic_twocol(
        prs,
        "H1 Method Note: Why We Do Not Lower Rope Mass Alone",
        "Newton core: spring force uses `ke` and `kd` directly, so the force law does not compensate for a smaller mass",
        newton_spring_force_png,
        "Newton core: particle integrate uses `inv_mass`, so mass-only scaling changes `ke / m`, `kd / m`, and the rope time scale",
        newton_mass_integrate_png,
    )
    _pic_twocol(
        prs,
        "H1: Short IR Pipeline Reconstructs The Right Newton Objects",
        "Importer: IR object arrays become native Newton particles and springs through `builder.add_particles` and `builder.add_spring`",
        import_particles_springs_png,
        "Importer: the bunny becomes one rigid body plus one attached mesh shape before `builder.finalize()` hands the model to Newton",
        import_body_mesh_png,
    )
    _pic_twocol(
        prs,
        "H1: Mesh Defines Contact; Rigid State Moves The Bunny",
        "Geometry kernel: mesh query writes contact point / normal in body-aware contact buffers",
        mesh_contact_state_png,
        "Contact kernel: the same total contact updates the rigid bunny wrench accumulator, not just the cloth particle",
        body_wrench_png,
    )
    _gif_single(
        prs,
        "H1 Result: Deformable-Rigid Interaction Is Validated",
        rigid_h1_gif,
        "Takeaway: object-only rope still exchanges force with a novel rigid bunny, so the bridge already works at the interaction level.",
    )
    _gif_twocol_large(
        prs,
        "H1 Addendum: Rope+Bunny Still Works At 1 kg And 5 kg Total Rope Mass",
        (
            f"Rope total mass = {float(rope_bunny_weight_1kg_summary['total_object_mass']):.0f} kg\n"
            f"bunny body_speed_max = {float(rope_bunny_weight_1kg_summary['body_speed_max']):.2f} m/s"
        ),
        rope_bunny_weight_1kg_gif,
        (
            f"Rope total mass = {float(rope_bunny_weight_5kg_summary['total_object_mass']):.0f} kg\n"
            f"bunny body_speed_max = {float(rope_bunny_weight_5kg_summary['body_speed_max']):.2f} m/s"
        ),
        rope_bunny_weight_5kg_gif,
        common_settings=(
            "bunny mass = 5 kg | drop height = 1.0 m | same close manual camera | "
            "`auto-set-weight` computes one shared `weight_scale` for mass + spring + contact"
        ),
    )
    _gif_twocol_large(
        prs,
        "H1 Addendum: Cloth+Bunny OFF Still Works At 1 kg And 5 kg Total Cloth Mass",
        (
            f"Cloth total mass = {float(cloth_bunny_weight_1kg_summary['total_object_mass']):.0f} kg\n"
            f"bunny body_speed_max = {float(cloth_bunny_weight_1kg_summary['body_speed_max']):.3f} m/s"
        ),
        cloth_bunny_weight_1kg_gif,
        (
            f"Cloth total mass = {float(cloth_bunny_weight_5kg_summary['total_object_mass']):.0f} kg\n"
            f"bunny body_speed_max = {float(cloth_bunny_weight_5kg_summary['body_speed_max']):.3f} m/s"
        ),
        cloth_bunny_weight_5kg_gif,
        common_settings=(
            "self-collision OFF | bunny mass = 5 kg | drop height = 0.5 m | same cloth IR | "
            "`auto-set-weight` computes one shared `weight_scale` for mass + spring + contact"
        ),
    )
    _gif_twocol_large(
        prs,
        "H1 Addendum: Cloth+Bunny OFF Still Works At 0.1 kg And 0.5 kg Total Cloth Mass",
        (
            f"Cloth total mass = {float(cloth_bunny_weight_0p1kg_summary['total_object_mass']):.1f} kg\n"
            f"bunny body_speed_max = {float(cloth_bunny_weight_0p1kg_summary['body_speed_max']):.3f} m/s"
        ),
        cloth_bunny_weight_0p1kg_gif,
        (
            f"Cloth total mass = {float(cloth_bunny_weight_0p5kg_summary['total_object_mass']):.1f} kg\n"
            f"bunny body_speed_max = {float(cloth_bunny_weight_0p5kg_summary['body_speed_max']):.3f} m/s"
        ),
        cloth_bunny_weight_0p5kg_gif,
        common_settings=(
            "self-collision OFF | bunny mass = 5 kg | drop height = 0.5 m | "
            "video ends 1 second after first cloth-bunny contact | "
            "`auto-set-weight` computes one shared `weight_scale` for mass + spring + contact"
        ),
    )
    for page_title, left_idx, right_idx in [
        ("H2 Counterexample: 1x/2x Radius Still Penetrates Bunny", 0, 1),
        ("H2 Counterexample: 3x/4x Radius Still Penetrates Bunny", 2, 3),
        ("H2 Counterexample: 8x/10x Radius Still Penetrates Thin Ears", 4, 5),
    ]:
        left_tag, left_gif, left_summary = cloth_bunny_radius_cases[left_idx]
        right_tag, right_gif, right_summary = cloth_bunny_radius_cases[right_idx]
        left_radius = 0.5 * float(left_summary["contact_collision_dist_used"])
        right_radius = 0.5 * float(right_summary["contact_collision_dist_used"])
        _gif_twocol_large(
            prs,
            page_title,
            f"{left_tag} radius ({left_radius:.4f} m)",
            left_gif,
            f"{right_tag} radius ({right_radius:.4f} m)",
            right_gif,
            common_settings=(
                "self-collision OFF | cloth total mass = 0.1 kg | bunny mass = 5 kg | "
                "drop height = 0.5 m | video ends 1 second after first cloth-bunny contact"
            ),
        )
    rigid_compare_radius = 0.5 * float(cloth_rigid_compare_bunny_summary["contact_collision_dist_used"])
    _gif_twocol_large(
        prs,
        "H2 Counterexample: Bunny Is Harder Than Thick Box At Same Radius",
        f"Cloth + bunny (radius = {rigid_compare_radius:.4f} m)",
        cloth_rigid_compare_bunny_gif,
        f"Cloth + box (radius = {rigid_compare_radius:.4f} m)",
        cloth_rigid_compare_box_gif,
        common_settings=(
            "self-collision OFF | cloth total mass = 0.1 kg | rigid mass = 5 kg | "
            "drop height = 0.5 m | video ends 1 second after first cloth-rigid contact"
        ),
    )
    for page_title, left_idx, right_idx in [
        ("H2 Counterexample: 1x/1.5x Bunny Size Still Penetrates", 0, 1),
        ("H2 Counterexample: 2x/3x Bunny Size Still Penetrates", 2, 3),
    ]:
        left_tag, left_gif, left_summary = cloth_bunny_size_cases[left_idx]
        right_tag, right_gif, right_summary = cloth_bunny_size_cases[right_idx]
        _gif_twocol_large(
            prs,
            page_title,
            f"{left_tag} bunny size (height = {float(left_summary['bunny_top_z']):.3f} m)",
            left_gif,
            f"{right_tag} bunny size (height = {float(right_summary['bunny_top_z']):.3f} m)",
            right_gif,
            common_settings=(
                "self-collision OFF | cloth total mass = 0.1 kg | bunny mass = 5 kg | "
                "particle radius = 0.0051 m | drop height = 0.5 m | "
                "video ends 1 second after first cloth-bunny contact"
            ),
        )
    left_tag, left_gif, left_summary = cloth_bunny_dt_cases[0]
    right_tag, right_gif, right_summary = cloth_bunny_dt_cases[1]
    _gif_twocol_large(
        prs,
        "H2 Counterexample: Smaller `dt` Still Does Not Eliminate Penetration",
        (
            f"{left_tag}: sim_dt = {float(left_summary['sim_dt']):.0e}, "
            f"substeps = {int(left_summary['substeps'])}"
        ),
        left_gif,
        (
            f"{right_tag}: sim_dt = {float(right_summary['sim_dt']):.0e}, "
            f"substeps = {int(right_summary['substeps'])}"
        ),
        right_gif,
        common_settings=(
            "self-collision OFF | cloth total mass = 0.1 kg | bunny mass = 5 kg | "
            "particle radius = 0.0051 m | drop height = 0.5 m | "
            "video ends 1 second after first cloth-bunny contact"
        ),
    )
    left_tag, left_gif, left_summary = cloth_bunny_ultralight_cases[0]
    right_tag, right_gif, right_summary = cloth_bunny_ultralight_cases[1]
    _gif_twocol_large(
        prs,
        "H2 Counterexample: Even 0.01 kg And 0.005 kg Still Penetrate Bunny",
        f"Cloth total mass = {left_tag}",
        left_gif,
        f"Cloth total mass = {right_tag}",
        right_gif,
        common_settings=(
            "self-collision OFF | bunny mass = 5 kg | particle radius = 0.0051 m | "
            "drop height = 0.5 m | video ends 1 second after first cloth-bunny contact"
        ),
    )

    _body(prs, "H2: Bunny Penetration Is Mainly A Support Problem", [
        "We test whether bunny penetration is caused mainly by external support quality rather than by missing mesh contact.",
        "The key suspects are `particle_radius`, local bunny geometry, and the strength of particle-shape support.",
        "H2 is a progress question: which hypotheses survived after controlled OFF experiments?",
    ])
    _pic_twocol(
        prs,
        "H2 Evidence: Continuous Mesh Contact Still Depends On `particle_radius`",
        "Detection: mesh contact still activates only when `d < margin + radius`",
        mesh_query_png,
        "Force: penalty depth still uses `c = signed distance - particle_radius`",
        newton_contact_png,
    )
    _body(prs, "H2: 64-Run OFF Grid Separates Three Main Effects", [
        "We sweep only three axes: `contact_collision_dist`, drop height, and bunny mass.",
        "Everything else stays fixed in OFF mode so the main effects remain comparable instead of anecdotal.",
        "This turns H2 into hypothesis testing rather than cherry-picking videos.",
    ])
    s = prs.slides.add_slide(_layout(prs, LY_BODY))
    _clear_placeholders(s)
    box = s.shapes.add_textbox(REF_TITLE_LEFT, REF_TITLE_TOP, REF_TITLE_W, REF_TITLE_H)
    _set_title_textbox(box, "H2: `contact_collision_dist` Dominates Across 64 OFF Runs", size_pt=28)
    left, top, width, height = _detach_body_placeholder(s)
    _add_pic(s, factor_chart, left, top, width, height)
    _add_label(
        s,
        int(left + 20),
        int(top + height - Inches(0.45)),
        int(width - 40),
        Inches(0.35),
        "Takeaway: larger `contact_collision_dist` (larger Newton `particle_radius`) helps most; bunny mass helps, but height alone does not solve the issue.",
        font_size=12,
        bold=False,
    )
    _gif_compare_large(
        prs,
        "H2: Bunny Now Has A Near-Working Point",
        "Bunny baseline: center drop with baseline particle-shape support still shows local penetration",
        p0_bunny_gif,
        "Best current bunny case: body-center target plus stronger particle-shape support reaches a near-working point",
        p0_bunny_final_gif,
        footer="Takeaway: latest decoupled alpha-scaled low-mass reruns also preserve rigid-rigid support, so the remaining gap is support quality, not rigid sinking.",
    )
    _picture_with_footer(
        prs,
        "H2 Mechanism: Trigger-Substep Force Diagnostic Separates Contact vs Spring",
        force_diag_png,
        [
            "white = outward normal | red = external contact | blue = spring-maintenance force | green = total acceleration",
            "Clean box control sanity check: wrong-direction nodes = 0 and inward-acceleration nodes = 0, so the diagnostic is ready to test whether bunny ears fail by normal direction or by insufficient contact magnitude.",
        ],
    )
    _boxed_bullets(prs, "H2 Result: Validated And Falsified Hypotheses", "Takeaway", [
        "Validated: larger `contact_collision_dist / particle_radius` clearly reduces penetration, and bunny mass helps but is not the main lever.",
        "Validated: thick box is easier because its support patch is broader and flatter than the bunny body or ears.",
        f"Validated: even 10x radius still leaves {float(thin_ear_10x_summary['max_penetration_depth_bunny_mesh_m']) * 1000:.2f} mm bunny-mesh penetration on the thin-ear geometry.",
        "Falsified: lower height alone or brute-force stiffness / damping increases are not robust fixes for bunny support.",
    ])

    _body(prs, "H3: Newton Self-Collision Is Not PhysTwin Parity", [
        "We test whether Newton native self-collision can honestly be presented as a PhysTwin parity result.",
        "This is a model-comparison question, not another cloth-penetration tuning question.",
        "Official Newton self-contact examples do exist, but the mainline examples use VBD rather than SemiImplicit.",
        "If H3 is true, `self-collision ON` should be treated only as a Newton-native ablation.",
    ])
    _pic_twocol(
        prs,
        "H3: PhysTwin And Newton Use Different Self-Collision Models",
        "PhysTwin: approach-gated contact plus impulse-style velocity update",
        phys_png,
        "Newton: overlap-triggered contact plus penalty-force accumulation",
        newton_self_png,
    )
    _boxed_bullets(prs, "H3 Result: `self-collision ON` Is Only A Newton-Native Ablation", "Takeaway", [
        "Takeaway: Newton native self-collision is not a PhysTwin parity target.",
        "Official Newton examples do show self-contact, but those examples are VBD-based, not SemiImplicit self-collision parity cases.",
        "`self-collision ON` can still be shown as a Newton-native ablation, but it should not be presented as the main engineering fix.",
        "So H3 narrows the story instead of opening another tuning branch.",
    ])

    _body(prs, "H4: Pair-Filtered Multi-Deformable Interaction Is Viable", [
        "We test whether multi-deformable interaction can work once pair filtering is made explicit.",
        "The hypothesis is not “rewrite the physics”, but “change pair selection and keep the Newton penalty model”.",
        "If H4 is true, we should already be able to show minimal positive demos.",
    ])
    _picture_with_footer(prs, "H4: We Changed Pair Selection, Not The Physics Model", cross_rope_png, [
        "Takeaway: we keep Newton's HashGrid and penalty model, and only change which particle pairs are allowed to interact.",
        "So H4 is a targeted pair-filtering intervention, not a new solver.",
    ])
    _gif_twocol(
        prs,
        "H4 Result: Filtered Cross-Object Contact Already Works In Two Demos",
        "Two ropes on ground\nrope particle mass = 1.0 kg | vertical gap = 0.35 m",
        rope_ground_gif,
        "One rope + grounded sloth\nrope drop height = 0.35 m | sloth rests on ground",
        multi_rope_gif,
        common_settings="No rigid bunny | gravity = 9.8 m/s^2 | drag on | Newton particle contact filtered to cross-object pairs only",
    )
    _gif_twocol_large(
        prs,
        "H4 Addendum: Two-Rope Weight Contrast Still Works",
        (
            f"Case 1: lower rope = {float(two_ropes_weight_case1_summary['rope_total_masses'][0]):g} kg, "
            f"upper rope = {float(two_ropes_weight_case1_summary['rope_total_masses'][1]):g} kg"
        ),
        two_ropes_weight_case1_gif,
        (
            f"Case 2: lower rope = {float(two_ropes_weight_case2_summary['rope_total_masses'][0]):g} kg, "
            f"upper rope = {float(two_ropes_weight_case2_summary['rope_total_masses'][1]):g} kg"
        ),
        two_ropes_weight_case2_gif,
        common_settings=(
            "No rigid bunny | lower rope rests on ground | upper rope drops in a cross layout | "
            "Newton particle contact filtered to cross-object pairs only | gravity = 9.8 m/s^2 | drag on | "
            "video ends 1 second after first cross-rope contact"
        ),
    )
    _gif_single(
        prs,
        "H4 Addendum: Two Crossed Ropes Also Work On A Rigid Box",
        two_ropes_box_stack_gif,
        (
            f"rope totals = {float(two_ropes_box_stack_summary['rope_total_masses'][0]):.1f} kg + "
            f"{float(two_ropes_box_stack_summary['rope_total_masses'][1]):.1f} kg | "
            f"one rigid box = {float(two_ropes_box_stack_summary['box_total_mass_kg']):.1f} kg | "
            f"box half extents = {float(two_ropes_box_stack_summary['box_half_extents'][0]):.2f} m | "
            f"drop height = {float(two_ropes_box_stack_summary['drop_height_m']):.2f} m"
        ),
    )
    _gif_single(
        prs,
        "H4 Addendum: Native Franka Can Lift And Release The Bridge Rope",
        robot_rope_franka_gif,
        (
            f"task = {str(robot_rope_franka_summary['task'])} | "
            f"native robot = {str(robot_rope_franka_summary['robot_geometry'])} | "
            f"first contact = {float(robot_rope_franka_summary['first_contact_time_s']) * 1000:.1f} ms | "
            f"contact duration = {float(robot_rope_franka_summary['contact_duration_s']) * 1000:.1f} ms | "
            f"rope COM displacement = {float(robot_rope_franka_summary['rope_com_displacement_m']) * 1000:.1f} mm | "
            f"peak midpoint lift = {float(robot_rope_franka_summary['rope_mid_segment_peak_z_delta_m']) * 1000:.1f} mm"
        ),
    )

    _body(prs, "H5: MPM One-Way Is Useful; Two-Way Is Not The Main Path", [
        "We test whether MPM should remain a mainline route for spring-mass two-way coupling.",
        "The positive case is one-way kinematic interaction; the real question is whether two-way support is architecturally promising.",
        "If H5 fails, MPM should stay as a side result, not the main story.",
    ])
    _gif_single(
        prs,
        "H5: One-Way MPM Works; Two-Way Should Be Deprioritized",
        mpm_sloth_gif,
        "Takeaway: controller-driven sloth already works as a one-way kinematic collider, but current two-way spring-mass support remains the architectural mismatch.",
    )
    _picture_with_footer(
        prs,
        "Newton Playground Profiling: Rope Replay Is Solver-Bound, Not Collision-Bound",
        rope_profile_png,
        [
            "Full replay result with rendering disabled: the dominant costs are `integrate_particles`, `spring_forces`, `write_kinematic_state`, and `drag_correction`.",
            "So the next optimization target is the internal rope spring-mass update path; broad-phase collision is effectively zero in this replay case.",
        ],
    )
    _boxed_bullets(prs, "Summary: Which Hypotheses Survived", "Takeaway", [
        "Validated: H1 bridge interaction, H2 external-support diagnosis, and H4 pair-filtered multi-deformable interaction.",
        "Extended: the same bridge rope now also supports a native Newton Franka Panda `lift and release` baseline without falling back to a proxy pusher.",
        "Reframed: H3 shows `self-collision ON` is only a Newton-native ablation, not PhysTwin parity evidence.",
        "Deprioritized: H5 shows MPM one-way is useful, but MPM two-way is not the main path for this project.",
    ])
    _body(prs, "Next Step: Fewer Slides, Stronger Evidence", [
        "Compress each hypothesis to one flow slide, one source-proof slide, and one best demo slide.",
        "Keep only the strongest bunny-support evidence and move backup sweeps or reference runs to transcript or appendix.",
        "Use the freed time to record a cleaner robot or fancy multi-deformable demo instead of adding more parameter pages.",
    ])


TRANSCRIPT_SOURCE_MD = MEETING_DIR / "transcript_source.md"
if TRANSCRIPT_SOURCE_MD.exists():
    TRANSCRIPT_MEETING = TRANSCRIPT_SOURCE_MD.read_text(encoding="utf-8")


def _markdown_to_pdf(md_text: str, html_path: Path, pdf_path: Path) -> None:
    html_body = markdown.markdown(md_text, extensions=["tables", "fenced_code"])
    html = f"""<!doctype html>
<html>
<head>
  <meta charset="utf-8">
  <style>
    body {{ font-family: 'DejaVu Sans', sans-serif; line-height: 1.55; margin: 36px 48px; color: #222; }}
    h1, h2, h3 {{ color: #1f3d5a; }}
    code {{ background: #f2f4f7; padding: 1px 4px; border-radius: 4px; }}
    pre {{ background: #f2f4f7; padding: 10px; overflow-x: auto; }}
    blockquote {{ color: #555; border-left: 4px solid #ccd5e0; padding-left: 12px; }}
    table {{ border-collapse: collapse; width: 100%; }}
    th, td {{ border: 1px solid #ccd5e0; padding: 6px 8px; }}
    a {{ color: #1f5d99; text-decoration: none; }}
    </style>
</head>
<body>{html_body}</body>
</html>"""
    html_path.write_text(html, encoding="utf-8")
    try:
        from weasyprint import HTML

        HTML(filename=str(html_path)).write_pdf(str(pdf_path))
        return
    except Exception:
        pass

    subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(pdf_path.parent),
            str(html_path),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def _renumber_transcript_headings(md_text: str) -> str:
    slide_no = 0

    def repl(match: re.Match[str]) -> str:
        nonlocal slide_no
        slide_no += 1
        title = match.group(1)
        return f"## Slide {slide_no} — {title}"

    return re.sub(r"^[ \t]*## Slide [^—]+ — (.+)$", repl, md_text, flags=re.MULTILINE)


def _remove_path(path: Path) -> None:
    if not path.exists():
        return
    if path.is_dir():
        shutil.rmtree(path)
    else:
        path.unlink()


def _load_json_or_fallback(path: Path, fallback: dict[str, float]) -> dict[str, float]:
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8"))
    return dict(fallback)


def _clean_tmp_slide_outputs() -> list[Path]:
    removed: list[Path] = []
    tmp = ROOT / "tmp"
    if not tmp.exists():
        return removed

    for child in tmp.iterdir():
        name = child.name.lower()
        if child.is_dir() and ("meeting_slides" in name or "imagecheck" in name):
            _remove_path(child)
            removed.append(child)
            continue
        if child.is_file():
            is_slide_file = name.startswith("bridge_meeting_") and child.suffix.lower() == ".pptx"
            is_transcript_file = name in {"transcript.md", "transcript.html", "transcript.pdf"}
            if is_slide_file or is_transcript_file:
                _remove_path(child)
                removed.append(child)
    return removed


def _reset_generated_outputs(out: Path) -> None:
    if not out.exists():
        return
    for path in (
        out / f"bridge_meeting_{FORMAL_SLIDE_TAG}.pptx",
        out / "transcript.md",
        out / "transcript.html",
        out / "transcript.pdf",
    ):
        _remove_path(path)

    for path in (out / "gif", out / "images"):
        _remove_path(path)


def main() -> int:
    args = parse_args()
    out = args.out_dir.resolve()
    removed_tmp = _clean_tmp_slide_outputs()
    out.mkdir(parents=True, exist_ok=True)
    _reset_generated_outputs(out)
    gif_dir = out / "gif"
    img_dir = out / "images"
    gif_dir.mkdir(parents=True, exist_ok=True)
    img_dir.mkdir(parents=True, exist_ok=True)

    rows = _load_rows(GRID_CSV) if GRID_CSV.exists() else []
    pptx_path = out / f"bridge_meeting_{FORMAL_SLIDE_TAG}.pptx"

    template_pptx = MEETING_TEMPLATE_PPTX.expanduser().resolve()
    if not template_pptx.exists():
        raise SystemExit(f"Missing slide template: {MEETING_TEMPLATE_PPTX}")
    shutil.copy2(template_pptx, pptx_path)
    prs = Presentation(str(template_pptx))
    build_deck_reference(prs, rows, gif_dir, img_dir)
    prs.save(str(pptx_path))

    transcript_text = _renumber_transcript_headings(TRANSCRIPT_MEETING)

    transcript_md = out / "transcript.md"
    transcript_md.write_text(transcript_text, encoding="utf-8")

    transcript_html = out / "transcript.html"
    transcript_pdf = out / "transcript.pdf"
    _markdown_to_pdf(transcript_text, transcript_html, transcript_pdf)

    print(f"PPTX: {pptx_path}")
    print(f"Transcript MD: {transcript_md}")
    print(f"Transcript PDF: {transcript_pdf}")
    print(f"Output Dir: {out}")
    print(f"Removed tmp slide outputs: {len(removed_tmp)}")
    print(f"Slides: {len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
